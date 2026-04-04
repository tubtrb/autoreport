"""Template profiling and fill planning for the public Autoreport contract."""

from __future__ import annotations

from hashlib import sha1
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from autoreport.models import (
    ImageSpec,
    PayloadSlide,
    ReportPayload,
    SlotOverride,
    TemplateContract,
    TemplatePatternContract,
    TemplateSectionContract,
    TemplateSlotContract,
    TEMPLATE_CONTRACT_VERSION,
)
from autoreport.outputs.errors import TemplateCompatibilityError
from autoreport.templates.autofill import (
    DiagnosticReport,
    FillPlan,
    FitStatus,
    PatternProfile,
    PlannedImageFill,
    PlannedSlide,
    PlannedTextFill,
    SlideDecoration,
    SlotContentKind,
    SlotDescriptor,
    TemplateProfile,
    fit_text_items_to_slot,
    fit_text_to_slot,
    sort_slots_in_reading_order,
)


TEMPLATE_NAME = "weekly_report"
BASIC_TEMPLATE_NAME = "autoreport_editorial"
MANUAL_TEMPLATE_NAME = "autoreport_manual"
SUPPORTED_TEMPLATE_NAMES = (BASIC_TEMPLATE_NAME, MANUAL_TEMPLATE_NAME, TEMPLATE_NAME)
BUILT_IN_TEMPLATE_ID = "autoreport-editorial-v1"
BUILT_IN_TEMPLATE_LABEL = "Autoreport Editorial"
MANUAL_TEMPLATE_ID = "autoreport-manual-v1"
MANUAL_TEMPLATE_LABEL = "Autoreport Manual"
TITLE_LAYOUT_INDEX = 0
BLANK_LAYOUT_INDEX = 6

TITLE_FONT_SIZE = 34
TITLE_MIN_FONT_SIZE = 20
SUBTITLE_FONT_SIZE = 20
SUBTITLE_MIN_FONT_SIZE = 14
SECTION_TITLE_FONT_SIZE = 24
SECTION_TITLE_MIN_FONT_SIZE = 20
SECTION_BODY_FONT_SIZE = 18
SECTION_BODY_MIN_FONT_SIZE = 13
CAPTION_FONT_SIZE = 14
CAPTION_MIN_FONT_SIZE = 11

BASIC_EDITORIAL_NAVY = (14, 40, 65)
BASIC_EDITORIAL_TEAL = (21, 96, 130)
BASIC_EDITORIAL_ORANGE = (233, 113, 50)
BASIC_EDITORIAL_LIGHT = (247, 248, 250)
BASIC_EDITORIAL_PALE = (234, 242, 246)
BASIC_EDITORIAL_LINE = (217, 224, 231)
BASIC_EDITORIAL_PALE_LINE = (227, 236, 242)
MANUAL_BG = (248, 250, 252)
MANUAL_TEXT_PANEL = (255, 255, 255)
MANUAL_FRAME = (226, 232, 240)
MANUAL_ACCENT = (22, 78, 99)
MANUAL_ACCENT_SOFT = (186, 230, 253)
MANUAL_ACTION = (15, 118, 110)
MANUAL_IMAGE_PANEL = (241, 245, 249)

IGNORED_TEXT_PLACEHOLDER_TYPES = frozenset(
    getattr(PP_PLACEHOLDER, name)
    for name in ("DATE", "FOOTER", "SLIDE_NUMBER")
    if hasattr(PP_PLACEHOLDER, name)
)
TITLE_PLACEHOLDER_TYPES = frozenset(
    getattr(PP_PLACEHOLDER, name)
    for name in ("TITLE", "CENTER_TITLE")
    if hasattr(PP_PLACEHOLDER, name)
)
TEXT_EXCLUDED_PLACEHOLDER_TYPES = (
    IGNORED_TEXT_PLACEHOLDER_TYPES
    | TITLE_PLACEHOLDER_TYPES
    | frozenset(
        getattr(PP_PLACEHOLDER, name)
        for name in ("SUBTITLE", "PICTURE")
        if hasattr(PP_PLACEHOLDER, name)
    )
)
PICTURE_PLACEHOLDER_TYPES = frozenset(
    getattr(PP_PLACEHOLDER, name)
    for name in ("PICTURE",)
    if hasattr(PP_PLACEHOLDER, name)
)


def profile_template(
    presentation: Presentation,
    *,
    template_path: Path | None,
    template_name: str = BASIC_TEMPLATE_NAME,
) -> TemplateProfile:
    """Profile either a built-in template family or a user template."""

    if template_name == BASIC_TEMPLATE_NAME and template_path is None:
        return profile_basic_template(presentation, template_path=None)
    if template_name == MANUAL_TEMPLATE_NAME and template_path is None:
        return profile_manual_template(presentation, template_path=None)

    return profile_user_template(presentation, template_path=template_path)


def profile_basic_template(
    presentation: Presentation,
    *,
    template_path: Path | None,
) -> TemplateProfile:
    """Build the built-in editorial profile with multiple generation patterns."""

    blank_layout = _get_layout(
        presentation,
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name="blank",
        template_path=template_path,
    )

    title_title_slot = _build_text_box_slot(
        presentation,
        layout_index=BLANK_LAYOUT_INDEX,
        slot_name="cover.title",
        alias="title",
        slot_type="title",
        x_ratio=0.315,
        y_ratio=0.251,
        width_ratio=0.469,
        height_ratio=0.160,
        preferred_font_size=TITLE_FONT_SIZE,
        min_font_size=TITLE_MIN_FONT_SIZE,
        allowed_kinds=(SlotContentKind.TITLE,),
    )
    title_subtitle_slot = _build_text_box_slot(
        presentation,
        layout_index=BLANK_LAYOUT_INDEX,
        slot_name="cover.subtitle_1",
        alias="subtitle_1",
        slot_type="text",
        x_ratio=0.319,
        y_ratio=0.451,
        width_ratio=0.353,
        height_ratio=0.127,
        preferred_font_size=SUBTITLE_FONT_SIZE,
        min_font_size=SUBTITLE_MIN_FONT_SIZE,
        allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
        orientation="stack",
        order=1,
    )
    title_pattern = PatternProfile(
        pattern_id="cover.editorial",
        kind="cover",
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name=blank_layout.name,
        slots=(title_title_slot, title_subtitle_slot),
        decorations=_build_basic_template_title_decorations(presentation),
    )

    contents_pattern = PatternProfile(
        pattern_id="contents.editorial",
        kind="contents",
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name=blank_layout.name,
        slots=(
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="contents.title",
                alias="title",
                slot_type="title",
                x_ratio=0.077,
                y_ratio=0.085,
                width_ratio=0.581,
                height_ratio=0.069,
                preferred_font_size=SECTION_TITLE_FONT_SIZE,
                min_font_size=SECTION_TITLE_MIN_FONT_SIZE,
                allowed_kinds=(SlotContentKind.TITLE,),
            ),
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="contents.body_1",
                alias="body_1",
                slot_type="text",
                x_ratio=0.081,
                y_ratio=0.237,
                width_ratio=0.566,
                height_ratio=0.597,
                preferred_font_size=SECTION_BODY_FONT_SIZE,
                min_font_size=SECTION_BODY_MIN_FONT_SIZE,
                allowed_kinds=(
                    SlotContentKind.PARAGRAPH_OR_BULLETS,
                    SlotContentKind.SHORT_FACT_OR_STATUS,
                ),
                orientation="stack",
                order=1,
            ),
        ),
        decorations=_build_basic_template_body_decorations(presentation),
    )

    slide_patterns = (
        PatternProfile(
            pattern_id="text.editorial",
            kind="text",
            layout_index=BLANK_LAYOUT_INDEX,
            layout_name=blank_layout.name,
            slots=(
                _build_text_box_slot(
                    presentation,
                    layout_index=BLANK_LAYOUT_INDEX,
                    slot_name="text.title",
                    alias="title",
                    slot_type="title",
                    x_ratio=0.077,
                    y_ratio=0.085,
                    width_ratio=0.581,
                    height_ratio=0.069,
                    preferred_font_size=SECTION_TITLE_FONT_SIZE,
                    min_font_size=SECTION_TITLE_MIN_FONT_SIZE,
                    allowed_kinds=(SlotContentKind.TITLE,),
                ),
                _build_text_box_slot(
                    presentation,
                    layout_index=BLANK_LAYOUT_INDEX,
                    slot_name="text.body_1",
                    alias="body_1",
                    slot_type="text",
                    x_ratio=0.081,
                    y_ratio=0.237,
                    width_ratio=0.566,
                    height_ratio=0.597,
                    preferred_font_size=SECTION_BODY_FONT_SIZE,
                    min_font_size=SECTION_BODY_MIN_FONT_SIZE,
                    allowed_kinds=(
                        SlotContentKind.PARAGRAPH_OR_BULLETS,
                        SlotContentKind.SHORT_FACT_OR_STATUS,
                    ),
                    orientation="stack",
                    order=1,
                ),
            ),
            decorations=_build_basic_template_body_decorations(presentation),
        ),
        PatternProfile(
            pattern_id="metrics.editorial",
            kind="metrics",
            layout_index=BLANK_LAYOUT_INDEX,
            layout_name=blank_layout.name,
            slots=(
                _build_text_box_slot(
                    presentation,
                    layout_index=BLANK_LAYOUT_INDEX,
                    slot_name="metrics.title",
                    alias="title",
                    slot_type="title",
                    x_ratio=0.077,
                    y_ratio=0.085,
                    width_ratio=0.581,
                    height_ratio=0.069,
                    preferred_font_size=SECTION_TITLE_FONT_SIZE,
                    min_font_size=SECTION_TITLE_MIN_FONT_SIZE,
                    allowed_kinds=(SlotContentKind.TITLE,),
                ),
                _build_text_box_slot(
                    presentation,
                    layout_index=BLANK_LAYOUT_INDEX,
                    slot_name="metrics.body_1",
                    alias="body_1",
                    slot_type="text",
                    x_ratio=0.081,
                    y_ratio=0.237,
                    width_ratio=0.566,
                    height_ratio=0.597,
                    preferred_font_size=SECTION_BODY_FONT_SIZE,
                    min_font_size=SECTION_BODY_MIN_FONT_SIZE,
                    allowed_kinds=(SlotContentKind.METRIC_LIST,),
                    orientation="stack",
                    order=1,
                ),
            ),
            decorations=_build_basic_template_body_decorations(presentation),
        ),
        *_build_basic_template_text_image_patterns(
            presentation,
            layout_name=blank_layout.name,
        ),
    )

    return TemplateProfile(
        template_name=BASIC_TEMPLATE_NAME,
        template_id=BUILT_IN_TEMPLATE_ID,
        template_label=BUILT_IN_TEMPLATE_LABEL,
        template_source="built_in",
        template_path=None,
        title_pattern=title_pattern,
        contents_pattern=contents_pattern,
        slide_patterns=slide_patterns,
    )


def profile_manual_template(
    presentation: Presentation,
    *,
    template_path: Path | None,
) -> TemplateProfile:
    """Build the built-in screenshot-first manual profile."""

    blank_layout = _get_layout(
        presentation,
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name="blank",
        template_path=template_path,
    )

    title_pattern = PatternProfile(
        pattern_id="cover.manual",
        kind="cover",
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name=blank_layout.name,
        slots=(
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="cover.title",
                alias="doc_title",
                slot_type="title",
                x_ratio=0.083,
                y_ratio=0.122,
                width_ratio=0.611,
                height_ratio=0.116,
                preferred_font_size=30,
                min_font_size=20,
                allowed_kinds=(SlotContentKind.TITLE,),
            ),
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="cover.doc_subtitle",
                alias="doc_subtitle",
                slot_type="text",
                x_ratio=0.083,
                y_ratio=0.274,
                width_ratio=0.611,
                height_ratio=0.114,
                preferred_font_size=18,
                min_font_size=13,
                allowed_kinds=(SlotContentKind.PARAGRAPH_OR_BULLETS, SlotContentKind.SHORT_FACT_OR_STATUS),
                orientation="stack",
                order=1,
                required=False,
            ),
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="cover.doc_version",
                alias="doc_version",
                slot_type="text",
                x_ratio=0.083,
                y_ratio=0.720,
                width_ratio=0.195,
                height_ratio=0.056,
                preferred_font_size=13,
                min_font_size=11,
                allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
                orientation="horizontal",
                order=2,
                required=False,
            ),
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="cover.author_or_owner",
                alias="author_or_owner",
                slot_type="text",
                x_ratio=0.298,
                y_ratio=0.720,
                width_ratio=0.255,
                height_ratio=0.056,
                preferred_font_size=13,
                min_font_size=11,
                allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
                orientation="horizontal",
                order=3,
                required=False,
            ),
        ),
        decorations=_build_manual_cover_decorations(presentation),
    )

    contents_pattern = PatternProfile(
        pattern_id="contents.manual",
        kind="contents",
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name=blank_layout.name,
        slots=(
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="contents.title",
                alias="contents_title",
                slot_type="title",
                x_ratio=0.083,
                y_ratio=0.088,
                width_ratio=0.611,
                height_ratio=0.074,
                preferred_font_size=24,
                min_font_size=18,
                allowed_kinds=(SlotContentKind.TITLE,),
            ),
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="contents.group_label",
                alias="contents_group_label",
                slot_type="text",
                x_ratio=0.083,
                y_ratio=0.194,
                width_ratio=0.380,
                height_ratio=0.052,
                preferred_font_size=14,
                min_font_size=11,
                allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
                orientation="stack",
                order=1,
                required=False,
            ),
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name="contents.body_1",
                alias="contents_entries",
                slot_type="text",
                x_ratio=0.083,
                y_ratio=0.286,
                width_ratio=0.547,
                height_ratio=0.462,
                preferred_font_size=18,
                min_font_size=13,
                allowed_kinds=(SlotContentKind.PARAGRAPH_OR_BULLETS, SlotContentKind.SHORT_FACT_OR_STATUS),
                orientation="stack",
                order=2,
            ),
        ),
        decorations=_build_manual_text_decorations(presentation),
    )

    slide_patterns = (
        PatternProfile(
            pattern_id="text.manual.section_break",
            kind="text",
            layout_index=BLANK_LAYOUT_INDEX,
            layout_name=blank_layout.name,
            slots=(
                _build_text_box_slot(
                    presentation,
                    layout_index=BLANK_LAYOUT_INDEX,
                    slot_name="text.title",
                    alias="section_title",
                    slot_type="title",
                    x_ratio=0.171,
                    y_ratio=0.316,
                    width_ratio=0.606,
                    height_ratio=0.112,
                    preferred_font_size=28,
                    min_font_size=18,
                    allowed_kinds=(SlotContentKind.TITLE,),
                ),
                _build_text_box_slot(
                    presentation,
                    layout_index=BLANK_LAYOUT_INDEX,
                    slot_name="text.section_no",
                    alias="section_no",
                    slot_type="text",
                    x_ratio=0.083,
                    y_ratio=0.316,
                    width_ratio=0.067,
                    height_ratio=0.112,
                    preferred_font_size=24,
                    min_font_size=18,
                    allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
                    orientation="stack",
                    order=1,
                ),
                _build_text_box_slot(
                    presentation,
                    layout_index=BLANK_LAYOUT_INDEX,
                    slot_name="text.section_subtitle",
                    alias="section_subtitle",
                    slot_type="text",
                    x_ratio=0.171,
                    y_ratio=0.470,
                    width_ratio=0.606,
                    height_ratio=0.100,
                    preferred_font_size=16,
                    min_font_size=12,
                    allowed_kinds=(SlotContentKind.PARAGRAPH_OR_BULLETS, SlotContentKind.SHORT_FACT_OR_STATUS),
                    orientation="stack",
                    order=2,
                    required=False,
                ),
            ),
            decorations=_build_manual_text_decorations(presentation),
        ),
        *_build_manual_procedure_patterns(
            presentation,
            layout_name=blank_layout.name,
        ),
    )

    return TemplateProfile(
        template_name=MANUAL_TEMPLATE_NAME,
        template_id=MANUAL_TEMPLATE_ID,
        template_label=MANUAL_TEMPLATE_LABEL,
        template_source="built_in",
        template_path=None,
        title_pattern=title_pattern,
        contents_pattern=contents_pattern,
        slide_patterns=slide_patterns,
    )


def _build_manual_procedure_patterns(
    presentation: Presentation,
    *,
    layout_name: str,
) -> tuple[PatternProfile, ...]:
    return (
        _build_manual_procedure_pattern(
            presentation,
            pattern_id="text_image.manual.procedure.one",
            layout_name=layout_name,
            body_region=(0.083, 0.372, 0.532, 0.330),
            image_regions=((0.664, 0.214, 0.260, 0.408),),
            caption_regions=((0.664, 0.654, 0.260, 0.066),),
        ),
        _build_manual_procedure_pattern(
            presentation,
            pattern_id="text_image.manual.procedure.two",
            layout_name=layout_name,
            body_region=(0.083, 0.340, 0.841, 0.108),
            image_regions=(
                (0.083, 0.516, 0.377, 0.208),
                (0.548, 0.516, 0.377, 0.208),
            ),
            caption_regions=(
                (0.083, 0.744, 0.377, 0.056),
                (0.548, 0.744, 0.377, 0.056),
            ),
        ),
        _build_manual_procedure_pattern(
            presentation,
            pattern_id="text_image.manual.procedure.three",
            layout_name=layout_name,
            body_region=(0.083, 0.340, 0.841, 0.094),
            image_regions=(
                (0.083, 0.500, 0.255, 0.188),
                (0.372, 0.500, 0.255, 0.188),
                (0.661, 0.500, 0.255, 0.188),
            ),
            caption_regions=(
                (0.083, 0.710, 0.255, 0.056),
                (0.372, 0.710, 0.255, 0.056),
                (0.661, 0.710, 0.255, 0.056),
            ),
        ),
    )


def _build_manual_procedure_pattern(
    presentation: Presentation,
    *,
    pattern_id: str,
    layout_name: str,
    body_region: tuple[float, float, float, float],
    image_regions: tuple[tuple[float, float, float, float], ...],
    caption_regions: tuple[tuple[float, float, float, float], ...],
) -> PatternProfile:
    slots: list[SlotDescriptor] = [
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.title",
            alias="step_title",
            slot_type="title",
            x_ratio=0.171,
            y_ratio=0.088,
            width_ratio=0.543,
            height_ratio=0.066,
            preferred_font_size=22,
            min_font_size=17,
            allowed_kinds=(SlotContentKind.TITLE,),
        ),
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.step_no",
            alias="step_no",
            slot_type="text",
            x_ratio=0.083,
            y_ratio=0.088,
            width_ratio=0.067,
            height_ratio=0.066,
            preferred_font_size=18,
            min_font_size=14,
            allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
            orientation="stack",
            order=1,
        ),
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.command_or_action",
            alias="command_or_action",
            slot_type="text",
            x_ratio=0.083,
            y_ratio=0.182,
            width_ratio=0.841,
            height_ratio=0.050,
            preferred_font_size=13,
            min_font_size=11,
            allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
            orientation="stack",
            order=2,
            required=False,
        ),
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.summary",
            alias="summary",
            slot_type="text",
            x_ratio=0.083,
            y_ratio=0.250,
            width_ratio=0.841,
            height_ratio=0.068,
            preferred_font_size=15,
            min_font_size=12,
            allowed_kinds=(SlotContentKind.PARAGRAPH_OR_BULLETS, SlotContentKind.SHORT_FACT_OR_STATUS),
            orientation="stack",
            order=3,
            required=False,
        ),
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.body_1",
            alias="detail_body",
            slot_type="text",
            x_ratio=body_region[0],
            y_ratio=body_region[1],
            width_ratio=body_region[2],
            height_ratio=body_region[3],
            preferred_font_size=15,
            min_font_size=11,
            allowed_kinds=(SlotContentKind.PARAGRAPH_OR_BULLETS,),
            orientation="stack",
            order=4,
        ),
    ]

    for index, region in enumerate(image_regions, start=1):
        slots.append(
            _build_geometry_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name=f"text_image.image_{index}",
                alias=f"image_{index}",
                slot_type="image",
                x_ratio=region[0],
                y_ratio=region[1],
                width_ratio=region[2],
                height_ratio=region[3],
                orientation="horizontal",
                order=index,
            )
        )

    for index, region in enumerate(caption_regions, start=1):
        slots.append(
            _build_text_box_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name=f"text_image.caption_{index}",
                alias=f"caption_{index}",
                slot_type="caption",
                x_ratio=region[0],
                y_ratio=region[1],
                width_ratio=region[2],
                height_ratio=region[3],
                preferred_font_size=12,
                min_font_size=10,
                allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
                orientation="horizontal",
                order=index,
                required=False,
            )
        )

    return PatternProfile(
        pattern_id=pattern_id,
        kind="text_image",
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name=layout_name,
        slots=tuple(slots),
        decorations=_build_manual_text_image_decorations(
            presentation,
            body_region=body_region,
            image_regions=image_regions,
            caption_regions=caption_regions,
        ),
    )


def _build_basic_template_text_image_patterns(
    presentation: Presentation,
    *,
    layout_name: str,
) -> tuple[PatternProfile, ...]:
    return (
        _build_basic_text_image_pattern(
            presentation,
            pattern_id="text_image.editorial",
            layout_name=layout_name,
            body_region=(0.081, 0.237, 0.566, 0.597),
            image_regions=((0.709, 0.197, 0.227, 0.470),),
            caption_region=(0.709, 0.700, 0.227, 0.095),
            image_orientation="stack",
        ),
        _build_basic_text_image_pattern(
            presentation,
            pattern_id="text_image.editorial.two_horizontal",
            layout_name=layout_name,
            body_region=(0.081, 0.210, 0.855, 0.180),
            image_regions=(
                (0.081, 0.445, 0.405, 0.255),
                (0.531, 0.445, 0.405, 0.255),
            ),
            caption_region=(0.081, 0.735, 0.855, 0.080),
            image_orientation="horizontal",
        ),
        _build_basic_text_image_pattern(
            presentation,
            pattern_id="text_image.editorial.two_vertical",
            layout_name=layout_name,
            body_region=(0.081, 0.237, 0.360, 0.525),
            image_regions=(
                (0.531, 0.237, 0.405, 0.235),
                (0.531, 0.505, 0.405, 0.235),
            ),
            caption_region=(0.531, 0.770, 0.405, 0.060),
            image_orientation="vertical",
        ),
        _build_basic_text_image_pattern(
            presentation,
            pattern_id="text_image.editorial.three_horizontal",
            layout_name=layout_name,
            body_region=(0.081, 0.205, 0.855, 0.155),
            image_regions=(
                (0.081, 0.420, 0.255, 0.225),
                (0.381, 0.420, 0.255, 0.225),
                (0.681, 0.420, 0.255, 0.225),
            ),
            caption_region=(0.081, 0.690, 0.855, 0.080),
            image_orientation="horizontal",
        ),
        _build_basic_text_image_pattern(
            presentation,
            pattern_id="text_image.editorial.three_vertical",
            layout_name=layout_name,
            body_region=(0.081, 0.237, 0.300, 0.560),
            image_regions=(
                (0.454, 0.205, 0.482, 0.150),
                (0.454, 0.390, 0.482, 0.150),
                (0.454, 0.575, 0.482, 0.150),
            ),
            caption_region=(0.454, 0.760, 0.482, 0.060),
            image_orientation="vertical",
        ),
    )


def _build_basic_text_image_pattern(
    presentation: Presentation,
    *,
    pattern_id: str,
    layout_name: str,
    body_region: tuple[float, float, float, float],
    image_regions: tuple[tuple[float, float, float, float], ...],
    caption_region: tuple[float, float, float, float],
    image_orientation: str,
) -> PatternProfile:
    slots: list[SlotDescriptor] = [
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.title",
            alias="title",
            slot_type="title",
            x_ratio=0.077,
            y_ratio=0.085,
            width_ratio=0.581,
            height_ratio=0.069,
            preferred_font_size=SECTION_TITLE_FONT_SIZE,
            min_font_size=SECTION_TITLE_MIN_FONT_SIZE,
            allowed_kinds=(SlotContentKind.TITLE,),
        ),
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.body_1",
            alias="body_1",
            slot_type="text",
            x_ratio=body_region[0],
            y_ratio=body_region[1],
            width_ratio=body_region[2],
            height_ratio=body_region[3],
            preferred_font_size=SECTION_BODY_FONT_SIZE,
            min_font_size=SECTION_BODY_MIN_FONT_SIZE,
            allowed_kinds=(SlotContentKind.PARAGRAPH_OR_BULLETS,),
            orientation="stack",
            order=1,
        ),
    ]

    for index, region in enumerate(image_regions, start=1):
        slots.append(
            _build_geometry_slot(
                presentation,
                layout_index=BLANK_LAYOUT_INDEX,
                slot_name=f"text_image.image_{index}",
                alias=f"image_{index}",
                slot_type="image",
                x_ratio=region[0],
                y_ratio=region[1],
                width_ratio=region[2],
                height_ratio=region[3],
                orientation=image_orientation,
                order=index,
            )
        )

    slots.append(
        _build_text_box_slot(
            presentation,
            layout_index=BLANK_LAYOUT_INDEX,
            slot_name="text_image.caption_1",
            alias="caption_1",
            slot_type="caption",
            x_ratio=caption_region[0],
            y_ratio=caption_region[1],
            width_ratio=caption_region[2],
            height_ratio=caption_region[3],
            preferred_font_size=CAPTION_FONT_SIZE,
            min_font_size=CAPTION_MIN_FONT_SIZE,
            allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
            orientation=image_orientation,
            order=1,
            required=False,
        )
    )

    return PatternProfile(
        pattern_id=pattern_id,
        kind="text_image",
        layout_index=BLANK_LAYOUT_INDEX,
        layout_name=layout_name,
        slots=tuple(slots),
        decorations=_build_basic_template_body_decorations(presentation),
    )


def profile_user_template(
    presentation: Presentation,
    *,
    template_path: Path | None,
) -> TemplateProfile:
    """Profile a user-supplied PowerPoint template into Autoreport patterns."""

    (
        title_layout_index,
        title_layout,
        title_placeholder,
        subtitle_placeholders,
    ) = _select_title_layout(
        presentation,
        template_path=template_path,
    )
    (
        body_layout_index,
        body_layout,
        body_title_placeholder,
        body_text_placeholders,
    ) = _select_body_layout(
        presentation,
        template_path=template_path,
        preferred_skip_layout_index=title_layout_index,
    )

    title_slots = (
        _build_placeholder_slot(
            title_placeholder,
            layout_index=title_layout_index,
            slot_name="cover.title",
            alias="title",
            slot_type="title",
            preferred_font_size=TITLE_FONT_SIZE,
            min_font_size=TITLE_MIN_FONT_SIZE,
            allowed_kinds=(SlotContentKind.TITLE,),
        ),
        *(
            _build_placeholder_slot(
                placeholder,
                layout_index=title_layout_index,
                slot_name=f"cover.subtitle_{index + 1}",
                alias=f"subtitle_{index + 1}",
                slot_type="text",
                preferred_font_size=SUBTITLE_FONT_SIZE,
                min_font_size=SUBTITLE_MIN_FONT_SIZE,
                allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
                orientation="stack",
                order=index + 1,
            )
            for index, placeholder in enumerate(subtitle_placeholders)
        ),
    )
    title_pattern = PatternProfile(
        pattern_id="cover.template",
        kind="cover",
        layout_index=title_layout_index,
        layout_name=title_layout.name,
        slots=title_slots,
    )

    text_slots = _build_text_slots_from_placeholders(
        body_text_placeholders,
        layout_index=body_layout_index,
        slot_prefix="text.body_",
        alias_prefix="body_",
        allowed_kinds=(
            SlotContentKind.PARAGRAPH_OR_BULLETS,
            SlotContentKind.METRIC_LIST,
            SlotContentKind.SHORT_FACT_OR_STATUS,
        ),
    )
    body_title_slot = _build_placeholder_slot(
        body_title_placeholder,
        layout_index=body_layout_index,
        slot_name="text.title",
        alias="title",
        slot_type="title",
        preferred_font_size=SECTION_TITLE_FONT_SIZE,
        min_font_size=SECTION_TITLE_MIN_FONT_SIZE,
        allowed_kinds=(SlotContentKind.TITLE,),
    )
    text_pattern = PatternProfile(
        pattern_id="text.template",
        kind="text",
        layout_index=body_layout_index,
        layout_name=body_layout.name,
        slots=(body_title_slot, *text_slots),
    )
    metrics_pattern = PatternProfile(
        pattern_id="metrics.template",
        kind="metrics",
        layout_index=body_layout_index,
        layout_name=body_layout.name,
        slots=(
            _clone_slot(body_title_slot, slot_name="metrics.title"),
            *(
                _clone_slot(
                    slot,
                    slot_name=slot.slot_name.replace("text.body_", "metrics.body_"),
                    alias=slot.alias,
                )
                for slot in text_slots
            ),
        ),
    )
    contents_pattern = PatternProfile(
        pattern_id="contents.template",
        kind="contents",
        layout_index=body_layout_index,
        layout_name=body_layout.name,
        slots=(
            _clone_slot(body_title_slot, slot_name="contents.title"),
            *(
                _clone_slot(
                    slot,
                    slot_name=slot.slot_name.replace("text.body_", "contents.body_"),
                    alias=slot.alias,
                )
                for slot in text_slots
            ),
        ),
    )

    slide_patterns: list[PatternProfile] = [text_pattern, metrics_pattern]
    mixed_pattern = _profile_text_image_pattern(
        presentation,
        template_path=template_path,
        preferred_skip_layout_index=title_layout_index,
    )
    if mixed_pattern is not None:
        slide_patterns.append(mixed_pattern)

    return TemplateProfile(
        template_name=(
            template_path.stem if template_path is not None else TEMPLATE_NAME
        ),
        template_id=_build_user_template_id(template_path, slide_patterns),
        template_label=(
            template_path.stem.replace("_", " ").title()
            if template_path is not None
            else "User Template"
        ),
        template_source="user_template",
        template_path=template_path,
        title_pattern=title_pattern,
        contents_pattern=contents_pattern,
        slide_patterns=tuple(slide_patterns),
    )


def export_template_contract(profile: TemplateProfile) -> TemplateContract:
    """Convert a profiled template into the public template-contract shape."""

    return TemplateContract(
        contract_version=TEMPLATE_CONTRACT_VERSION,
        template_id=profile.template_id,
        template_label=profile.template_label,
        template_source=profile.template_source,
        title_slide=_export_section_contract(profile.title_pattern),
        contents_slide=_export_section_contract(profile.contents_pattern),
        slide_patterns=tuple(
            _export_pattern_contract(pattern)
            for pattern in profile.slide_patterns
        ),
    )


def build_report_fill_plan(
    payload: ReportPayload,
    template_profile: TemplateProfile,
    *,
    image_refs: dict[str, Path] | None = None,
) -> tuple[FillPlan, DiagnosticReport]:
    """Map a validated payload into concrete slides for the current template."""

    image_refs = image_refs or {}
    fill_plan = FillPlan()
    diagnostics = DiagnosticReport()
    _add_font_risk_warnings(diagnostics, template_profile)

    fill_plan.slides.append(
        _plan_title_slide(
            payload.title_slide.title,
            payload.title_slide.subtitle,
            payload.title_slide.slot_values,
            template_profile.title_pattern,
            diagnostics,
        )
    )

    included_titles = [
        slide.title
        for slide in payload.slides
        if slide.include_in_contents
    ]
    if payload.contents.enabled:
        fill_plan.slides.extend(
            _plan_textual_pattern_slides(
                slide_title="Contents",
                items=included_titles or ["Overview"],
                pattern=template_profile.contents_pattern,
                diagnostics=diagnostics,
                kind="contents",
                slot_overrides=_build_named_slot_overrides(
                    pattern=template_profile.contents_pattern,
                    slot_values=payload.contents.slot_values,
                    skip_aliases=_pattern_body_aliases(template_profile.contents_pattern),
                ),
            )
        )

    for payload_slide in payload.slides:
        pattern = _select_pattern_for_payload_slide(
            template_profile,
            payload_slide,
        )
        fill_plan.slides.extend(
            _plan_payload_slide(
                payload_slide,
                pattern=pattern,
                diagnostics=diagnostics,
                image_refs=image_refs,
            )
        )

    return fill_plan, diagnostics


def _plan_title_slide(
    title: str,
    subtitle_items: list[str],
    slot_values: dict[str, str],
    pattern: PatternProfile,
    diagnostics: DiagnosticReport,
) -> PlannedSlide:
    title_slot = pattern.slots_by_type("title")[0]
    title_text = slot_values.get(title_slot.alias or title_slot.slot_name, title)
    title_fit = fit_text_to_slot(title_text, title_slot)
    _record_fit_diagnostics(
        diagnostics,
        slide_title=title_text,
        fit_result=title_fit,
        label="cover title",
    )

    subtitle_fills: list[PlannedTextFill] = []
    remaining_subtitle_items = list(subtitle_items)
    remaining_subtitle_slots: list[SlotDescriptor] = []
    for slot in pattern.slots_by_type("text"):
        alias = slot.alias or slot.slot_name
        if alias not in slot_values:
            remaining_subtitle_slots.append(slot)
            continue
        fill = _build_named_text_fill(slot, slot_values[alias])
        subtitle_fills.append(fill)
        if fill.fit_result is not None:
            _record_fit_diagnostics(
                diagnostics,
                slide_title=title_text,
                fit_result=fill.fit_result,
                label=alias,
            )
    planned_subtitle_fills, _ = _plan_items_across_slots(
        remaining_subtitle_items,
        remaining_subtitle_slots,
        reserve_one_item_per_remaining_slot=True,
    )
    subtitle_fills.extend(planned_subtitle_fills)
    for index, fill in enumerate(subtitle_fills, start=1):
        if fill.fit_result is not None:
            _record_fit_diagnostics(
                diagnostics,
                slide_title=title_text,
                fit_result=fill.fit_result,
                label=f"cover subtitle slot {index}",
            )

    return PlannedSlide(
        pattern_id=pattern.pattern_id,
        kind=pattern.kind,
        layout_name=pattern.layout_name,
        layout_index=pattern.layout_index,
        slide_title=title_text,
        decorations=list(pattern.decorations),
        text_fills=[
            PlannedTextFill(
                slot=title_slot,
                text=title_text,
                font_size=title_fit.font_size,
                fit_result=title_fit,
            ),
            *subtitle_fills,
        ],
    )


def _plan_payload_slide(
    payload_slide: PayloadSlide,
    *,
    pattern: PatternProfile,
    diagnostics: DiagnosticReport,
    image_refs: dict[str, Path],
) -> list[PlannedSlide]:
    if payload_slide.kind == "metrics":
        items = [metric.as_text() for metric in payload_slide.items]
        return _plan_textual_pattern_slides(
            slide_title=payload_slide.title,
            items=items,
            pattern=pattern,
            diagnostics=diagnostics,
            kind=payload_slide.kind,
            slot_overrides=payload_slide.slot_overrides,
        )

    if payload_slide.kind == "text":
        return _plan_textual_pattern_slides(
            slide_title=payload_slide.title,
            items=payload_slide.body,
            pattern=pattern,
            diagnostics=diagnostics,
            kind=payload_slide.kind,
            slot_overrides=payload_slide.slot_overrides,
        )

    return _plan_text_image_slides(
        payload_slide,
        pattern=pattern,
        diagnostics=diagnostics,
        image_refs=image_refs,
    )


def _plan_textual_pattern_slides(
    *,
    slide_title: str,
    items: list[str],
    pattern: PatternProfile,
    diagnostics: DiagnosticReport,
    kind: str,
    slot_overrides: dict[str, object] | None = None,
) -> list[PlannedSlide]:
    slot_overrides = slot_overrides or {}
    title_slot = pattern.slots_by_type("title")[0]
    title_override = slot_overrides.get(title_slot.slot_name)
    display_title = slide_title
    if title_override is not None and getattr(title_override, "text", None):
        display_title = "\n".join(title_override.text)

    body_slots = list(pattern.slots_by_type("text"))
    override_fills: list[PlannedTextFill] = []
    remaining_body_slots: list[SlotDescriptor] = []
    for slot in body_slots:
        override = slot_overrides.get(slot.slot_name)
        if override is None:
            remaining_body_slots.append(slot)
            continue
        override_fill = _build_text_override_fill(slot, override)
        override_fills.append(override_fill)
        if override_fill.fit_result is not None:
            _record_fit_diagnostics(
                diagnostics,
                slide_title=slide_title,
                fit_result=override_fill.fit_result,
                label=slot.alias or slot.slot_name,
            )

    planned_slides: list[PlannedSlide] = []
    remaining_items = list(items)
    continuation_index = 0

    while remaining_items or continuation_index == 0:
        continuation = continuation_index > 0
        resolved_slide_title = (
            f"{slide_title} (cont.)"
            if continuation
            else slide_title
        )
        resolved_display_title = (
            f"{display_title} (cont.)"
            if continuation
            else display_title
        )
        title_fit = fit_text_to_slot(resolved_display_title, title_slot)
        _record_fit_diagnostics(
            diagnostics,
            slide_title=resolved_slide_title,
            fit_result=title_fit,
            label="title",
        )

        body_fills: list[PlannedTextFill] = list(override_fills)
        if remaining_items and remaining_body_slots:
            planned_body_fills, remaining_items = _plan_items_across_slots(
                remaining_items,
                remaining_body_slots,
                reserve_one_item_per_remaining_slot=True,
            )
            for fill in planned_body_fills:
                if fill.fit_result is not None:
                    _record_fit_diagnostics(
                        diagnostics,
                        slide_title=resolved_slide_title,
                        fit_result=fill.fit_result,
                        label=fill.slot.alias or fill.slot.slot_name,
                    )
            body_fills.extend(planned_body_fills)
        elif continuation or (remaining_items and not remaining_body_slots):
            break

        planned_slides.append(
            PlannedSlide(
                pattern_id=pattern.pattern_id,
                kind=kind,
                layout_name=pattern.layout_name,
                layout_index=pattern.layout_index,
                slide_title=resolved_slide_title,
                continuation=continuation,
                decorations=list(pattern.decorations),
                text_fills=[
                    PlannedTextFill(
                        slot=title_slot,
                        text=resolved_display_title,
                        font_size=title_fit.font_size,
                        fit_result=title_fit,
                    ),
                    *body_fills,
                ],
            )
        )

        if not remaining_items:
            break
        continuation_index += 1

    return planned_slides


def _plan_text_image_slides(
    payload_slide: PayloadSlide,
    *,
    pattern: PatternProfile,
    diagnostics: DiagnosticReport,
    image_refs: dict[str, Path],
) -> list[PlannedSlide]:
    caption_slots = sort_slots_in_reading_order(list(pattern.slots_by_type("caption")))
    image_slots = sort_slots_in_reading_order(list(pattern.slots_by_type("image")))
    slot_overrides = payload_slide.slot_overrides

    image_fills: list[PlannedImageFill] = []
    for image_slot in image_slots:
        image_override = slot_overrides.get(image_slot.slot_name)
        resolved_image = None
        if image_override is not None and image_override.image is not None:
            resolved_image = image_override.image
        elif image_slot.order == 1 and payload_slide.image is not None:
            resolved_image = payload_slide.image
        if resolved_image is not None:
            image_fills.append(
                PlannedImageFill(
                    slot=image_slot,
                    image_path=_resolve_image_path(resolved_image, image_refs),
                    fit=resolved_image.fit,
                )
            )

    caption_fills: list[PlannedTextFill] = []
    for caption_slot in caption_slots:
        caption_override = slot_overrides.get(caption_slot.slot_name)
        caption_values = None
        if caption_override is not None and caption_override.text is not None:
            caption_values = caption_override.text
        elif caption_slot.order == 1 and payload_slide.caption:
            caption_values = [payload_slide.caption]
        if caption_values:
            caption_text = "\n".join(caption_values)
            caption_fit = fit_text_to_slot(caption_text, caption_slot)
            caption_fills.append(
                PlannedTextFill(
                    slot=caption_slot,
                    text=caption_text,
                    font_size=caption_fit.font_size,
                    fit_result=caption_fit,
                )
            )
            _record_fit_diagnostics(
                diagnostics,
                slide_title=payload_slide.title,
                fit_result=caption_fit,
                label=caption_slot.alias or caption_slot.slot_name,
            )

    slides = _plan_textual_pattern_slides(
        slide_title=payload_slide.title,
        items=payload_slide.body,
        pattern=pattern,
        diagnostics=diagnostics,
        kind=payload_slide.kind,
        slot_overrides=slot_overrides,
    )
    if slides:
        # Keep the image-caption pair anchored to the primary slide so
        # continuation slides only show the remaining narrative content.
        primary_slide = slides[0]
        primary_slide.image_fills.extend(image_fills)
        primary_slide.text_fills.extend(caption_fills)
    return slides


def _select_pattern_for_payload_slide(
    template_profile: TemplateProfile,
    payload_slide: PayloadSlide,
) -> PatternProfile:
    if payload_slide.pattern_id is not None:
        pattern = template_profile.get_pattern(payload_slide.pattern_id)
        if pattern is None:
            raise ValueError(f"Unknown pattern id: {payload_slide.pattern_id}")
        return pattern

    candidates = template_profile.patterns_for_kind(payload_slide.kind)
    if not candidates:
        raise ValueError(
            f"Template '{template_profile.template_id}' does not support "
            f"slide kind '{payload_slide.kind}'."
        )
    return candidates[0]


def _resolve_image_path(image: ImageSpec, image_refs: dict[str, Path]) -> Path:
    if image.path is not None:
        return image.path
    if image.ref is not None and image.ref in image_refs:
        return image_refs[image.ref]
    raise ValueError("Image payload could not be resolved to a file path.")


def _build_text_override_fill(slot: SlotDescriptor, override) -> PlannedTextFill:
    values = list(override.text or [])
    if slot.slot_type in {"title", "caption"}:
        text = "\n".join(values)
        fit_result = fit_text_to_slot(text, slot)
        return PlannedTextFill(
            slot=slot,
            text=text,
            font_size=fit_result.font_size,
            fit_result=fit_result,
        )

    fit_result = fit_text_items_to_slot(values, slot)
    return PlannedTextFill(
        slot=slot,
        items=values[: fit_result.consumed_items],
        font_size=fit_result.font_size,
        fit_result=fit_result,
    )


def _build_named_text_fill(slot: SlotDescriptor, value: str) -> PlannedTextFill:
    if slot.slot_type in {"title", "caption"}:
        fit_result = fit_text_to_slot(value, slot)
        return PlannedTextFill(
            slot=slot,
            text=value,
            font_size=fit_result.font_size,
            fit_result=fit_result,
        )
    fit_result = fit_text_items_to_slot([value], slot)
    return PlannedTextFill(
        slot=slot,
        items=[value][: fit_result.consumed_items],
        font_size=fit_result.font_size,
        fit_result=fit_result,
    )


def _build_named_slot_overrides(
    *,
    pattern: PatternProfile,
    slot_values: dict[str, str],
    skip_aliases: set[str] | None = None,
) -> dict[str, SlotOverride]:
    skip_aliases = set() if skip_aliases is None else set(skip_aliases)
    overrides: dict[str, SlotOverride] = {}
    for slot in pattern.slots:
        if slot.slot_type == "image":
            continue
        alias = slot.alias or slot.slot_name
        if alias in skip_aliases or alias not in slot_values:
            continue
        overrides[slot.slot_name] = SlotOverride(
            slot_id=slot.slot_name,
            text=[slot_values[alias]],
        )
    return overrides


def _plan_items_across_slots(
    items: list[str],
    slots: tuple[SlotDescriptor, ...] | list[SlotDescriptor],
    *,
    reserve_one_item_per_remaining_slot: bool,
) -> tuple[list[PlannedTextFill], list[str]]:
    fills: list[PlannedTextFill] = []
    remaining_items = list(items)
    ordered_slots = sort_slots_in_reading_order(list(slots))

    for index, slot in enumerate(ordered_slots):
        if not remaining_items:
            break

        remaining_slot_count = len(ordered_slots) - index - 1
        max_items_for_slot = len(remaining_items)
        if reserve_one_item_per_remaining_slot and remaining_slot_count > 0:
            max_items_for_slot = max(1, len(remaining_items) - remaining_slot_count)
        candidate_items = remaining_items[:max_items_for_slot]
        fit_result = fit_text_items_to_slot(candidate_items, slot)
        consumed_items = candidate_items[: fit_result.consumed_items]
        fills.append(
            PlannedTextFill(
                slot=slot,
                items=consumed_items,
                font_size=fit_result.font_size,
                fit_result=fit_result,
            )
        )
        remaining_items = remaining_items[len(consumed_items) :]

    return fills, remaining_items


def _export_section_contract(pattern: PatternProfile) -> TemplateSectionContract:
    return TemplateSectionContract(
        pattern_id=pattern.pattern_id,
        layout_name=pattern.layout_name,
        slots=tuple(_export_slot_contract(slot) for slot in pattern.slots),
    )


def _export_pattern_contract(pattern: PatternProfile) -> TemplatePatternContract:
    return TemplatePatternContract(
        pattern_id=pattern.pattern_id,
        kind=pattern.kind,
        layout_name=pattern.layout_name,
        slots=tuple(_export_slot_contract(slot) for slot in pattern.slots),
        image_count=len(pattern.slots_by_type("image")),
        image_layout=_derive_pattern_image_layout(pattern),
        caption_slots=len(pattern.slots_by_type("caption")),
        body_slot_count=len(_pattern_body_slots(pattern)),
    )


def _derive_pattern_image_layout(pattern: PatternProfile) -> str:
    image_slots = list(pattern.slots_by_type("image"))
    if not image_slots:
        return "stack"
    return image_slots[0].orientation or "stack"


def _pattern_body_slots(pattern: PatternProfile) -> tuple[SlotDescriptor, ...]:
    return tuple(
        slot
        for slot in pattern.slots_by_type("text")
        if slot.slot_name.startswith(f"{pattern.kind}.body_")
    )


def _pattern_body_aliases(pattern: PatternProfile) -> set[str]:
    return {
        slot.alias or slot.slot_name
        for slot in _pattern_body_slots(pattern)
    }


def _export_slot_contract(slot: SlotDescriptor) -> TemplateSlotContract:
    return TemplateSlotContract(
        slot_id=slot.slot_name,
        alias=slot.alias or slot.slot_name,
        slot_type=slot.slot_type,
        required=slot.required,
        orientation=slot.orientation,
        order=slot.order,
    )


def _build_user_template_id(
    template_path: Path | None,
    patterns: list[PatternProfile],
) -> str:
    digest = sha1()
    if template_path is not None and template_path.exists():
        digest.update(template_path.read_bytes())
    else:
        digest.update(b"default-template")
    for pattern in patterns:
        digest.update(pattern.pattern_id.encode("utf-8"))
        digest.update(pattern.layout_name.encode("utf-8"))
        for slot in pattern.slots:
            digest.update(slot.slot_name.encode("utf-8"))
            digest.update((slot.slot_type or "").encode("utf-8"))
            digest.update(str(slot.placeholder_index).encode("utf-8"))
    return f"template-{digest.hexdigest()[:12]}"


def _build_text_slots_from_placeholders(
    placeholders,
    *,
    layout_index: int,
    slot_prefix: str,
    alias_prefix: str,
    allowed_kinds: tuple[SlotContentKind, ...],
) -> tuple[SlotDescriptor, ...]:
    provisional = [
        _build_placeholder_slot(
            placeholder,
            layout_index=layout_index,
            slot_name="__tmp__",
            alias="__tmp__",
            slot_type="text",
            preferred_font_size=SECTION_BODY_FONT_SIZE,
            min_font_size=SECTION_BODY_MIN_FONT_SIZE,
            allowed_kinds=allowed_kinds,
        )
        for placeholder in placeholders
    ]
    ordered_placeholders = sort_slots_in_reading_order(provisional)
    orientation = _infer_orientation(ordered_placeholders)
    slots: list[SlotDescriptor] = []
    for index, slot in enumerate(ordered_placeholders, start=1):
        slots.append(
            _clone_slot(
                slot,
                slot_name=f"{slot_prefix}{index}",
                alias=f"{alias_prefix}{index}",
                orientation=orientation,
                order=index,
            )
        )
    return tuple(slots)


def _profile_text_image_pattern(
    presentation: Presentation,
    *,
    template_path: Path | None,
    preferred_skip_layout_index: int,
) -> PatternProfile | None:
    best_candidate = None
    for layout_index, layout in enumerate(presentation.slide_layouts):
        if layout_index == preferred_skip_layout_index:
            continue
        try:
            title_placeholder = _pick_body_title_placeholder(
                layout,
                layout_name="text_image",
                template_path=template_path,
            )
            text_placeholders = _list_body_text_placeholders(
                layout,
                excluded_idx=title_placeholder.placeholder_format.idx,
            )
            image_placeholders = _list_image_placeholders(layout)
            if not text_placeholders or not image_placeholders:
                continue
            candidate = (
                len(image_placeholders),
                len(text_placeholders),
                sum(
                    placeholder.width * placeholder.height
                    for placeholder in image_placeholders
                ),
                -layout_index,
                layout_index,
                layout,
                title_placeholder,
                text_placeholders,
                image_placeholders,
            )
            if best_candidate is None or candidate > best_candidate:
                best_candidate = candidate
        except TemplateCompatibilityError:
            continue

    if best_candidate is None:
        return None

    (
        _,
        _,
        _,
        _,
        layout_index,
        layout,
        title_placeholder,
        text_placeholders,
        image_placeholders,
    ) = best_candidate

    ordered_text_slots = _build_text_slots_from_placeholders(
        text_placeholders[:1],
        layout_index=layout_index,
        slot_prefix="text_image.body_",
        alias_prefix="body_",
        allowed_kinds=(SlotContentKind.PARAGRAPH_OR_BULLETS,),
    )
    caption_slots: tuple[SlotDescriptor, ...] = ()
    if len(text_placeholders) > 1:
        caption_slots = _build_text_slots_from_placeholders(
            text_placeholders[1:2],
            layout_index=layout_index,
            slot_prefix="text_image.caption_",
            alias_prefix="caption_",
            allowed_kinds=(SlotContentKind.SHORT_FACT_OR_STATUS,),
        )
        caption_slots = tuple(
            _clone_slot(slot, slot_type="caption", required=False)
            for slot in caption_slots
        )

    provisional_image_slots = [
        _build_placeholder_slot(
            image_placeholder,
            layout_index=layout_index,
            slot_name=f"__image_{index}",
            alias=f"__image_{index}",
            slot_type="image",
            preferred_font_size=SECTION_BODY_FONT_SIZE,
            min_font_size=SECTION_BODY_MIN_FONT_SIZE,
            allowed_kinds=(),
        )
        for index, image_placeholder in enumerate(image_placeholders)
    ]
    ordered_image_slots = sort_slots_in_reading_order(provisional_image_slots)
    image_orientation = _infer_orientation(ordered_image_slots)
    image_slots = tuple(
        _clone_slot(
            slot,
            slot_name=f"text_image.image_{index + 1}",
            alias=f"image_{index + 1}",
            slot_type="image",
            orientation=image_orientation,
            order=index + 1,
        )
        for index, slot in enumerate(ordered_image_slots)
    )

    return PatternProfile(
        pattern_id="text_image.template",
        kind="text_image",
        layout_index=layout_index,
        layout_name=layout.name,
        slots=(
            _build_placeholder_slot(
                title_placeholder,
                layout_index=layout_index,
                slot_name="text_image.title",
                alias="title",
                slot_type="title",
                preferred_font_size=SECTION_TITLE_FONT_SIZE,
                min_font_size=SECTION_TITLE_MIN_FONT_SIZE,
                allowed_kinds=(SlotContentKind.TITLE,),
            ),
            *ordered_text_slots,
            *image_slots,
            *caption_slots,
        ),
    )


def _get_layout(
    presentation: Presentation,
    *,
    layout_index: int,
    layout_name: str,
    template_path: Path | None,
):
    try:
        return presentation.slide_layouts[layout_index]
    except IndexError as exc:
        raise TemplateCompatibilityError(
            template_path,
            f"missing '{layout_name}' slide layout at index {layout_index}",
        ) from exc


def _select_title_layout(
    presentation: Presentation,
    *,
    template_path: Path | None,
):
    best_candidate = None
    for layout_index, layout in enumerate(presentation.slide_layouts):
        try:
            title_placeholder = _pick_title_text_placeholder(
                layout,
                layout_name="title",
                template_path=template_path,
            )
            subtitle_placeholders = _pick_title_secondary_text_placeholders(
                layout,
                excluded_idx=title_placeholder.placeholder_format.idx,
                anchor_top=title_placeholder.top,
                layout_name="title",
                template_path=template_path,
            )
            candidate = (
                _score_title_layout_candidate(
                    title_placeholder=title_placeholder,
                    secondary_placeholders=subtitle_placeholders,
                ),
                layout_index,
                layout,
                title_placeholder,
                subtitle_placeholders,
            )
            if best_candidate is None or candidate > best_candidate:
                best_candidate = candidate
        except TemplateCompatibilityError:
            continue

    if best_candidate is None:
        raise TemplateCompatibilityError(
            template_path,
            "no compatible title layout exposing both title and subtitle placeholders",
        )

    _, layout_index, layout, title_placeholder, subtitle_placeholders = best_candidate
    return layout_index, layout, title_placeholder, subtitle_placeholders


def _select_body_layout(
    presentation: Presentation,
    *,
    template_path: Path | None,
    preferred_skip_layout_index: int | None = None,
):
    candidate_indices = [
        index
        for index, _ in enumerate(presentation.slide_layouts)
        if index != preferred_skip_layout_index
    ]
    best_candidate = None

    for layout_index in candidate_indices:
        layout = presentation.slide_layouts[layout_index]
        try:
            body_title_placeholder = _pick_body_title_placeholder(
                layout,
                layout_name="body",
                template_path=template_path,
            )
            body_content_placeholders = _pick_body_content_placeholders(
                layout,
                excluded_idx=body_title_placeholder.placeholder_format.idx,
                layout_name="body",
                template_path=template_path,
            )
            candidate = (
                len(body_content_placeholders),
                sum(
                    placeholder.width * placeholder.height
                    for placeholder in body_content_placeholders
                ),
                -layout_index,
                layout_index,
                layout,
                body_title_placeholder,
                body_content_placeholders,
            )
            if best_candidate is None or candidate > best_candidate:
                best_candidate = candidate
        except TemplateCompatibilityError:
            continue

    if best_candidate is None:
        raise TemplateCompatibilityError(
            template_path,
            "no compatible body layout exposing both title and text placeholders",
        )

    _, _, _, layout_index, layout, title_placeholder, content_placeholders = best_candidate
    return layout_index, layout, title_placeholder, content_placeholders


def _find_title_placeholder(layout):
    for placeholder in getattr(layout, "placeholders", []):
        placeholder_type = placeholder.placeholder_format.type
        if (
            placeholder_type in TITLE_PLACEHOLDER_TYPES
            and hasattr(placeholder, "text_frame")
        ):
            return placeholder
    return None


def _pick_title_text_placeholder(
    layout,
    *,
    layout_name: str,
    template_path: Path | None,
):
    placeholder = _find_title_placeholder(layout)
    if placeholder is not None:
        return placeholder

    candidates = _list_text_placeholders(layout)
    if not candidates:
        raise TemplateCompatibilityError(
            template_path,
            f"'{layout_name}' layout has no text placeholder",
        )
    return max(candidates, key=lambda shape: shape.width * shape.height)


def _pick_body_title_placeholder(
    layout,
    *,
    layout_name: str,
    template_path: Path | None,
):
    placeholder = _find_title_placeholder(layout)
    if placeholder is not None:
        return placeholder

    candidates = _list_text_placeholders(layout)
    if len(candidates) < 2:
        raise TemplateCompatibilityError(
            template_path,
            f"'{layout_name}' layout has no title-like text placeholder",
        )
    return min(candidates, key=lambda shape: (shape.top, shape.left))


def _pick_title_secondary_text_placeholders(
    layout,
    *,
    excluded_idx: int,
    anchor_top: int,
    layout_name: str,
    template_path: Path | None,
):
    candidates = _list_text_placeholders(layout, excluded_idx=excluded_idx)
    if not candidates:
        raise TemplateCompatibilityError(
            template_path,
            f"'{layout_name}' layout has no secondary text placeholder",
        )

    lower_candidates = [
        shape
        for shape in candidates
        if shape.top >= anchor_top
    ]
    scoped_candidates = lower_candidates or candidates
    max_area = max(shape.width * shape.height for shape in scoped_candidates)
    prominent_candidates = [
        shape
        for shape in scoped_candidates
        if (shape.width * shape.height) >= (max_area * 0.45)
    ]
    return tuple(
        sorted(prominent_candidates, key=lambda shape: (shape.top, shape.left))
    )


def _pick_body_content_placeholders(
    layout,
    *,
    excluded_idx: int,
    layout_name: str,
    template_path: Path | None,
):
    candidates = _list_body_text_placeholders(layout, excluded_idx=excluded_idx)
    if not candidates:
        raise TemplateCompatibilityError(
            template_path,
            f"'{layout_name}' layout has no secondary text placeholder",
        )

    max_area = max(shape.width * shape.height for shape in candidates)
    prominent_candidates = [
        shape
        for shape in candidates
        if (shape.width * shape.height) >= (max_area * 0.6)
    ] or [max(candidates, key=lambda shape: shape.width * shape.height)]
    return tuple(
        sorted(prominent_candidates, key=lambda shape: (shape.top, shape.left))
    )


def _list_text_placeholders(layout, *, excluded_idx: int | None = None):
    return [
        shape
        for shape in layout.placeholders
        if hasattr(shape, "text_frame")
        and shape.placeholder_format.type not in IGNORED_TEXT_PLACEHOLDER_TYPES
        and (
            excluded_idx is None
            or shape.placeholder_format.idx != excluded_idx
        )
    ]


def _list_body_text_placeholders(layout, *, excluded_idx: int | None = None):
    return [
        shape
        for shape in layout.placeholders
        if hasattr(shape, "text_frame")
        and shape.placeholder_format.type not in TEXT_EXCLUDED_PLACEHOLDER_TYPES
        and (
            excluded_idx is None
            or shape.placeholder_format.idx != excluded_idx
        )
    ]


def _list_image_placeholders(layout):
    return [
        shape
        for shape in layout.placeholders
        if shape.placeholder_format.type in PICTURE_PLACEHOLDER_TYPES
    ]


def _score_title_layout_candidate(
    *,
    title_placeholder,
    secondary_placeholders,
) -> tuple[int, int, int, int, int]:
    title_area = title_placeholder.width * title_placeholder.height
    max_secondary_area = max(
        shape.width * shape.height for shape in secondary_placeholders
    )
    has_explicit_subtitle = any(
        getattr(PP_PLACEHOLDER, "SUBTITLE", None)
        == shape.placeholder_format.type
        for shape in secondary_placeholders
    )
    has_explicit_title = title_placeholder.placeholder_format.type in TITLE_PLACEHOLDER_TYPES
    return (
        int(has_explicit_subtitle),
        title_area - max_secondary_area,
        int(has_explicit_title),
        title_area,
        -len(secondary_placeholders),
    )


def _build_placeholder_slot(
    placeholder,
    *,
    layout_index: int,
    slot_name: str,
    alias: str,
    slot_type: str,
    preferred_font_size: int,
    min_font_size: int,
    allowed_kinds: tuple[SlotContentKind, ...],
    orientation: str | None = None,
    order: int | None = None,
    required: bool = True,
) -> SlotDescriptor:
    paragraph = (
        placeholder.text_frame.paragraphs[0]
        if hasattr(placeholder, "text_frame")
        else None
    )
    return SlotDescriptor(
        slot_name=slot_name,
        alias=alias,
        slot_type=slot_type,
        required=required,
        orientation=orientation,
        order=order,
        layout_index=layout_index,
        placeholder_index=placeholder.placeholder_format.idx,
        x=placeholder.left,
        y=placeholder.top,
        width=placeholder.width,
        height=placeholder.height,
        preferred_font_size=preferred_font_size,
        min_font_size=min_font_size,
        allowed_kinds=allowed_kinds,
        explicit_font_name=(None if paragraph is None else paragraph.font.name),
        placeholder_type=str(placeholder.placeholder_format.type),
    )


def _build_text_box_slot(
    presentation: Presentation,
    *,
    layout_index: int,
    slot_name: str,
    alias: str,
    slot_type: str,
    x_ratio: float,
    y_ratio: float,
    width_ratio: float,
    height_ratio: float,
    preferred_font_size: int,
    min_font_size: int,
    allowed_kinds: tuple[SlotContentKind, ...],
    orientation: str | None = None,
    order: int | None = None,
    required: bool = True,
) -> SlotDescriptor:
    return SlotDescriptor(
        slot_name=slot_name,
        alias=alias,
        slot_type=slot_type,
        required=required,
        orientation=orientation,
        order=order,
        layout_index=layout_index,
        placeholder_index=None,
        x=int(presentation.slide_width * x_ratio),
        y=int(presentation.slide_height * y_ratio),
        width=int(presentation.slide_width * width_ratio),
        height=int(presentation.slide_height * height_ratio),
        preferred_font_size=preferred_font_size,
        min_font_size=min_font_size,
        allowed_kinds=allowed_kinds,
    )


def _build_geometry_slot(
    presentation: Presentation,
    *,
    layout_index: int,
    slot_name: str,
    alias: str,
    slot_type: str,
    x_ratio: float,
    y_ratio: float,
    width_ratio: float,
    height_ratio: float,
    orientation: str | None = None,
    order: int | None = None,
    required: bool = True,
) -> SlotDescriptor:
    return SlotDescriptor(
        slot_name=slot_name,
        alias=alias,
        slot_type=slot_type,
        required=required,
        orientation=orientation,
        order=order,
        layout_index=layout_index,
        placeholder_index=None,
        x=int(presentation.slide_width * x_ratio),
        y=int(presentation.slide_height * y_ratio),
        width=int(presentation.slide_width * width_ratio),
        height=int(presentation.slide_height * height_ratio),
        placeholder_type="geometry",
    )


def _clone_slot(
    slot: SlotDescriptor,
    *,
    slot_name: str,
    alias: str | None = None,
    slot_type: str | None = None,
    orientation: str | None = None,
    order: int | None = None,
    required: bool | None = None,
) -> SlotDescriptor:
    return SlotDescriptor(
        slot_name=slot_name,
        alias=(slot.alias if alias is None else alias),
        slot_type=(slot.slot_type if slot_type is None else slot_type),
        required=(slot.required if required is None else required),
        orientation=(slot.orientation if orientation is None else orientation),
        order=(slot.order if order is None else order),
        layout_index=slot.layout_index,
        placeholder_index=slot.placeholder_index,
        x=slot.x,
        y=slot.y,
        width=slot.width,
        height=slot.height,
        preferred_font_size=slot.preferred_font_size,
        min_font_size=slot.min_font_size,
        allowed_kinds=slot.allowed_kinds,
        explicit_font_name=slot.explicit_font_name,
        priority=slot.priority,
        placeholder_type=slot.placeholder_type,
    )


def _infer_orientation(slots: list[SlotDescriptor] | tuple[SlotDescriptor, ...]) -> str:
    if len(slots) <= 1:
        return "stack"

    ordered = sort_slots_in_reading_order(list(slots))
    x_span = max(slot.x for slot in ordered) - min(slot.x for slot in ordered)
    y_span = max(slot.y for slot in ordered) - min(slot.y for slot in ordered)
    max_height = max(slot.height for slot in ordered)
    max_width = max(slot.width for slot in ordered)

    if y_span <= max_height // 4:
        return "horizontal"
    if x_span <= max_width // 4:
        return "stack"
    return "vertical"


def _build_manual_cover_decorations(
    presentation: Presentation,
) -> tuple[SlideDecoration, ...]:
    return (
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=1.000,
            height_ratio=1.000,
            fill_rgb=MANUAL_BG,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=0.028,
            height_ratio=1.000,
            fill_rgb=MANUAL_ACCENT,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.720,
            y_ratio=0.106,
            width_ratio=0.214,
            height_ratio=0.694,
            fill_rgb=MANUAL_ACCENT_SOFT,
            line_rgb=MANUAL_FRAME,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.083,
            y_ratio=0.664,
            width_ratio=0.470,
            height_ratio=0.004,
            fill_rgb=MANUAL_ACTION,
        ),
    )


def _build_manual_text_decorations(
    presentation: Presentation,
) -> tuple[SlideDecoration, ...]:
    return (
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=1.000,
            height_ratio=1.000,
            fill_rgb=MANUAL_BG,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=1.000,
            height_ratio=0.034,
            fill_rgb=MANUAL_ACCENT,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.064,
            y_ratio=0.158,
            width_ratio=0.862,
            height_ratio=0.648,
            fill_rgb=MANUAL_TEXT_PANEL,
            line_rgb=MANUAL_FRAME,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.083,
            y_ratio=0.176,
            width_ratio=0.126,
            height_ratio=0.004,
            fill_rgb=MANUAL_ACTION,
        ),
    )


def _build_manual_text_image_decorations(
    presentation: Presentation,
    *,
    body_region: tuple[float, float, float, float],
    image_regions: tuple[tuple[float, float, float, float], ...],
    caption_regions: tuple[tuple[float, float, float, float], ...],
) -> tuple[SlideDecoration, ...]:
    decorations: list[SlideDecoration] = [
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=1.000,
            height_ratio=1.000,
            fill_rgb=MANUAL_BG,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=1.000,
            height_ratio=0.034,
            fill_rgb=MANUAL_ACCENT,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.064,
            y_ratio=0.172,
            width_ratio=0.862,
            height_ratio=0.040,
            fill_rgb=MANUAL_ACCENT_SOFT,
            line_rgb=MANUAL_FRAME,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=body_region[0] - 0.012,
            y_ratio=body_region[1] - 0.020,
            width_ratio=body_region[2] + 0.024,
            height_ratio=body_region[3] + 0.040,
            fill_rgb=MANUAL_TEXT_PANEL,
            line_rgb=MANUAL_FRAME,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.083,
            y_ratio=0.176,
            width_ratio=0.126,
            height_ratio=0.004,
            fill_rgb=MANUAL_ACTION,
        ),
    ]
    for region in image_regions:
        decorations.append(
            _build_ratio_decoration(
                presentation,
                x_ratio=region[0],
                y_ratio=region[1],
                width_ratio=region[2],
                height_ratio=region[3],
                fill_rgb=MANUAL_IMAGE_PANEL,
                line_rgb=MANUAL_FRAME,
            )
        )
    for region in caption_regions:
        decorations.append(
            _build_ratio_decoration(
                presentation,
                x_ratio=region[0],
                y_ratio=region[1],
                width_ratio=region[2],
                height_ratio=region[3],
                fill_rgb=MANUAL_TEXT_PANEL,
                line_rgb=MANUAL_FRAME,
            )
        )
    return tuple(decorations)


def _build_basic_template_title_decorations(
    presentation: Presentation,
) -> tuple[SlideDecoration, ...]:
    return (
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=0.251,
            height_ratio=1.000,
            fill_rgb=BASIC_EDITORIAL_NAVY,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.821,
            y_ratio=0.000,
            width_ratio=0.179,
            height_ratio=1.000,
            fill_rgb=(243, 244, 246),
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.793,
            y_ratio=0.000,
            width_ratio=0.016,
            height_ratio=1.000,
            fill_rgb=BASIC_EDITORIAL_ORANGE,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.315,
            y_ratio=0.210,
            width_ratio=0.195,
            height_ratio=0.005,
            fill_rgb=BASIC_EDITORIAL_TEAL,
        ),
    )


def _build_basic_template_body_decorations(
    presentation: Presentation,
) -> tuple[SlideDecoration, ...]:
    return (
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=1.000,
            height_ratio=1.000,
            fill_rgb=BASIC_EDITORIAL_LIGHT,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.000,
            y_ratio=0.000,
            width_ratio=1.000,
            height_ratio=0.037,
            fill_rgb=BASIC_EDITORIAL_NAVY,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.046,
            y_ratio=0.093,
            width_ratio=0.012,
            height_ratio=0.793,
            fill_rgb=BASIC_EDITORIAL_TEAL,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.068,
            y_ratio=0.197,
            width_ratio=0.604,
            height_ratio=0.664,
            fill_rgb=(255, 255, 255),
            line_rgb=BASIC_EDITORIAL_LINE,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.709,
            y_ratio=0.197,
            width_ratio=0.227,
            height_ratio=0.173,
            fill_rgb=BASIC_EDITORIAL_PALE,
            line_rgb=BASIC_EDITORIAL_PALE_LINE,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.709,
            y_ratio=0.411,
            width_ratio=0.227,
            height_ratio=0.256,
            fill_rgb=BASIC_EDITORIAL_NAVY,
        ),
        _build_ratio_decoration(
            presentation,
            x_ratio=0.077,
            y_ratio=0.176,
            width_ratio=0.090,
            height_ratio=0.005,
            fill_rgb=BASIC_EDITORIAL_ORANGE,
        ),
    )


def _build_ratio_decoration(
    presentation: Presentation,
    *,
    x_ratio: float,
    y_ratio: float,
    width_ratio: float,
    height_ratio: float,
    fill_rgb: tuple[int, int, int],
    line_rgb: tuple[int, int, int] | None = None,
) -> SlideDecoration:
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    return SlideDecoration(
        shape_type="rectangle",
        x=round(slide_width * x_ratio),
        y=round(slide_height * y_ratio),
        width=round(slide_width * width_ratio),
        height=max(1, round(slide_height * height_ratio)),
        fill_rgb=fill_rgb,
        line_rgb=line_rgb,
    )


def _add_font_risk_warnings(
    diagnostics: DiagnosticReport,
    template_profile: TemplateProfile,
) -> None:
    if template_profile.template_path is None:
        return

    for pattern in (
        template_profile.title_pattern,
        template_profile.contents_pattern,
        *template_profile.slide_patterns,
    ):
        for slot in pattern.slots:
            if slot.slot_type == "image":
                continue
            if slot.explicit_font_name is None:
                diagnostics.add_warning(
                    "font-substitution-risk",
                    "Template slot "
                    f"'{slot.slot_name}' does not pin a font face; output may vary "
                    "across machines.",
                )


def _record_fit_diagnostics(
    diagnostics: DiagnosticReport,
    *,
    slide_title: str,
    fit_result,
    label: str,
) -> None:
    if fit_result.status == FitStatus.SHRINK:
        diagnostics.add_warning(
            "font-shrink",
            f"{label} was shrunk to {fit_result.font_size}pt to fit the template.",
            slide_title=slide_title,
        )
    if fit_result.status in (FitStatus.SPILL, FitStatus.OVERFLOW):
        diagnostics.add_warning(
            "overflow-spill",
            f"{label} exceeded the slot budget and continued onto another slide.",
            slide_title=slide_title,
        )
    if fit_result.out_of_bounds_risk:
        diagnostics.add_warning(
            "out-of-bounds-risk",
            f"{label} may still exceed the safe text box bounds at "
            f"{fit_result.font_size}pt.",
            slide_title=slide_title,
        )
