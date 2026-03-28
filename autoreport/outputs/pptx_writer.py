"""PowerPoint output support for rendered weekly reports."""

from __future__ import annotations

from pathlib import Path
from typing import Any
from zipfile import BadZipFile

from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.util import Pt

from autoreport.outputs.errors import (
    OutputWriteError,
    TemplateCompatibilityError,
    TemplateLoadError,
    TemplateNotFoundError,
    TemplateReadError,
)
from autoreport.templates.autofill import FillPlan, PlannedSlide
from autoreport.templates.weekly_report import profile_weekly_template


class PowerPointWriter:
    """Write report content into a `.pptx` presentation."""

    def write(
        self,
        *,
        template_path: Path | None,
        output_path: Path,
        context: dict[str, Any],
    ) -> Path:
        """Write a PowerPoint file using a template and prepared context."""

        presentation = self._load_presentation(template_path)
        self._ensure_weekly_template_compatibility(
            presentation,
            template_path=template_path,
        )
        self._clear_slides(presentation)

        for slide_context in context.get("slides", []):
            layout = slide_context["layout"]
            if layout == "title":
                self._add_title_slide(
                    presentation,
                    title=slide_context["title"],
                    subtitle=slide_context["subtitle"],
                )
                continue

            if layout == "bullets":
                self._add_bullet_slide(
                    presentation,
                    title=slide_context["title"],
                    items=slide_context["items"],
                )
                continue

            raise ValueError(f"Unsupported slide layout: {layout}")

        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            presentation.save(str(output_path))
        except OSError as exc:
            raise OutputWriteError(output_path) from exc

        return output_path

    def write_fill_plan(
        self,
        *,
        presentation: Presentation,
        output_path: Path,
        fill_plan: FillPlan,
    ) -> Path:
        """Write a presentation from a template-aware slide plan."""

        self._clear_slides(presentation)
        for planned_slide in fill_plan.slides:
            self._write_planned_slide(presentation, planned_slide)

        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            presentation.save(str(output_path))
        except OSError as exc:
            raise OutputWriteError(output_path) from exc

        return output_path

    def _load_presentation(self, template_path: Path | None) -> Presentation:
        """Load an existing template or fall back to the default presentation."""

        if template_path is None:
            return Presentation()

        if not template_path.exists() or not template_path.is_file():
            raise TemplateNotFoundError(template_path)

        try:
            return Presentation(str(template_path))
        except (BadZipFile, KeyError, PackageNotFoundError, ValueError) as exc:
            raise TemplateLoadError(template_path) from exc
        except OSError as exc:
            raise TemplateReadError(template_path) from exc

    def _ensure_weekly_template_compatibility(
        self,
        presentation: Presentation,
        *,
        template_path: Path | None,
    ) -> None:
        """Verify the template exposes the layouts/placeholders weekly reports use."""

        profile_weekly_template(
            presentation,
            template_path=template_path,
        )

    def _assert_layout_access(
        self,
        presentation: Presentation,
        *,
        layout_index: int,
        placeholder_index: int,
        layout_name: str,
        template_path: Path | None,
    ) -> None:
        """Check that the expected layout and placeholder can be instantiated."""

        try:
            layout = presentation.slide_layouts[layout_index]
        except IndexError as exc:
            raise TemplateCompatibilityError(
                template_path,
                f"missing '{layout_name}' slide layout at index {layout_index}",
            ) from exc

        try:
            slide = presentation.slides.add_slide(layout)
            if slide.shapes.title is None:
                raise TemplateCompatibilityError(
                    template_path,
                    f"'{layout_name}' layout has no title placeholder",
                )
            placeholder = slide.placeholders[placeholder_index]
            if not hasattr(placeholder, "text_frame"):
                raise TemplateCompatibilityError(
                    template_path,
                    f"'{layout_name}' layout placeholder {placeholder_index} "
                    f"does not support text",
                )
        except TemplateCompatibilityError:
            raise
        except (AttributeError, IndexError, KeyError) as exc:
            raise TemplateCompatibilityError(
                template_path,
                f"'{layout_name}' layout is missing placeholder {placeholder_index}",
            ) from exc

    def _clear_slides(self, presentation: Presentation) -> None:
        """Remove existing slides while preserving the template theme."""

        slide_id_list = list(presentation.slides._sldIdLst)
        for slide_id in slide_id_list:
            relationship_id = slide_id.rId
            presentation.part.drop_rel(relationship_id)
            presentation.slides._sldIdLst.remove(slide_id)

    def _add_title_slide(
        self,
        presentation: Presentation,
        *,
        title: str,
        subtitle: str,
    ) -> None:
        """Add the opening title slide."""

        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def _add_bullet_slide(
        self,
        presentation: Presentation,
        *,
        title: str,
        items: list[str],
    ) -> None:
        """Add a content slide with bullet items."""

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for index, item in enumerate(items):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = item

    def _write_planned_slide(
        self,
        presentation: Presentation,
        planned_slide: PlannedSlide,
    ) -> None:
        """Render one planned slide into the target presentation."""

        slide = presentation.slides.add_slide(
            presentation.slide_layouts[planned_slide.layout_index]
        )
        title_placeholder = slide.placeholders[
            planned_slide.title_slot.placeholder_index
        ]
        title_placeholder.text = planned_slide.title_text
        if planned_slide.title_font_size is not None:
            self._apply_font_size(
                title_placeholder.text_frame,
                planned_slide.title_font_size,
            )

        if planned_slide.body_slot is None:
            return

        body_placeholder = slide.placeholders[
            planned_slide.body_slot.placeholder_index
        ]
        if planned_slide.subtitle_text is not None:
            body_placeholder.text = planned_slide.subtitle_text
            if planned_slide.body_font_size is not None:
                self._apply_font_size(
                    body_placeholder.text_frame,
                    planned_slide.body_font_size,
                )
            return

        text_frame = body_placeholder.text_frame
        text_frame.clear()
        for index, item in enumerate(planned_slide.body_items):
            paragraph = (
                text_frame.paragraphs[0]
                if index == 0
                else text_frame.add_paragraph()
            )
            paragraph.text = item
            if planned_slide.body_font_size is not None:
                paragraph.font.size = Pt(planned_slide.body_font_size)

    def _apply_font_size(self, text_frame, font_size: int) -> None:
        """Apply one font size across all paragraphs in a text frame."""

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)

