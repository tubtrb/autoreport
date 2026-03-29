"""Branch-local validation helper for the PptxGenJS template-authoring spike."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import platform
import shutil
import subprocess
from tempfile import TemporaryDirectory
from typing import Any

from pptx import Presentation

from autoreport.engine.generator import generate_report_from_mapping
from autoreport.template_flow import inspect_template_contract
from autoreport.templates.weekly_report import profile_template

from experiments.v04.prototypes.pptxgenjs_template_spike.demo_payloads import (
    build_demo_payload_document,
)


REPO_ROOT = Path(__file__).resolve().parents[4]
PROTOTYPE_ROOT = Path(__file__).resolve().parent
GENERATED_DIR = PROTOTYPE_ROOT / "generated"
VALIDATION_DIR = REPO_ROOT / "experiments" / "v04" / "validation"
VALIDATION_JSON_PATH = VALIDATION_DIR / "pptxgenjs-template-authoring-evidence.json"
VALIDATION_REPORT_PATH = VALIDATION_DIR / "pptxgenjs-template-authoring-report.md"
TEMPLATE_FILES = (
    "v04-minimal-text-template.pptx",
    "v04-stacked-text-template.pptx",
    "v04-text-image-template.pptx",
)


def _repo_relative_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(REPO_ROOT.resolve()))
    except ValueError:
        return str(path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Validate the PptxGenJS-authored template fixtures."
    )
    parser.add_argument(
        "--pretty",
        action="store_true",
        help="Pretty-print the JSON output.",
    )
    parser.add_argument(
        "--write-validation",
        action="store_true",
        help="Write the JSON and markdown artifacts into `experiments/v04/validation/`.",
    )
    return parser.parse_args()


def _command_output(args: list[str]) -> str:
    executable = shutil.which(args[0])
    if executable is None and platform.system() == "Windows":
        fallback_name = f"{args[0]}.cmd"
        executable = shutil.which(fallback_name)
    if executable is None:
        return ""

    completed = subprocess.run(
        [executable, *args[1:]],
        cwd=REPO_ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )
    if completed.returncode != 0:
        return ""
    return completed.stdout.strip()


def _placeholder_layout_dump(template_path: Path) -> list[dict[str, Any]]:
    presentation = Presentation(str(template_path))
    layouts: list[dict[str, Any]] = []
    for index, layout in enumerate(presentation.slide_layouts):
        placeholders: list[dict[str, Any]] = []
        for placeholder in layout.placeholders:
            placeholder_format = placeholder.placeholder_format
            placeholders.append(
                {
                    "idx": placeholder_format.idx,
                    "type": str(placeholder_format.type).split()[0].lower(),
                    "name": placeholder.name,
                    "left": placeholder.left,
                    "top": placeholder.top,
                    "width": placeholder.width,
                    "height": placeholder.height,
                }
            )
        layouts.append(
            {
                "index": index,
                "name": layout.name,
                "placeholders": placeholders,
            }
        )
    return layouts


def _generated_shape_dump(presentation: Presentation) -> list[list[dict[str, Any]]]:
    slides: list[list[dict[str, Any]]] = []
    for slide in presentation.slides:
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            try:
                placeholder_format = shape.placeholder_format
            except Exception:
                placeholder_format = None
            if placeholder_format is None:
                placeholder_type = "none"
                placeholder_index: int | str = "none"
            else:
                placeholder_type = str(placeholder_format.type).split()[0].lower()
                placeholder_index = placeholder_format.idx
            if not text:
                continue
            shapes.append(
                {
                    "placeholder_index": placeholder_index,
                    "placeholder_type": placeholder_type,
                    "text": text,
                }
            )
        slides.append(shapes)
    return slides


def _pattern_summary(pattern: dict[str, Any]) -> dict[str, Any]:
    text_slots = [
        slot for slot in pattern["slots"]
        if slot["slot_type"] in {"title", "text", "caption"}
    ]
    image_slots = [
        slot for slot in pattern["slots"]
        if slot["slot_type"] == "image"
    ]
    return {
        "kind": pattern.get("kind", "system"),
        "layout_name": pattern["layout_name"],
        "text_slot_count": len(text_slots),
        "text_orientations": [
            slot["orientation"] for slot in text_slots
            if slot.get("orientation") is not None
        ],
        "image_slot_count": len(image_slots),
        "image_orientations": [
            slot["orientation"] for slot in image_slots
            if slot.get("orientation") is not None
        ],
    }


def _contract_summary(contract_document: dict[str, Any]) -> dict[str, Any]:
    template_contract = contract_document["template_contract"]
    summary_slides = [
        {"section": "title_slide", **_pattern_summary(template_contract["title_slide"])},
        {"section": "contents_slide", **_pattern_summary(template_contract["contents_slide"])},
    ]
    for pattern in template_contract["slide_patterns"]:
        summary_slides.append({"section": pattern["kind"], **_pattern_summary(pattern)})
    return {
        "template_id": template_contract["template_id"],
        "template_source": template_contract["template_source"],
        "slides": summary_slides,
    }


def _profile_image_placeholder_indices(template_path: Path) -> set[int]:
    profile = profile_template(
        Presentation(str(template_path)),
        template_path=template_path,
    )
    indices: set[int] = set()
    for pattern in profile.slide_patterns:
        if pattern.kind != "text_image":
            continue
        for slot in pattern.slots:
            if slot.slot_type == "image" and slot.placeholder_index is not None:
                indices.add(slot.placeholder_index)
    return indices


def validate_template(template_path: Path) -> dict[str, Any]:
    result: dict[str, Any] = {
        "template_path": _repo_relative_path(template_path),
        "template_filename": template_path.name,
        "inspection_success": False,
        "contract_success": False,
        "generation_success": False,
        "contract_summary": None,
        "placeholder_dump": _placeholder_layout_dump(template_path),
    }
    try:
        contract = inspect_template_contract(template_path=template_path)
        result["inspection_success"] = True
        result["inspection_template_id"] = contract.template_id
    except Exception as exc:  # pragma: no cover - evidence capture path
        result["inspection_error"] = f"{type(exc).__name__}: {exc}"
        return result

    try:
        contract = inspect_template_contract(template_path=template_path)
        contract_document = contract.to_dict()
        result["contract_success"] = True
        result["contract_template_id"] = contract.template_id
        result["contract_summary"] = _contract_summary(contract_document)
        result["image_placeholder_indices"] = sorted(
            _profile_image_placeholder_indices(template_path)
        )
    except Exception as exc:  # pragma: no cover - evidence capture path
        result["contract_error"] = f"{type(exc).__name__}: {exc}"
        return result

    try:
        with TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / f"{template_path.stem}-generated.pptx"
            payload_document, image_refs = build_demo_payload_document(contract)
            generated_path = generate_report_from_mapping(
                payload_document,
                output_path=output_path,
                template_path=template_path,
                image_refs=image_refs,
            )
            generated_presentation = Presentation(str(generated_path))
            generated_shape_dump = _generated_shape_dump(generated_presentation)
            result["generation_success"] = True
            result["generated_output_path"] = f"tmp/{generated_path.name}"
            result["generated_slide_titles"] = [
                (
                    slide.shapes.title.text
                    if slide.shapes.title is not None
                    else next(
                        (
                            shape["text"].splitlines()[0]
                            for shape in generated_shape_dump[index]
                            if shape["text"]
                        ),
                        "(untitled)",
                    )
                )
                for index, slide in enumerate(generated_presentation.slides)
            ]
            result["generated_slide_count"] = len(generated_presentation.slides)
            result["generated_shape_dump"] = generated_shape_dump
            image_placeholder_indices = set(result["image_placeholder_indices"])
            result["text_written_to_image_placeholder"] = any(
                shape["placeholder_index"] in image_placeholder_indices
                for slide in result["generated_shape_dump"]
                for shape in slide
            )
    except Exception as exc:  # pragma: no cover - evidence capture path
        result["generation_error"] = f"{type(exc).__name__}: {exc}"

    return result


def _resolve_verdict(results: list[dict[str, Any]]) -> str:
    by_name = {result["template_filename"]: result for result in results}
    minimal = by_name["v04-minimal-text-template.pptx"]
    stacked = by_name["v04-stacked-text-template.pptx"]
    text_image = by_name["v04-text-image-template.pptx"]

    text_only_success = all(
        result["inspection_success"]
        and result["contract_success"]
        and result["generation_success"]
        for result in (minimal, stacked)
    )
    image_success = (
        text_image["inspection_success"]
        and text_image["contract_success"]
        and text_image["generation_success"]
    )
    image_pattern_available = any(
        slide["section"] == "text_image"
        for slide in text_image.get("contract_summary", {}).get("slides", [])
    )
    image_collision = bool(text_image.get("text_written_to_image_placeholder"))

    if text_only_success and image_success and image_pattern_available and not image_collision:
        return "viable"
    if text_only_success and image_success and (image_collision or not image_pattern_available):
        return "viable_text_only"
    return "not_viable"


def render_markdown_report(payload: dict[str, Any]) -> str:
    lines = [
        "# PptxGenJS Template-Authoring Spike Report",
        "",
        "## Scope",
        "",
        "- Spike type: template authoring only",
        "- Runtime replacement: no",
        "- Python runtime changes: none",
        "- Working root: `codex/v0.4-incubator2`",
        "- Evidence file: `experiments/v04/validation/pptxgenjs-template-authoring-evidence.json`",
        "",
        "## Environment",
        "",
        f"- Node: `{payload['versions']['node']}`",
        f"- npm: `{payload['versions']['npm']}`",
        f"- Python: `{payload['versions']['python']}`",
        "",
        "## Template Outputs",
    ]
    lines.extend(
        f"- `experiments/v04/prototypes/pptxgenjs_template_spike/generated/{result['template_filename']}`"
        for result in payload["templates"]
    )
    lines.extend(["", "## Results", ""])
    for result in payload["templates"]:
        lines.append(f"### `{result['template_filename']}`")
        lines.append("")
        lines.append(
            f"- Inspection: {'success' if result['inspection_success'] else 'failure'}"
        )
        lines.append(
            f"- Contract extraction: {'success' if result['contract_success'] else 'failure'}"
        )
        lines.append(
            f"- Generation smoke: {'success' if result['generation_success'] else 'failure'}"
        )
        if result["contract_summary"] is not None:
            lines.append(
                f"- `template_id`: `{result['contract_summary']['template_id']}`"
            )
            lines.append("- Contract summary:")
            for slide in result["contract_summary"]["slides"]:
                lines.append(
                    "  - "
                    f"{slide['section']}: text={slide['text_slot_count']} "
                    f"image={slide['image_slot_count']}"
                )
        if result.get("generated_slide_titles"):
            lines.append("- Generated deck reopened with visible titles:")
            for title in result["generated_slide_titles"]:
                lines.append(f"  - `{title}`")
        if result.get("text_written_to_image_placeholder"):
            lines.append(
                "- Observation: an image-capable placeholder still received text."
            )
        lines.append("")
    lines.extend(
        [
            "## Final Verdict",
            "",
            f"`{payload['verdict']}`",
            "",
            "## Why",
            "",
        ]
    )
    if payload["verdict"] == "viable":
        lines.append("- All three fixtures inspect and generate cleanly, including the mixed text-image template.")
    elif payload["verdict"] == "viable_text_only":
        lines.append("- The text-only fixtures are stable, but the mixed template still allows text into the image-capable placeholder.")
    else:
        lines.append("- The current authoring matrix is not stable enough to treat as active evidence.")
    return "\n".join(lines) + "\n"


def write_validation_artifacts(payload: dict[str, Any]) -> None:
    VALIDATION_DIR.mkdir(parents=True, exist_ok=True)
    VALIDATION_JSON_PATH.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    VALIDATION_REPORT_PATH.write_text(
        render_markdown_report(payload),
        encoding="utf-8",
    )


def main() -> int:
    args = parse_args()
    payload = {
        "versions": {
            "node": _command_output(["node", "--version"]),
            "npm": _command_output(["npm", "--version"]),
            "python": platform.python_version(),
        },
        "templates": [
            validate_template(GENERATED_DIR / filename)
            for filename in TEMPLATE_FILES
        ],
    }
    payload["verdict"] = _resolve_verdict(payload["templates"])
    if args.write_validation:
        write_validation_artifacts(payload)
    print(json.dumps(payload, indent=2 if args.pretty else None, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
