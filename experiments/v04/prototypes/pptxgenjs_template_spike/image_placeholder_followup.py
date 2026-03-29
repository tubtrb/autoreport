"""Follow-up diagnostics for the v0.4 PptxGenJS image placeholder spike."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import platform
import shutil
import subprocess
from tempfile import TemporaryDirectory
from typing import Any
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation

from autoreport.engine.generator import generate_report_from_mapping
from autoreport.template_flow import inspect_template_contract
from autoreport.templates.weekly_report import profile_template

from experiments.v04.prototypes.pptxgenjs_template_spike.demo_payloads import (
    build_demo_payload_document,
)


REPO_ROOT = Path(__file__).resolve().parents[4]
PROTOTYPE_ROOT = Path(__file__).resolve().parent
GENERATED_DIR = PROTOTYPE_ROOT / "generated"
VALIDATION_DIR = REPO_ROOT / "experiments" / "v04" / "validation"
LANDSCAPE_IMAGE_PATH = PROTOTYPE_ROOT / "assets" / "landscape.png"
CONTROL_TEMPLATE_PATH = GENERATED_DIR / "v04-control-text-picture-template.pptx"
VALIDATION_JSON_PATH = VALIDATION_DIR / "pptxgenjs-image-placeholder-followup-evidence.json"
VALIDATION_REPORT_PATH = VALIDATION_DIR / "pptxgenjs-image-placeholder-followup-report.md"
FOLLOWUP_TEMPLATE_PATHS = {
    "baseline": GENERATED_DIR / "v04-text-image-template.pptx",
    "pic_token": GENERATED_DIR / "v04-text-image-pic-token-template.pptx",
    "compact_image": GENERATED_DIR / "v04-text-image-compact-image-template.pptx",
    "control": CONTROL_TEMPLATE_PATH,
}
FOLLOWUP_TEMPLATE_CATEGORIES = {
    "baseline": "baseline",
    "pic_token": "variant",
    "compact_image": "variant",
    "control": "control",
}
FOLLOWUP_TEMPLATE_INTENTS = {
    "baseline": "PptxGenJS baseline image placeholder recipe",
    "pic_token": "PptxGenJS raw `pic` placeholder token attempt",
    "compact_image": "PptxGenJS geometry workaround with lower image prominence",
    "control": "python-pptx + OOXML control fixture with a real `pic` placeholder",
}
IMAGE_PLACEHOLDER_TYPES = {
    "object",
    "picture",
    "bitmap",
    "vertical_object",
}
PPTX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _repo_relative_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(REPO_ROOT.resolve()))
    except ValueError:
        return str(path)


def _scrub_embedded_image_paths(template_path: Path) -> None:
    with zipfile.ZipFile(template_path, "r") as source_zip:
        files = {
            info.filename: source_zip.read(info.filename)
            for info in source_zip.infolist()
        }

    replacements = {
        str(LANDSCAPE_IMAGE_PATH.resolve()): "assets/landscape.png",
    }
    changed = False
    for filename, data in list(files.items()):
        if not filename.endswith(".xml"):
            continue
        text = data.decode("utf-8", errors="ignore")
        updated = text
        for source, target in replacements.items():
            updated = updated.replace(source, target)
        if updated == text:
            continue
        files[filename] = updated.encode("utf-8")
        changed = True

    if not changed:
        return

    with zipfile.ZipFile(
        template_path,
        "w",
        compression=zipfile.ZIP_DEFLATED,
    ) as output_zip:
        for filename, data in files.items():
            output_zip.writestr(filename, data)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Diagnose the v0.4 mixed image/text placeholder behavior."
    )
    parser.add_argument(
        "--pretty",
        action="store_true",
        help="Pretty-print the JSON payload.",
    )
    parser.add_argument(
        "--write-validation",
        action="store_true",
        help="Write the machine-readable evidence and markdown report.",
    )
    return parser.parse_args()


def _command_output(args: list[str]) -> str:
    executable = shutil.which(args[0])
    if executable is None and platform.system() == "Windows":
        executable = shutil.which(f"{args[0]}.cmd")
    if executable is None:
        return ""

    completed = subprocess.run(
        [executable, *args[1:]],
        cwd=REPO_ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )
    if completed.returncode != 0:
        return ""
    return completed.stdout.strip()


def _emu(inches: float) -> int:
    return round(inches * 914400)


def _disable_non_target_layout(layout_root: ET.Element) -> None:
    for placeholder in layout_root.findall(
        "p:cSld/p:spTree/p:sp/p:nvSpPr/p:nvPr/p:ph",
        PPTX_NS,
    ):
        idx = placeholder.get("idx")
        if idx not in {None, "10", "11", "12"}:
            placeholder.set("type", "dt")


def _set_geometry(shape: ET.Element, *, x: int, y: int, w: int, h: int) -> None:
    xfrm = shape.find("p:spPr/a:xfrm", PPTX_NS)
    if xfrm is None:
        return
    off = xfrm.find("a:off", PPTX_NS)
    ext = xfrm.find("a:ext", PPTX_NS)
    if off is None or ext is None:
        return
    off.set("x", str(x))
    off.set("y", str(y))
    ext.set("cx", str(w))
    ext.set("cy", str(h))


def _rewrite_control_layouts(template_path: Path) -> None:
    with zipfile.ZipFile(template_path, "r") as source_zip:
        files = {
            info.filename: source_zip.read(info.filename)
            for info in source_zip.infolist()
        }

    for part_index in range(2, 12):
        layout_part = f"ppt/slideLayouts/slideLayout{part_index}.xml"
        if layout_part not in files:
            continue
        if part_index == 5:
            continue
        layout_root = ET.fromstring(files[layout_part])
        _disable_non_target_layout(layout_root)
        files[layout_part] = ET.tostring(
            layout_root,
            encoding="utf-8",
            xml_declaration=True,
        )

    body_layout_part = "ppt/slideLayouts/slideLayout5.xml"
    body_layout_root = ET.fromstring(files[body_layout_part])
    c_sld = body_layout_root.find("p:cSld", PPTX_NS)
    if c_sld is not None:
        c_sld.set("name", "V04 Control Mixed Body")

    for shape in body_layout_root.findall("p:cSld/p:spTree/p:sp", PPTX_NS):
        placeholder = shape.find("p:nvSpPr/p:nvPr/p:ph", PPTX_NS)
        if placeholder is None:
            continue

        idx = placeholder.get("idx")
        if idx == "1":
            placeholder.set("type", "body")
            _set_geometry(
                shape,
                x=_emu(0.8),
                y=_emu(1.55),
                w=_emu(5.9),
                h=_emu(4.7),
            )
            continue
        if idx == "2":
            placeholder.set("type", "dt")
            continue
        if idx == "3":
            placeholder.set("type", "body")
            _set_geometry(
                shape,
                x=_emu(7.15),
                y=_emu(5.3),
                w=_emu(5.0),
                h=_emu(0.7),
            )
            continue
        if idx == "4":
            placeholder.set("type", "pic")
            _set_geometry(
                shape,
                x=_emu(7.15),
                y=_emu(1.55),
                w=_emu(5.0),
                h=_emu(3.55),
            )

    files[body_layout_part] = ET.tostring(
        body_layout_root,
        encoding="utf-8",
        xml_declaration=True,
    )

    with zipfile.ZipFile(
        template_path,
        "w",
        compression=zipfile.ZIP_DEFLATED,
    ) as output_zip:
        for filename, data in files.items():
            output_zip.writestr(filename, data)


def _slide_placeholders_by_idx(slide) -> dict[int, Any]:
    placeholders: dict[int, Any] = {}
    for placeholder in slide.placeholders:
        placeholders[placeholder.placeholder_format.idx] = placeholder
    return placeholders


def generate_control_fixture() -> Path:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.save(str(CONTROL_TEMPLATE_PATH))
    _rewrite_control_layouts(CONTROL_TEMPLATE_PATH)

    presentation = Presentation(str(CONTROL_TEMPLATE_PATH))
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_placeholders = _slide_placeholders_by_idx(title_slide)
    title_slide.shapes.title.text = "Preview Title"
    if 1 in title_placeholders:
        title_placeholders[1].text = "Preview Subtitle"

    body_slide = presentation.slides.add_slide(presentation.slide_layouts[4])
    body_placeholders = _slide_placeholders_by_idx(body_slide)
    body_placeholders[0].text = "Product Highlights"
    body_placeholders[1].text = "Primary text placeholder preview"
    body_placeholders[3].text = "Caption placeholder preview"
    body_placeholders[4].insert_picture(str(LANDSCAPE_IMAGE_PATH))
    presentation.save(str(CONTROL_TEMPLATE_PATH))
    _scrub_embedded_image_paths(CONTROL_TEMPLATE_PATH)
    return CONTROL_TEMPLATE_PATH


def _normalize_placeholder_type(placeholder_type: Any) -> str:
    return str(placeholder_type).split()[0].lower()


def _shape_geometry(shape: ET.Element) -> dict[str, int | None]:
    xfrm = shape.find("p:spPr/a:xfrm", PPTX_NS)
    if xfrm is None:
        xfrm = shape.find("p:xfrm", PPTX_NS)
    if xfrm is None:
        return {
            "x": None,
            "y": None,
            "w": None,
            "h": None,
        }
    off = xfrm.find("a:off", PPTX_NS)
    ext = xfrm.find("a:ext", PPTX_NS)
    return {
        "x": int(off.get("x")) if off is not None else None,
        "y": int(off.get("y")) if off is not None else None,
        "w": int(ext.get("cx")) if ext is not None else None,
        "h": int(ext.get("cy")) if ext is not None else None,
    }


def _extract_part_placeholder_dump(part_name: str, xml_bytes: bytes) -> dict[str, Any]:
    root = ET.fromstring(xml_bytes)
    c_sld = root.find("p:cSld", PPTX_NS)
    placeholders: list[dict[str, Any]] = []
    for shape in root.findall("p:cSld/p:spTree/*", PPTX_NS):
        placeholder = shape.find("p:nvSpPr/p:nvPr/p:ph", PPTX_NS)
        if placeholder is None:
            placeholder = shape.find("p:nvPicPr/p:nvPr/p:ph", PPTX_NS)
        if placeholder is None:
            continue

        c_nv_pr = shape.find("p:nvSpPr/p:cNvPr", PPTX_NS)
        if c_nv_pr is None:
            c_nv_pr = shape.find("p:nvPicPr/p:cNvPr", PPTX_NS)
        geometry = _shape_geometry(shape)
        placeholders.append(
            {
                "shape_tag": shape.tag.split("}")[-1],
                "shape_name": c_nv_pr.get("name") if c_nv_pr is not None else "",
                "idx": int(placeholder.get("idx")) if placeholder.get("idx") else None,
                "type": placeholder.get("type"),
                "orientation": placeholder.get("orient"),
                "size_hint": placeholder.get("sz"),
                "has_text_body": shape.find("p:txBody", PPTX_NS) is not None,
                "has_blip_fill": shape.find("p:blipFill", PPTX_NS) is not None,
                **geometry,
            }
        )

    return {
        "part_name": part_name,
        "name": c_sld.get("name") if c_sld is not None else "",
        "placeholders": placeholders,
    }


def _ooxml_placeholder_dump(template_path: Path) -> dict[str, list[dict[str, Any]]]:
    with zipfile.ZipFile(template_path, "r") as archive:
        dumps = {"masters": [], "layouts": [], "slides": []}
        for filename in archive.namelist():
            if not filename.endswith(".xml"):
                continue
            if filename.startswith("ppt/slideMasters/"):
                dumps["masters"].append(
                    _extract_part_placeholder_dump(
                        filename,
                        archive.read(filename),
                    )
                )
                continue
            if filename.startswith("ppt/slideLayouts/"):
                dumps["layouts"].append(
                    _extract_part_placeholder_dump(
                        filename,
                        archive.read(filename),
                    )
                )
                continue
            if filename.startswith("ppt/slides/"):
                dumps["slides"].append(
                    _extract_part_placeholder_dump(
                        filename,
                        archive.read(filename),
                    )
                )
        for key in dumps:
            dumps[key].sort(key=lambda item: item["part_name"])
        return dumps


def _python_layout_dump(template_path: Path) -> list[dict[str, Any]]:
    presentation = Presentation(str(template_path))
    layouts: list[dict[str, Any]] = []
    for index, layout in enumerate(presentation.slide_layouts):
        placeholders: list[dict[str, Any]] = []
        for placeholder in layout.placeholders:
            placeholder_format = placeholder.placeholder_format
            placeholders.append(
                {
                    "idx": placeholder_format.idx,
                    "type": _normalize_placeholder_type(placeholder_format.type),
                    "name": placeholder.name,
                    "has_text_frame": hasattr(placeholder, "text_frame"),
                    "left": placeholder.left,
                    "top": placeholder.top,
                    "width": placeholder.width,
                    "height": placeholder.height,
                }
            )
        layouts.append(
            {
                "index": index,
                "name": layout.name,
                "placeholders": placeholders,
            }
        )
    return layouts


def _generated_shape_dump(presentation: Presentation) -> list[list[dict[str, Any]]]:
    slides: list[list[dict[str, Any]]] = []
    for slide in presentation.slides:
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if not text:
                continue
            placeholder_format = getattr(shape, "placeholder_format", None)
            if placeholder_format is None:
                placeholder_type = "none"
                placeholder_index: int | str = "none"
            else:
                placeholder_type = _normalize_placeholder_type(placeholder_format.type)
                placeholder_index = placeholder_format.idx
            shapes.append(
                {
                    "placeholder_index": placeholder_index,
                    "placeholder_type": placeholder_type,
                    "text": text,
                }
            )
        slides.append(shapes)
    return slides


def _profile_slot_summary(slot: Any) -> dict[str, Any]:
    return {
        "slot_name": slot.slot_name,
        "placeholder_index": slot.placeholder_index,
        "placeholder_type": _normalize_placeholder_type(slot.placeholder_type),
        "orientation": slot.orientation,
        "order": slot.order,
        "required": slot.required,
    }


def _visible_slide_titles(presentation: Presentation) -> list[str]:
    titles: list[str] = []
    for slide in presentation.slides:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            titles.append(slide.shapes.title.text)
            continue
        fallback = "(untitled)"
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                fallback = text.splitlines()[0]
                break
        titles.append(fallback)
    return titles


def _requested_image_type(deck_id: str) -> str:
    if deck_id == "pic_token":
        return "pic"
    if deck_id == "control":
        return "pic"
    return "image"


def _resolve_verdict(decks: list[dict[str, Any]]) -> tuple[str, list[str], list[str]]:
    by_id = {deck["deck_id"]: deck for deck in decks}
    baseline = by_id["baseline"]
    pic_token = by_id["pic_token"]
    compact = by_id["compact_image"]
    control = by_id["control"]

    authored_decks = (baseline, pic_token, compact)
    authored_problem = any(
        (
            not deck.get("contract_success", False)
            or not deck.get("contract_summary", {}).get("text_image_pattern_available", False)
            or not deck.get("generation_success", False)
            or deck.get("generation_summary", {}).get("text_written_to_image_placeholder", False)
        )
        for deck in authored_decks
    )
    control_problem = (
        (
            not control.get("contract_success", False)
            or not control.get("contract_summary", {}).get("text_image_pattern_available", False)
            or not control.get("generation_success", False)
            or control.get("generation_summary", {}).get("text_written_to_image_placeholder", False)
        )
    )

    if not authored_problem and not control_problem:
        verdict = "runtime_hardened"
    elif control_problem and authored_problem:
        verdict = "mixed"
    elif control_problem:
        verdict = "profiler_issue"
    elif authored_problem:
        verdict = "authoring_issue"
    else:
        verdict = "inconclusive"

    reasons: list[str] = []
    next_steps = []
    if verdict == "runtime_hardened":
        reasons.extend(
            [
                (
                    "all four comparison decks now inspect into a reusable `text_image` pattern "
                    "under the rebased runtime"
                ),
                (
                    "actual profiled image placeholder indices remain separate from body and caption "
                    "indices, and the generated decks no longer write text into those image placeholders"
                ),
                (
                    "the PptxGenJS fixtures still emit `OBJECT` placeholders instead of raw OOXML `pic`, "
                    "but the hardened runtime now treats those placeholders as image-capable slots"
                ),
            ]
        )
        next_steps.append(
            "Use the mixed text-image fixtures as workflow automation alpha inputs and keep runtime changes scoped to regression-backed profiling fixes."
        )
    elif verdict == "mixed":
        reasons.extend(
            [
                "both the authoring recipe and the runtime boundary still contribute to mixed-layout instability.",
                "at least one authored fixture and the control fixture still diverge on inspection or generation behavior.",
            ]
        )
        next_steps.extend(
            [
                (
                    "Document a short-term authoring recipe: keep the primary body placeholder "
                    "materially larger than the image placeholder, and do not rely on PptxGenJS "
                    "placeholder type tokens to emit a true `pic` placeholder."
                ),
                (
                    "Plan a minimal runtime patch in the main repo to stop `_list_text_placeholders` "
                    "from admitting image-capable placeholders as body text candidates."
                ),
            ]
        )
    elif verdict == "authoring_issue":
        reasons.extend(
            [
                "the control fixture is stable, but at least one PptxGenJS-authored fixture still fails to export or render a safe mixed contract.",
                "the remaining mixed-layout gap is specific to how the template was authored, not to the current runtime boundary.",
            ]
        )
        next_steps.append(
            "Publish a safe PptxGenJS authoring recipe and keep the runtime untouched."
        )
    elif verdict == "profiler_issue":
        reasons.extend(
            [
                "the control fixture still fails or collides while the authored variants remain stable enough to inspect or render.",
                "the remaining blocker sits in the runtime profiling boundary rather than in the template authoring recipe.",
            ]
        )
        next_steps.append(
            "Move to a runtime-only follow-up that repairs mixed-layout profiling before rerunning the authoring comparison."
        )
    else:
        reasons.append(
            "the refreshed post-patch evidence does not cleanly separate authoring behavior from runtime behavior."
        )
        next_steps.append(
            "Hold the status at `viable_text_only` until another control axis can separate authoring from profiler behavior."
        )

    return verdict, reasons, next_steps


def _deck_report_summary(deck: dict[str, Any]) -> str:
    if not deck.get("contract_success"):
        return (
            f"- `{deck['deck_id']}`\n"
            f"  path: `{deck['template_path']}`\n"
            f"  contract: `failure`\n"
            f"  error: `{deck.get('contract_error', 'unknown')}`"
        )

    body_text_slots = ", ".join(
        f"{slot['placeholder_index']}:{slot['placeholder_type']}:{slot['orientation']}"
        for slot in deck["contract_summary"]["body_text_slots"]
    ) or "(none)"
    body_image_slots = ", ".join(
        f"{slot['placeholder_index']}:{slot['placeholder_type']}:{slot['orientation']}"
        for slot in deck["contract_summary"]["body_image_slots"]
    ) or "(none)"
    collision = (
        "yes"
        if deck["generation_summary"]["text_written_to_image_placeholder"]
        else "no"
    )
    return (
        f"- `{deck['deck_id']}`\n"
        f"  path: `{deck['template_path']}`\n"
        f"  body layout: `{deck['contract_summary']['body_layout_name']}`\n"
        f"  body text slots: `{body_text_slots}`\n"
        f"  body image slots: `{body_image_slots}`\n"
        f"  text written into image placeholder: `{collision}`"
    )


def render_markdown_report(evidence: dict[str, Any]) -> str:
    commands = [
        "npm ci",
        "node generate_templates.mjs",
        ".\\.venv\\Scripts\\python.exe experiments/v04/prototypes/pptxgenjs_template_spike/image_placeholder_followup.py --write-validation --pretty",
        ".\\.venv\\Scripts\\python.exe -m unittest tests.test_pptxgenjs_template_spike tests.test_pptxgenjs_image_placeholder_followup",
        ".\\.venv\\Scripts\\python.exe -m unittest tests.test_workflow_automation_spike tests.test_v04_prototype_scaffolds tests.test_workflow_automation_sandbox",
        ".\\.venv\\Scripts\\python.exe -m unittest tests.test_generator",
    ]
    deck_lines = "\n".join(_deck_report_summary(deck) for deck in evidence["decks"])
    versions = evidence["versions"]
    verdict = evidence["verdict"]
    reasons = "\n".join(f"- {reason}" for reason in evidence["verdict_reasons"])
    next_steps = "\n".join(
        f"- {step}"
        for step in evidence["recommended_next_steps"]
    )
    return f"""# v0.4 PptxGenJS Image Placeholder Follow-up Report

## Summary

- Goal: reassess the mixed image/text placeholder comparison after rebasing onto the latest `v0.3-master` runtime and applying the incubator hardening patches.
- Final verdict: `{verdict}`
- Scope: branch-local diagnostics only under `experiments/v04/`

## Environment

- Node: `{versions['node']}`
- npm: `{versions['npm']}`
- Python: `{versions['python']}`

## Commands

{chr(10).join(f"- `{command}`" for command in commands)}

## Deck Matrix

{deck_lines}

## Verdict Reasons

{reasons}

## Next Steps

{next_steps}
"""


def diagnose_template(deck_id: str, template_path: Path) -> dict[str, Any]:
    ooxml_dump = _ooxml_placeholder_dump(template_path)
    python_layout_dump = _python_layout_dump(template_path)

    result: dict[str, Any] = {
        "deck_id": deck_id,
        "category": FOLLOWUP_TEMPLATE_CATEGORIES[deck_id],
        "intent": FOLLOWUP_TEMPLATE_INTENTS[deck_id],
        "requested_image_placeholder_type": _requested_image_type(deck_id),
        "template_path": _repo_relative_path(template_path),
        "inspection_success": False,
        "contract_success": False,
        "generation_success": False,
        "ooxml_dump": ooxml_dump,
        "python_layout_dump": python_layout_dump,
    }

    try:
        inspection = inspect_template_contract(template_path=template_path)
        result["inspection_success"] = True
        result["inspection_template_id"] = inspection.template_id
    except Exception as exc:  # pragma: no cover - evidence capture path
        result["inspection_error"] = f"{type(exc).__name__}: {exc}"
        return result

    try:
        contract = inspect_template_contract(template_path=template_path)
        profile = profile_template(
            Presentation(str(template_path)),
            template_path=template_path,
        )
        contract_document = contract.to_dict()["template_contract"]
        result["contract_success"] = True
        result["contract_template_id"] = contract.template_id
        body_pattern = next(
            (
                pattern
                for pattern in contract_document["slide_patterns"]
                if pattern["kind"] == "text_image"
            ),
            None,
        )
        profiled_body_pattern = next(
            (
                pattern
                for pattern in profile.slide_patterns
                if pattern.kind == "text_image"
            ),
            None,
        )
        contents_slide = contract_document["contents_slide"]
        result["contract_summary"] = {
            "template_id": contract_document["template_id"],
            "text_image_pattern_available": body_pattern is not None,
            "body_layout_name": (
                body_pattern["layout_name"]
                if body_pattern is not None
                else contents_slide["layout_name"]
            ),
            "body_text_slots": [
                _profile_slot_summary(slot)
                for slot in (
                    profiled_body_pattern.slots
                    if profiled_body_pattern is not None
                    else []
                )
                if slot.slot_type == "text"
            ],
            "body_image_slots": [
                _profile_slot_summary(slot)
                for slot in (
                    profiled_body_pattern.slots
                    if profiled_body_pattern is not None
                    else []
                )
                if slot.slot_type == "image"
            ],
            "contents_text_slots": [
                {
                    "slot_id": slot["slot_id"],
                    "orientation": slot["orientation"],
                    "order": slot["order"],
                    "required": slot["required"],
                }
                for slot in contents_slide["slots"]
                if slot["slot_type"] == "text"
            ],
        }
    except Exception as exc:  # pragma: no cover - evidence capture path
        result["contract_error"] = f"{type(exc).__name__}: {exc}"
        return result

    body_layout_name = result["contract_summary"]["body_layout_name"]
    body_layout_ooxml = next(
        (
            layout
            for layout in ooxml_dump["layouts"]
            if layout["name"] == body_layout_name
        ),
        None,
    )
    body_layout_python = next(
        (
            layout
            for layout in python_layout_dump
            if layout["name"] == body_layout_name
        ),
        None,
    )
    result["body_layout_ooxml"] = body_layout_ooxml
    result["body_layout_python"] = body_layout_python

    try:
        with TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / f"{template_path.stem}-generated.pptx"
            payload_document, image_refs = build_demo_payload_document(contract)
            generated_path = generate_report_from_mapping(
                payload_document,
                output_path=output_path,
                template_path=template_path,
                image_refs=image_refs,
            )
            generated_presentation = Presentation(str(generated_path))
            generated_shape_dump = _generated_shape_dump(generated_presentation)
            image_placeholder_indices = {
                slot["placeholder_index"]
                for slot in result["contract_summary"]["body_image_slots"]
                if slot["placeholder_index"] is not None
            }
            image_placeholder_text_hits = [
                {
                    "slide_index": slide_index,
                    "slide_title": _visible_slide_titles(generated_presentation)[slide_index],
                    **shape,
                }
                for slide_index, slide_shapes in enumerate(generated_shape_dump)
                for shape in slide_shapes
                if shape["placeholder_index"] in image_placeholder_indices
            ]
            result["generation_success"] = True
            result["generation_summary"] = {
                "generated_slide_count": len(generated_presentation.slides),
                "generated_slide_titles": _visible_slide_titles(generated_presentation),
                "generated_shape_dump": generated_shape_dump,
                "image_placeholder_indices": sorted(image_placeholder_indices),
                "image_placeholder_text_hits": image_placeholder_text_hits,
                "text_written_to_image_placeholder": bool(image_placeholder_text_hits),
            }
    except Exception as exc:  # pragma: no cover - evidence capture path
        result["generation_error"] = f"{type(exc).__name__}: {exc}"

    return result


def collect_followup_evidence() -> dict[str, Any]:
    generate_control_fixture()
    decks = [
        diagnose_template(deck_id, template_path)
        for deck_id, template_path in FOLLOWUP_TEMPLATE_PATHS.items()
    ]
    verdict, reasons, next_steps = _resolve_verdict(decks)
    return {
        "versions": {
            "node": _command_output(["node", "--version"]),
            "npm": _command_output(["npm", "--version"]),
            "python": platform.python_version(),
        },
        "verdict": verdict,
        "verdict_reasons": reasons,
        "recommended_next_steps": next_steps,
        "decks": decks,
    }


def write_validation_artifacts(evidence: dict[str, Any]) -> None:
    VALIDATION_DIR.mkdir(parents=True, exist_ok=True)
    VALIDATION_JSON_PATH.write_text(
        json.dumps(evidence, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    VALIDATION_REPORT_PATH.write_text(
        render_markdown_report(evidence),
        encoding="utf-8",
    )


def main() -> int:
    args = parse_args()
    evidence = collect_followup_evidence()
    if args.write_validation:
        write_validation_artifacts(evidence)
    print(json.dumps(evidence, indent=2 if args.pretty else None, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
