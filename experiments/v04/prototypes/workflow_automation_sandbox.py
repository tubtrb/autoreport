"""Sandbox runner for v0.4 workflow automation rehearsals.

This module stays branch-local and reuses the current contract-first runtime
without wiring the incubator flow into the public CLI or web entrypoints.
"""

from __future__ import annotations

import argparse
from copy import deepcopy
from dataclasses import dataclass
from datetime import datetime
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
import yaml

from autoreport.engine.generator import generate_report_from_mapping
from autoreport.loader import load_yaml
from autoreport.template_flow import (
    PUBLIC_BUILT_IN_TEMPLATE_NAME,
    inspect_template_contract,
    scaffold_report_payload,
    serialize_document,
)
from autoreport.validator import validate_payload

from .pptxgenjs_template_spike.demo_payloads import DEFAULT_IMAGE_REFS
from .workflow_automation_reporting import (
    sandbox_command,
    sandbox_input_summary,
    write_worker_status,
)
from .workflow_automation_spike import (
    AutomationTrigger,
    build_template_report_automation_plan,
    plan_to_markdown,
    summarize_manual_gates,
)

DEFAULT_SANDBOX_CONFIG = "run-config.yaml"
DEFAULT_REPORT_BRIEF = "report-brief.yaml"
DEFAULT_OUTPUT_FILENAME = "autoreport_workflow_preview.pptx"


@dataclass(frozen=True)
class SandboxRunOptions:
    """Configuration values that shape one sandbox automation rehearsal."""

    template_name: str
    template_path: Path | None
    include_text_shaping: bool
    include_human_review: bool
    include_publish_handoff: bool


@dataclass(frozen=True)
class SandboxLayout:
    """Resolved directories and files inside one sandbox root."""

    root: Path
    inputs_dir: Path
    plans_dir: Path
    contracts_dir: Path
    drafts_dir: Path
    artifacts_dir: Path
    reviews_dir: Path
    logs_dir: Path

    @classmethod
    def from_root(cls, root: Path) -> "SandboxLayout":
        return cls(
            root=root,
            inputs_dir=root / "inputs",
            plans_dir=root / "plans",
            contracts_dir=root / "contracts",
            drafts_dir=root / "drafts",
            artifacts_dir=root / "artifacts",
            reviews_dir=root / "reviews",
            logs_dir=root / "logs",
        )

    def ensure_dirs(self) -> None:
        for path in (
            self.inputs_dir,
            self.plans_dir,
            self.contracts_dir,
            self.drafts_dir,
            self.artifacts_dir,
            self.reviews_dir,
            self.logs_dir,
        ):
            path.mkdir(parents=True, exist_ok=True)


@dataclass(frozen=True)
class SandboxRunResult:
    """Small result payload returned after a sandbox rehearsal completes."""

    sandbox_root: Path
    template_name: str
    template_reference: str
    template_id: str
    artifact_path: Path
    review_path: Path
    slide_titles: list[str]
    manual_gates: list[str]
    completed_at: str

    def to_mapping(self) -> dict[str, Any]:
        return {
            "sandbox_root": str(self.sandbox_root.resolve()),
            "template_name": self.template_name,
            "template_reference": self.template_reference,
            "template_id": self.template_id,
            "artifact_path": str(self.artifact_path.resolve()),
            "review_path": str(self.review_path.resolve()),
            "slide_titles": list(self.slide_titles),
            "manual_gates": list(self.manual_gates),
            "completed_at": self.completed_at,
        }


def load_sandbox_configuration(sandbox_root: Path) -> tuple[AutomationTrigger, SandboxRunOptions]:
    """Load the sandbox trigger and plan options from the input config file."""

    config_path = sandbox_root / "inputs" / DEFAULT_SANDBOX_CONFIG
    raw = load_yaml(config_path)

    trigger_raw = raw.get("trigger", {})
    if not isinstance(trigger_raw, dict):
        raise ValueError("Field 'trigger' must be an object.")

    trigger = AutomationTrigger(
        name=str(trigger_raw.get("name", sandbox_root.name)).strip() or sandbox_root.name,
        source=str(trigger_raw.get("source", "sandbox")).strip() or "sandbox",
        cadence=str(trigger_raw.get("cadence", "on_demand")).strip() or "on_demand",
    )

    options_raw = raw.get("plan_options", {})
    if not isinstance(options_raw, dict):
        raise ValueError("Field 'plan_options' must be an object.")

    built_in = str(
        raw.get("built_in", PUBLIC_BUILT_IN_TEMPLATE_NAME)
    ).strip() or PUBLIC_BUILT_IN_TEMPLATE_NAME
    template_path_raw = raw.get("template_path")
    template_path: Path | None = None
    if template_path_raw:
        template_path = (sandbox_root / str(template_path_raw)).resolve()
    options = SandboxRunOptions(
        template_name=built_in,
        template_path=template_path,
        include_text_shaping=bool(options_raw.get("include_text_shaping", False)),
        include_human_review=bool(options_raw.get("include_human_review", True)),
        include_publish_handoff=bool(options_raw.get("include_publish_handoff", False)),
    )
    return trigger, options


def run_workflow_automation_sandbox(
    sandbox_root: Path,
    *,
    repo_root: Path | None = None,
) -> SandboxRunResult:
    """Run one branch-local workflow automation rehearsal inside a sandbox."""

    layout = SandboxLayout.from_root(sandbox_root)
    layout.ensure_dirs()
    brief_path = layout.inputs_dir / DEFAULT_REPORT_BRIEF
    config_path = layout.inputs_dir / DEFAULT_SANDBOX_CONFIG
    command = sandbox_command(sandbox_root, repo_root=repo_root)
    options: SandboxRunOptions | None = None
    input_summary = (
        f"sandbox={sandbox_root.resolve()}; "
        f"config={config_path.resolve()}; "
        f"brief={brief_path.resolve()}"
    )

    try:
        trigger, options = load_sandbox_configuration(sandbox_root)
        input_summary = sandbox_input_summary(
            sandbox_root=sandbox_root,
            config_path=config_path,
            brief_path=brief_path,
            template_name=options.template_name,
            template_path=options.template_path,
        )
        plan = build_template_report_automation_plan(
            trigger,
            include_text_shaping=options.include_text_shaping,
            include_human_review=options.include_human_review,
            include_publish_handoff=options.include_publish_handoff,
        )

        brief_document = load_yaml(brief_path)
        contract = inspect_template_contract(
            template_path=options.template_path,
            built_in=(
                None
                if options.template_path is not None
                else options.template_name
            ),
        )
        contract_document = contract.to_dict()
        template_id = contract.template_id

        payload_skeleton = scaffold_report_payload(contract)
        payload_skeleton_document = payload_skeleton.to_dict()
        draft_payload_document = _build_payload_document(
            brief_document=brief_document,
            payload_skeleton_document=payload_skeleton_document,
            template_id=template_id,
        )
        image_refs = _collect_image_refs(draft_payload_document)
        validated_payload = validate_payload(
            draft_payload_document,
            contract,
            available_image_refs=image_refs.keys(),
        )
        runtime_template_name = _runtime_template_name(options)

        artifact_path = layout.artifacts_dir / DEFAULT_OUTPUT_FILENAME
        generate_report_from_mapping(
            validated_payload.to_dict(),
            output_path=artifact_path,
            template_name=runtime_template_name,
            template_path=options.template_path,
            image_refs=image_refs,
        )

        completed_at = datetime.now().astimezone().isoformat(timespec="seconds")
        slide_titles = _extract_slide_titles(artifact_path)
        review_path = layout.reviews_dir / "review-notes.md"

        _write_text(
            layout.plans_dir / "automation-plan.md",
            plan_to_markdown(plan) + "\n",
        )
        _write_json(
            layout.contracts_dir / "template-contract.json",
            contract_document,
        )
        _write_yaml(
            layout.contracts_dir / "template-contract.yaml",
            contract_document,
        )
        _write_text(
            layout.drafts_dir / "payload-skeleton.yaml",
            serialize_document(payload_skeleton_document, fmt="yaml"),
        )
        _write_json(
            layout.drafts_dir / "payload-skeleton.json",
            payload_skeleton_document,
        )
        _write_yaml(
            layout.drafts_dir / "draft-report-payload.yaml",
            draft_payload_document,
        )
        _write_json(
            layout.drafts_dir / "draft-report-payload.json",
            draft_payload_document,
        )
        _write_text(
            review_path,
            _build_review_notes(
                artifact_path=artifact_path,
                template_name=options.template_name,
                template_id=template_id,
                slide_titles=slide_titles,
                manual_gates=summarize_manual_gates(plan),
                completed_at=completed_at,
            ),
        )

        result = SandboxRunResult(
            sandbox_root=sandbox_root,
            template_name=_template_display_name(options),
            template_reference=_template_reference(options),
            template_id=template_id,
            artifact_path=artifact_path,
            review_path=review_path,
            slide_titles=slide_titles,
            manual_gates=summarize_manual_gates(plan),
            completed_at=completed_at,
        )
        summary_path = layout.logs_dir / "run-summary.json"
        _write_json(
            summary_path,
            result.to_mapping(),
        )
    except Exception as exc:
        write_worker_status(
            repo_root=repo_root,
            status="blocked",
            task_summary=(
                "Workflow automation sandbox failed before producing a reviewable artifact."
            ),
            evidence_input=input_summary,
            evidence_command=command,
            artifact_paths=[
                config_path,
                brief_path,
                *_existing_output_paths(layout),
            ],
            visible_result=f"Sandbox run failed: {type(exc).__name__}: {exc}",
            remaining_gap=(
                "Fix the failing template inspection, payload drafting, validation, "
                f"or generation step and rerun the sandbox: {type(exc).__name__}: {exc}"
            ),
        )
        raise

    write_worker_status(
        repo_root=repo_root,
        status=(
            "ready_for_review"
            if options is not None and options.include_human_review
            else "in_progress"
        ),
        task_summary=(
            "Workflow automation sandbox generated a reviewable artifact."
            if options is not None and options.include_human_review
            else "Workflow automation sandbox generated an artifact without a manual review gate."
        ),
        evidence_input=input_summary,
        evidence_command=command,
        artifact_paths=[
            result.artifact_path,
            result.review_path,
            layout.logs_dir / "run-summary.json",
            layout.contracts_dir / "template-contract.json",
            layout.drafts_dir / "draft-report-payload.yaml",
        ],
        visible_result=(
            f"Generated {len(result.slide_titles)} slides with primary artifact at "
            f"{result.artifact_path.resolve()}."
        ),
        remaining_gap=(
            "Human review and explicit finalize step are still required before master handoff."
            if options is not None and options.include_human_review
            else "Explicit finalize step is still required if this run should become a master-review handoff."
        ),
    )
    return result


def _extract_slide_titles(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    titles: list[str] = []
    for slide in presentation.slides:
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.text.strip():
            titles.append(title_shape.text)
            continue

        fallback_title = "(untitled)"
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                fallback_title = text.splitlines()[0]
                break
        titles.append(fallback_title)
    return titles


def _build_payload_document(
    *,
    brief_document: dict[str, Any],
    payload_skeleton_document: dict[str, Any],
    template_id: str,
) -> dict[str, Any]:
    payload_root = deepcopy(payload_skeleton_document["report_payload"])
    brief_root = brief_document.get("report_payload", brief_document)

    if not isinstance(brief_root, dict):
        raise ValueError("Sandbox report brief must be a mapping.")

    title_slide = brief_root.get("title_slide")
    if title_slide is not None:
        payload_root["title_slide"] = title_slide

    contents = brief_root.get("contents")
    if contents is not None:
        payload_root["contents"] = contents

    slides = brief_root.get("slides")
    if slides is not None:
        payload_root["slides"] = slides

    payload_root["template_id"] = template_id
    return {"report_payload": payload_root}


def _collect_image_refs(payload_document: dict[str, Any]) -> dict[str, Path]:
    payload_root = payload_document["report_payload"]
    refs: dict[str, Path] = {}
    for slide in payload_root.get("slides", []):
        if not isinstance(slide, dict):
            continue
        image = slide.get("image")
        if not isinstance(image, dict):
            continue
        ref = image.get("ref")
        if isinstance(ref, str) and ref in DEFAULT_IMAGE_REFS:
            refs[ref] = DEFAULT_IMAGE_REFS[ref]
    return refs


def _build_review_notes(
    *,
    artifact_path: Path,
    template_name: str,
    template_id: str,
    slide_titles: list[str],
    manual_gates: list[str],
    completed_at: str,
) -> str:
    lines = [
        "# Review Notes",
        "",
        f"- Completed at: {completed_at}",
        f"- Template name: {template_name}",
        f"- Template id: {template_id}",
        f"- Artifact path: {artifact_path.resolve()}",
        f"- Slide count: {len(slide_titles)}",
        "",
        "## Visible checks",
    ]
    lines.extend(f"- {title}" for title in slide_titles)
    lines.append("")
    lines.append("## Manual gates")
    if manual_gates:
        lines.extend(f"- {gate}" for gate in manual_gates)
    else:
        lines.append("- No explicit manual gates configured for this sandbox run.")
    lines.append("")
    lines.append("## Outcome")
    if manual_gates:
        lines.append(
            "- The artifact was generated successfully and is ready for human review inside the sandbox."
        )
    else:
        lines.append(
            "- The artifact was generated successfully. No manual review gate was configured for this sandbox run."
        )
    return "\n".join(lines) + "\n"


def _template_display_name(options: SandboxRunOptions) -> str:
    if options.template_path is not None:
        return options.template_path.stem
    return options.template_name


def _runtime_template_name(options: SandboxRunOptions) -> str:
    if options.template_path is not None:
        return PUBLIC_BUILT_IN_TEMPLATE_NAME
    return options.template_name


def _template_reference(options: SandboxRunOptions) -> str:
    if options.template_path is not None:
        return str(options.template_path.resolve())
    return options.template_name


def _existing_output_paths(layout: SandboxLayout) -> list[Path]:
    return [
        layout.plans_dir / "automation-plan.md",
        layout.contracts_dir / "template-contract.json",
        layout.contracts_dir / "template-contract.yaml",
        layout.drafts_dir / "payload-skeleton.yaml",
        layout.drafts_dir / "payload-skeleton.json",
        layout.drafts_dir / "draft-report-payload.yaml",
        layout.drafts_dir / "draft-report-payload.json",
        layout.artifacts_dir / DEFAULT_OUTPUT_FILENAME,
        layout.reviews_dir / "review-notes.md",
        layout.logs_dir / "run-summary.json",
    ]


def _write_text(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def _write_json(path: Path, payload: dict[str, Any]) -> None:
    _write_text(
        path,
        json.dumps(payload, indent=2, ensure_ascii=False) + "\n",
    )


def _write_yaml(path: Path, payload: dict[str, Any]) -> None:
    _write_text(
        path,
        yaml.safe_dump(payload, sort_keys=False, allow_unicode=True),
    )


def build_parser() -> argparse.ArgumentParser:
    """Create the sandbox runner parser."""

    parser = argparse.ArgumentParser(
        description="Run a branch-local v0.4 workflow automation sandbox."
    )
    parser.add_argument(
        "--sandbox",
        required=True,
        help="Path to the sandbox root directory.",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    """Run one sandbox rehearsal and print a small JSON summary."""

    parser = build_parser()
    args = parser.parse_args(argv)
    result = run_workflow_automation_sandbox(Path(args.sandbox))
    print(json.dumps(result.to_mapping(), indent=2, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
