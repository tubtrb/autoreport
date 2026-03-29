"""Branch-local tests for the v0.4 image placeholder follow-up spike."""

from __future__ import annotations

from functools import lru_cache
import shutil
import tempfile
import unittest
from pathlib import Path
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation

from autoreport.engine.generator import generate_report_from_mapping
from autoreport.template_flow import inspect_template_contract
from experiments.v04.prototypes.pptxgenjs_template_spike.image_placeholder_followup import (
    CONTROL_TEMPLATE_PATH,
    FOLLOWUP_TEMPLATE_PATHS,
    PPTX_NS,
    collect_followup_evidence,
)
from experiments.v04.prototypes.pptxgenjs_template_spike.demo_payloads import (
    build_demo_payload_document,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
PIC_TOKEN_TEMPLATE_PATH = FOLLOWUP_TEMPLATE_PATHS["pic_token"]
COMPACT_TEMPLATE_PATH = FOLLOWUP_TEMPLATE_PATHS["compact_image"]


def _shape_text_hits(template_path: Path, *, image_placeholder_indices: set[int]) -> list[dict[str, str | int]]:
    contract = inspect_template_contract(template_path=template_path)
    payload_document, image_refs = build_demo_payload_document(contract)
    temp_root = Path(tempfile.mkdtemp())
    try:
        output_path = temp_root / f"{template_path.stem}-output.pptx"
        generated_path = generate_report_from_mapping(
            payload_document,
            output_path=output_path,
            template_path=template_path,
            image_refs=image_refs,
        )
        presentation = Presentation(str(generated_path))
        hits: list[dict[str, str | int]] = []
        for slide_index, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if not text:
                    continue
                placeholder_format = getattr(shape, "placeholder_format", None)
                if placeholder_format is None:
                    continue
                if placeholder_format.idx not in image_placeholder_indices:
                    continue
                hits.append(
                    {
                        "slide_index": slide_index,
                        "placeholder_index": placeholder_format.idx,
                        "placeholder_type": str(placeholder_format.type).split()[0].lower(),
                        "text": text,
                    }
                )
        return hits
    finally:
        shutil.rmtree(temp_root, ignore_errors=True)


def _ooxml_layout_placeholders(template_path: Path, *, layout_part: str) -> list[dict[str, str | int | None]]:
    with zipfile.ZipFile(template_path, "r") as archive:
        root = ET.fromstring(archive.read(layout_part))
    placeholders: list[dict[str, str | int | None]] = []
    for shape in root.findall("p:cSld/p:spTree/*", PPTX_NS):
        placeholder = shape.find("p:nvSpPr/p:nvPr/p:ph", PPTX_NS)
        if placeholder is None:
            placeholder = shape.find("p:nvPicPr/p:nvPr/p:ph", PPTX_NS)
        if placeholder is None:
            continue
        placeholders.append(
            {
                "idx": int(placeholder.get("idx")) if placeholder.get("idx") else None,
                "type": placeholder.get("type"),
            }
        )
    return placeholders


@lru_cache(maxsize=1)
def _followup_evidence() -> dict:
    return collect_followup_evidence()


class PptxGenJSImagePlaceholderFollowupTestCase(unittest.TestCase):
    """Verify the branch-local follow-up evidence for mixed image/text placeholders."""

    def test_followup_fixture_files_exist(self) -> None:
        for template_path in FOLLOWUP_TEMPLATE_PATHS.values():
            self.assertTrue(template_path.exists(), msg=str(template_path))

    def test_pic_token_variant_does_not_emit_a_raw_pic_placeholder(self) -> None:
        placeholders = _ooxml_layout_placeholders(
            PIC_TOKEN_TEMPLATE_PATH,
            layout_part="ppt/slideLayouts/slideLayout3.xml",
        )
        image_placeholder = next(
            placeholder
            for placeholder in placeholders
            if placeholder["idx"] == 102
        )

        self.assertIsNone(image_placeholder["type"])

    def test_compact_variant_avoids_text_being_written_into_the_image_placeholder(self) -> None:
        hits = _shape_text_hits(
            COMPACT_TEMPLATE_PATH,
            image_placeholder_indices={102},
        )

        self.assertEqual(hits, [])

    def test_control_fixture_preserves_pic_placeholder_and_followup_captures_no_collision(self) -> None:
        placeholders = _ooxml_layout_placeholders(
            CONTROL_TEMPLATE_PATH,
            layout_part="ppt/slideLayouts/slideLayout5.xml",
        )
        image_placeholder = next(
            placeholder
            for placeholder in placeholders
            if placeholder["idx"] == 4
        )
        evidence = _followup_evidence()
        control_deck = next(
            deck
            for deck in evidence["decks"]
            if deck["deck_id"] == "control"
        )

        self.assertEqual(image_placeholder["type"], "pic")
        self.assertTrue(control_deck["inspection_success"])
        self.assertTrue(control_deck["contract_success"])
        self.assertTrue(control_deck["generation_success"])
        self.assertFalse(
            control_deck["generation_summary"]["text_written_to_image_placeholder"]
        )

    def test_followup_verdict_is_runtime_hardened(self) -> None:
        evidence = _followup_evidence()

        self.assertEqual(evidence["verdict"], "runtime_hardened")


if __name__ == "__main__":
    unittest.main()
