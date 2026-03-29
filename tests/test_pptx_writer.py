"""Tests for PowerPoint writer output and template loading behavior."""

from __future__ import annotations

import base64
import shutil
import unittest
from pathlib import Path
from unittest.mock import patch
from uuid import uuid4

from pptx import Presentation

from autoreport.outputs.pptx_writer import PowerPointWriter, TemplateLoadError, TemplateReadError
from autoreport.templates.autofill import (
    FillPlan,
    PlannedImageFill,
    PlannedSlide,
    PlannedTextFill,
    SlotDescriptor,
)


TEST_TEMP_ROOT = Path("tests") / "_tmp"
PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jXioAAAAASUVORK5CYII="
)


def make_test_dir() -> Path:
    TEST_TEMP_ROOT.mkdir(exist_ok=True)
    test_dir = TEST_TEMP_ROOT / uuid4().hex
    test_dir.mkdir()
    return test_dir


class PowerPointWriterTestCase(unittest.TestCase):
    """Verify writer-level output and file loading behavior."""

    def test_load_presentation_raises_template_load_error_for_invalid_file(self) -> None:
        writer = PowerPointWriter()
        test_dir = make_test_dir()
        try:
            template_path = test_dir / "invalid-template.pptx"
            template_path.write_text("not a pptx", encoding="utf-8")

            with self.assertRaises(TemplateLoadError):
                writer._load_presentation(template_path)
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

    def test_load_presentation_raises_template_read_error_for_os_errors(self) -> None:
        writer = PowerPointWriter()
        test_dir = make_test_dir()
        try:
            template_path = test_dir / "template.pptx"
            Presentation().save(str(template_path))

            with patch(
                "autoreport.outputs.pptx_writer.Presentation",
                side_effect=PermissionError("locked"),
            ):
                with self.assertRaises(TemplateReadError):
                    writer._load_presentation(template_path)
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

    def test_write_fill_plan_supports_text_box_and_image_slots(self) -> None:
        writer = PowerPointWriter()
        test_dir = make_test_dir()
        try:
            image_path = test_dir / "workflow.png"
            image_path.write_bytes(PNG_BYTES)
            output_path = test_dir / "writer-output.pptx"
            presentation = Presentation()
            fill_plan = FillPlan(
                slides=[
                    PlannedSlide(
                        pattern_id="text_image.editorial",
                        kind="text_image",
                        layout_name="Blank",
                        layout_index=6,
                        slide_title="Autoreport",
                        text_fills=[
                            PlannedTextFill(
                                slot=SlotDescriptor(
                                    slot_name="title",
                                    alias="title",
                                    slot_type="title",
                                    layout_index=6,
                                    placeholder_index=None,
                                    x=600000,
                                    y=600000,
                                    width=5000000,
                                    height=700000,
                                ),
                                text="Autoreport",
                                font_size=28,
                            ),
                            PlannedTextFill(
                                slot=SlotDescriptor(
                                    slot_name="body_1",
                                    alias="body_1",
                                    slot_type="text",
                                    layout_index=6,
                                    placeholder_index=None,
                                    x=600000,
                                    y=1600000,
                                    width=5000000,
                                    height=1800000,
                                ),
                                items=[
                                    "Generate editable PowerPoint decks.",
                                    "Fill template slots deterministically.",
                                ],
                                font_size=18,
                            ),
                        ],
                        image_fills=[
                            PlannedImageFill(
                                slot=SlotDescriptor(
                                    slot_name="image_1",
                                    alias="image_1",
                                    slot_type="image",
                                    layout_index=6,
                                    placeholder_index=None,
                                    x=6500000,
                                    y=1600000,
                                    width=2200000,
                                    height=1800000,
                                ),
                                image_path=image_path,
                                fit="contain",
                            )
                        ],
                    )
                ]
            )

            writer.write_fill_plan(
                presentation=presentation,
                output_path=output_path,
                fill_plan=fill_plan,
            )

            written = Presentation(str(output_path))
            texts = [
                shape.text
                for shape in written.slides[0].shapes
                if hasattr(shape, "text") and shape.text
            ]
            image_shapes = [
                shape
                for shape in written.slides[0].shapes
                if hasattr(shape, "image")
            ]
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertIn("Autoreport", texts)
        self.assertIn("Generate editable PowerPoint decks.", "\n".join(texts))
        self.assertGreaterEqual(len(image_shapes), 1)


if __name__ == "__main__":
    unittest.main()
