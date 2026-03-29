"""Tests for contract-first generation and template inspection."""

from __future__ import annotations

import base64
import shutil
import unittest
from pathlib import Path
from uuid import uuid4
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation

from autoreport.engine.generator import (
    generate_report,
    generate_report_from_mapping,
    prepare_generation_artifacts_from_mapping,
)
from autoreport.models import ReportRequest
from autoreport.template_flow import inspect_template_contract, scaffold_payload
from autoreport.templates.weekly_report import profile_template


TEST_TEMP_ROOT = Path("tests") / "_tmp"
PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jXioAAAAASUVORK5CYII="
)
PPTX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def make_test_dir() -> Path:
    TEST_TEMP_ROOT.mkdir(exist_ok=True)
    test_dir = TEST_TEMP_ROOT / uuid4().hex
    test_dir.mkdir()
    return test_dir


def build_authoring_payload(
    *,
    image_paths: list[Path] | None = None,
    image_orientation: str = "auto",
    include_contents: bool = True,
) -> dict[str, object]:
    payload = scaffold_payload(
        inspect_template_contract(built_in="autoreport_editorial"),
        include_text_image=(image_paths is not None),
    ).to_dict()
    if image_paths is None:
        payload["authoring_payload"]["contents"]["enabled"] = include_contents
        return payload

    payload["authoring_payload"]["contents"]["enabled"] = include_contents
    text_image_slide = payload["authoring_payload"]["slides"][2]
    text_image_slide["context"]["summary"] = "Compare visuals against the authored narrative."
    text_image_slide["assets"]["images"] = [
        {"path": str(path), "fit": "contain"} for path in image_paths
    ]
    text_image_slide["layout_request"]["image_count"] = len(image_paths)
    text_image_slide["layout_request"]["image_orientation"] = image_orientation
    return payload


def build_runtime_text_image_payload(
    *,
    image_path: Path,
    text_image_body: list[str],
    template_id: str = "autoreport-editorial-v1",
    include_contents: bool = True,
) -> dict[str, object]:
    return {
        "report_payload": {
            "payload_version": "autoreport.payload.v1",
            "template_id": template_id,
            "title_slide": {
                "title": "Autoreport",
                "subtitle": ["Template-aware PPTX autofill engine"],
            },
            "contents": {"enabled": include_contents},
            "slides": [
                {
                    "kind": "text_image",
                    "title": "Why It Matters",
                    "include_in_contents": True,
                    "body": text_image_body,
                    "image": {"path": str(image_path), "fit": "contain"},
                    "caption": "Workflow preview",
                    "slot_overrides": {},
                }
            ],
        }
    }


def build_text_image_spill_body() -> list[str]:
    return [
        (
            f"Bullet {index}: template-aware autofill keeps editorial structure "
            "stable while long narrative content spills onto a continuation slide "
            "for release inspection."
        )
        for index in range(1, 11)
    ]


def _emu(inches: float) -> int:
    return round(inches * 914400)


def _set_geometry(shape: ET.Element, *, x: int, y: int, w: int, h: int) -> None:
    xfrm = shape.find("p:spPr/a:xfrm", PPTX_NS)
    if xfrm is None:
        return
    off = xfrm.find("a:off", PPTX_NS)
    ext = xfrm.find("a:ext", PPTX_NS)
    if off is None or ext is None:
        return
    off.set("x", str(x))
    off.set("y", str(y))
    ext.set("cx", str(w))
    ext.set("cy", str(h))


def _disable_non_target_layout(layout_root: ET.Element) -> None:
    for placeholder in layout_root.findall(
        "p:cSld/p:spTree/p:sp/p:nvSpPr/p:nvPr/p:ph",
        PPTX_NS,
    ):
        idx = placeholder.get("idx")
        if idx not in {None, "10", "11", "12"}:
            placeholder.set("type", "dt")


def _rewrite_mixed_layout(template_path: Path, *, image_placeholder_type: str) -> None:
    with zipfile.ZipFile(template_path, "r") as source_zip:
        files = {
            info.filename: source_zip.read(info.filename)
            for info in source_zip.infolist()
        }

    for part_index in range(2, 12):
        layout_part = f"ppt/slideLayouts/slideLayout{part_index}.xml"
        if layout_part not in files or part_index == 5:
            continue
        layout_root = ET.fromstring(files[layout_part])
        _disable_non_target_layout(layout_root)
        files[layout_part] = ET.tostring(
            layout_root,
            encoding="utf-8",
            xml_declaration=True,
        )

    body_layout_part = "ppt/slideLayouts/slideLayout5.xml"
    body_layout_root = ET.fromstring(files[body_layout_part])
    c_sld = body_layout_root.find("p:cSld", PPTX_NS)
    if c_sld is not None:
        c_sld.set(
            "name",
            f"Autoreport Mixed {image_placeholder_type.upper()} Body",
        )

    for shape in body_layout_root.findall("p:cSld/p:spTree/p:sp", PPTX_NS):
        placeholder = shape.find("p:nvSpPr/p:nvPr/p:ph", PPTX_NS)
        if placeholder is None:
            continue

        idx = placeholder.get("idx")
        if idx == "1":
            placeholder.set("type", "body")
            _set_geometry(
                shape,
                x=_emu(0.8),
                y=_emu(1.55),
                w=_emu(5.9),
                h=_emu(4.7),
            )
            continue
        if idx == "2":
            placeholder.set("type", "dt")
            continue
        if idx == "3":
            placeholder.set("type", "body")
            _set_geometry(
                shape,
                x=_emu(7.15),
                y=_emu(5.3),
                w=_emu(5.0),
                h=_emu(0.7),
            )
            continue
        if idx == "4":
            placeholder.set("type", image_placeholder_type)
            _set_geometry(
                shape,
                x=_emu(7.15),
                y=_emu(1.55),
                w=_emu(5.0),
                h=_emu(3.55),
            )

    files[body_layout_part] = ET.tostring(
        body_layout_root,
        encoding="utf-8",
        xml_declaration=True,
    )

    with zipfile.ZipFile(
        template_path,
        "w",
        compression=zipfile.ZIP_DEFLATED,
    ) as output_zip:
        for filename, data in files.items():
            output_zip.writestr(filename, data)


def build_mixed_user_template(
    test_dir: Path,
    *,
    image_placeholder_type: str,
) -> Path:
    template_path = test_dir / f"mixed-{image_placeholder_type}.pptx"
    Presentation().save(str(template_path))
    _rewrite_mixed_layout(
        template_path,
        image_placeholder_type=image_placeholder_type,
    )
    return template_path


def _placeholders_by_idx(slide):
    return {
        placeholder.placeholder_format.idx: placeholder
        for placeholder in slide.placeholders
    }


class GeneratorTestCase(unittest.TestCase):
    """Verify contract-first generation creates the expected PowerPoint structure."""

    def test_prepare_generation_artifacts_exposes_editorial_contract(self) -> None:
        artifacts = prepare_generation_artifacts_from_mapping(
            build_authoring_payload(),
            presentation=Presentation(),
        )

        self.assertEqual(artifacts.template_contract.template_id, "autoreport-editorial-v1")
        self.assertEqual(
            [pattern.pattern_id for pattern in artifacts.template_contract.slide_patterns],
            [
                "text.editorial",
                "metrics.editorial",
                "text_image.editorial",
                "text_image.editorial.two_horizontal",
                "text_image.editorial.two_vertical",
                "text_image.editorial.three_horizontal",
                "text_image.editorial.three_vertical",
            ],
        )
        self.assertEqual(
            [slide.slide_title for slide in artifacts.fill_plan.slides],
            ["Autoreport", "Contents", "What It Does", "Adoption Snapshot"],
        )
        self.assertEqual(len(artifacts.diagnostic_report.errors), 0)

    def test_generate_report_from_mapping_creates_editorial_presentation_from_authoring_payload(self) -> None:
        test_dir = make_test_dir()
        try:
            output_path = test_dir / "output" / "autoreport_demo.pptx"
            generated_path = generate_report_from_mapping(
                build_authoring_payload(),
                output_path=output_path,
            )
            presentation = Presentation(str(output_path))
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(generated_path, output_path)
        self.assertEqual(len(presentation.slides), 4)

    def test_generate_report_from_mapping_supports_single_image_authoring(self) -> None:
        test_dir = make_test_dir()
        try:
            image_path = test_dir / "workflow.png"
            image_path.write_bytes(PNG_BYTES)
            output_path = test_dir / "output" / "autoreport_text_image.pptx"

            generate_report_from_mapping(
                build_authoring_payload(image_paths=[image_path]),
                output_path=output_path,
            )
            presentation = Presentation(str(output_path))
            image_shapes = [
                shape
                for shape in presentation.slides[-1].shapes
                if hasattr(shape, "image")
            ]
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(len(presentation.slides), 5)
        self.assertGreaterEqual(len(image_shapes), 1)

    def test_generate_report_from_mapping_supports_two_horizontal_images(self) -> None:
        test_dir = make_test_dir()
        try:
            image_paths = []
            for name in ("one.png", "two.png"):
                path = test_dir / name
                path.write_bytes(PNG_BYTES)
                image_paths.append(path)
            output_path = test_dir / "output" / "autoreport_two_horizontal.pptx"

            generate_report_from_mapping(
                build_authoring_payload(
                    image_paths=image_paths,
                    image_orientation="horizontal",
                ),
                output_path=output_path,
            )
            presentation = Presentation(str(output_path))
            image_shapes = [
                shape
                for shape in presentation.slides[-1].shapes
                if hasattr(shape, "image")
            ]
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(len(presentation.slides), 5)
        self.assertEqual(len(image_shapes), 2)

    def test_generate_report_from_mapping_supports_three_vertical_images(self) -> None:
        test_dir = make_test_dir()
        try:
            image_paths = []
            for name in ("one.png", "two.png", "three.png"):
                path = test_dir / name
                path.write_bytes(PNG_BYTES)
                image_paths.append(path)
            output_path = test_dir / "output" / "autoreport_three_vertical.pptx"

            generate_report_from_mapping(
                build_authoring_payload(
                    image_paths=image_paths,
                    image_orientation="vertical",
                    include_contents=False,
                ),
                output_path=output_path,
            )
            presentation = Presentation(str(output_path))
            image_shapes = [
                shape
                for shape in presentation.slides[-1].shapes
                if hasattr(shape, "image")
            ]
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(len(presentation.slides), 4)
        self.assertEqual(len(image_shapes), 3)

    def test_runtime_text_image_continuation_does_not_repeat_media(self) -> None:
        test_dir = make_test_dir()
        try:
            image_path = test_dir / "workflow.png"
            image_path.write_bytes(PNG_BYTES)
            artifacts = prepare_generation_artifacts_from_mapping(
                build_runtime_text_image_payload(
                    image_path=image_path,
                    text_image_body=build_text_image_spill_body(),
                ),
                presentation=Presentation(),
            )
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        text_image_slides = [
            slide
            for slide in artifacts.generation_summary.slides
            if slide.kind == "text_image"
        ]

        self.assertEqual(
            [slide.slide_title for slide in text_image_slides],
            ["Why It Matters", "Why It Matters (cont.)"],
        )
        self.assertEqual(text_image_slides[0].image_slot_names, ("text_image.image_1",))
        self.assertEqual(text_image_slides[1].image_slot_names, ())

    def test_generate_report_uses_payload_file_and_default_output_directory(self) -> None:
        test_dir = make_test_dir()
        previous_cwd = Path.cwd()
        try:
            payload_path = (test_dir / "autoreport_payload.yaml").resolve()
            payload_path.write_text(
                __import__("yaml").safe_dump(
                    build_authoring_payload(),
                    sort_keys=False,
                    allow_unicode=True,
                ),
                encoding="utf-8",
            )
            expected_output_path = (test_dir / "output" / "autoreport_payload.pptx").resolve()

            __import__("os").chdir(test_dir)
            generated_path = generate_report(ReportRequest(source_path=payload_path))
            generated_exists = expected_output_path.exists()
        finally:
            __import__("os").chdir(previous_cwd)
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(generated_path, Path("output") / "autoreport_payload.pptx")
        self.assertTrue(generated_exists)

    def test_user_template_is_profiled_into_the_same_contract_family(self) -> None:
        test_dir = make_test_dir()
        try:
            template_path = test_dir / "sample-template.pptx"
            Presentation().save(str(template_path))
            contract = inspect_template_contract(template_path=template_path)
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(contract.template_source, "user_template")
        self.assertIn("text", [pattern.kind for pattern in contract.slide_patterns])
        self.assertIn("metrics", [pattern.kind for pattern in contract.slide_patterns])
        self.assertIn("text_image", [pattern.kind for pattern in contract.slide_patterns])

    def test_pic_mixed_user_template_exports_a_stable_text_image_contract(self) -> None:
        test_dir = make_test_dir()
        try:
            template_path = build_mixed_user_template(
                test_dir,
                image_placeholder_type="pic",
            )
            first_contract = inspect_template_contract(template_path=template_path)
            second_contract = inspect_template_contract(template_path=template_path)
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(first_contract.template_source, "user_template")
        self.assertEqual(first_contract.template_id, second_contract.template_id)
        mixed_pattern = next(
            pattern
            for pattern in first_contract.slide_patterns
            if pattern.kind == "text_image"
        )
        slot_types = [slot.slot_type for slot in mixed_pattern.slots]
        self.assertIn("image", slot_types)
        self.assertIn("caption", slot_types)
        self.assertIn("text", slot_types)

    def test_object_placeholder_is_classified_as_image_not_text(self) -> None:
        test_dir = make_test_dir()
        try:
            template_path = build_mixed_user_template(
                test_dir,
                image_placeholder_type="obj",
            )
            contract = inspect_template_contract(template_path=template_path)
            profile = profile_template(
                Presentation(str(template_path)),
                template_path=template_path,
            )
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertIn("text_image", [pattern.kind for pattern in contract.slide_patterns])
        mixed_pattern = next(
            pattern
            for pattern in profile.slide_patterns
            if pattern.kind == "text_image"
        )
        image_slot = mixed_pattern.slots_by_type("image")[0]
        text_placeholder_indexes = {
            slot.placeholder_index
            for slot in mixed_pattern.slots
            if slot.slot_type in {"title", "text", "caption"}
        }
        self.assertIsNotNone(image_slot.placeholder_index)
        self.assertNotIn(image_slot.placeholder_index, text_placeholder_indexes)

    def test_generate_report_from_mapping_supports_mixed_user_template(self) -> None:
        test_dir = make_test_dir()
        try:
            template_path = build_mixed_user_template(
                test_dir,
                image_placeholder_type="pic",
            )
            contract = inspect_template_contract(template_path=template_path)
            image_path = test_dir / "workflow.png"
            image_path.write_bytes(PNG_BYTES)
            output_path = test_dir / "output" / "autoreport_user_template_text_image.pptx"

            generated_path = generate_report_from_mapping(
                build_runtime_text_image_payload(
                    image_path=image_path,
                    text_image_body=["Teams keep their own brand language."],
                    template_id=contract.template_id,
                    include_contents=False,
                ),
                output_path=output_path,
                template_path=template_path,
            )
            presentation = Presentation(str(output_path))
            mixed_slide = presentation.slides[1]
            texts = [
                shape.text
                for shape in mixed_slide.shapes
                if hasattr(shape, "text") and shape.text
            ]
            image_shapes = [
                shape
                for shape in mixed_slide.shapes
                if hasattr(shape, "image")
            ]
            placeholders = _placeholders_by_idx(mixed_slide)
        finally:
            shutil.rmtree(test_dir, ignore_errors=True)

        self.assertEqual(generated_path, output_path)
        self.assertEqual(len(presentation.slides), 2)
        self.assertGreaterEqual(len(image_shapes), 1)
        self.assertIn("Why It Matters", texts)
        self.assertIn("Teams keep their own brand language.", texts)
        self.assertIn("Workflow preview", texts)
        self.assertEqual(placeholders[1].text, "Teams keep their own brand language.")
        self.assertEqual(placeholders[3].text, "Workflow preview")


if __name__ == "__main__":
    unittest.main()
