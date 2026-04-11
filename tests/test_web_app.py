"""Tests for the public FastAPI contract-first demo application."""

from __future__ import annotations

import base64
import gc
from io import BytesIO
import json
from pathlib import Path
import tempfile
import unittest
from unittest.mock import patch
import warnings

from fastapi.testclient import TestClient
from pptx import Presentation

from autoreport.loader import parse_yaml_text
from autoreport.template_flow import serialize_document
from autoreport.web.app import (
    MANUAL_PROCEDURE_EXAMPLE_YAML,
    MEDIA_TYPE_PPTX,
    PROMPTED_MANUAL_PROCEDURE_EXAMPLE_YAML,
    app,
)


PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jXioAAAAASUVORK5CYII="
)

VALID_AUTHORING_PAYLOAD_YAML = """
authoring_payload:
  payload_version: autoreport.authoring.v1
  template_id: autoreport-editorial-v1
  deck_context:
    audience: executives
  title_slide:
    title: Autoreport
    subtitle:
      - Template-aware PPTX autofill engine
  contents:
    enabled: true
  slides:
    - slide_no: 1
      goal: What It Does
      include_in_contents: true
      context:
        summary: Generate editable PowerPoint decks from structured inputs.
      layout_request:
        kind: text
        image_orientation: auto
""".strip()

VALID_TEXT_IMAGE_AUTHORING_PAYLOAD_YAML = """
authoring_payload:
  payload_version: autoreport.authoring.v1
  template_id: autoreport-editorial-v1
  deck_context:
    audience: executives
  title_slide:
    title: Autoreport
    subtitle:
      - Template-aware PPTX autofill engine
  contents:
    enabled: true
  slides:
    - slide_no: 1
      goal: Visual Proof
      include_in_contents: true
      context:
        summary: Pair narrative context with an uploaded image.
        caption: Workflow preview
      assets:
        images:
          - ref: image_1
            fit: contain
      layout_request:
        kind: text_image
        image_orientation: auto
""".strip()

VALID_REPORT_CONTENT_WITH_IMAGE_YAML = """
report_content:
  title_slide:
    pattern_id: cover.editorial
    slots:
      title: Autoreport
      subtitle_1: |
        Template-aware PPTX autofill engine
  slides:
    - pattern_id: text_image.editorial
      slots:
        title: Visual Proof
        body_1: |
          Explain what the visual should prove.
        image_1: image_1
        caption_1: Example caption
""".strip()

VALID_REPORT_CONTENT_YAML = """
report_content:
  title_slide:
    pattern_id: cover.editorial
    slots:
      title: Autoreport
      subtitle_1: |
        Template-aware PPTX autofill engine
  contents_slide:
    pattern_id: contents.editorial
    slots:
      title: Contents
      body_1: |
        1. What It Does
  slides:
    - pattern_id: text.editorial
      slots:
        title: What It Does
        body_1: |
          Generate editable PowerPoint decks from structured inputs.
""".strip()

INVALID_MANUAL_PATTERN_REPORT_CONTENT_YAML = """
report_content:
  title_slide:
    pattern_id: cover.manual
    slots:
      doc_title: HBM4 Architecture & Usage Guide
      doc_subtitle: Screenshot-first guide for HBM4 validation
      doc_version: v1.0.0
      author_or_owner: Memory Architecture Team
  slides:
    - pattern_id: text.manual.section_break
      slots:
        section_no: "2."
        section_title: HBM4 Overview
        section_subtitle: Basic concept and key improvements over HBM3.
    - pattern_id: image.manual.step
      slots:
        step_no: "2-1"
        step_title: Identify HBM4 Structure
        command_or_action: Review HBM4 stack diagram
        summary: Understand TSV-based 3D stacking and wide I/O interface.
        image_1: image_1
        caption_1: HBM4 vertical stack structure
        detail_body: |
          HBM4 uses a 3D-stacked DRAM architecture with TSVs.
""".strip()

BROKEN_MANUAL_INDENTATION_REPORT_CONTENT_YAML = """
report_content:
title_slide:
pattern_id: cover.manual
slots:
doc_title: Autoreport PowerPoint User Guide
doc_subtitle: |
Screenshot-first guide for the public web manual mode
doc_version: v0.4.2
author_or_owner: Autoreport Team
contents_slide:
pattern_id: contents.manual
slots:
contents_title: Contents
contents_group_label: Procedure Overview
slides:
- pattern_id: text.manual.section_break
slots:
section_no: "1."
section_title: Review The Manual Starter
section_subtitle: Start with the built-in manual procedure starter before editing content.
- pattern_id: text_image.manual.procedure.one
slots:
step_no: "1.1"
step_title: Review The Starter Example
command_or_action: "Action: open the starter example and confirm the selected template."
summary: Use one screenshot to show the starting editor state.
detail_body: |
Review the starter YAML, note the built-in template mode, and confirm
the page is ready before moving to the next step.
image_1: image_1
caption_1: Starter example loaded in the editor
""".strip()

TRUNCATED_CLAUDE_REPORT_CONTENT_YAML = """
report_content:
  title_slide:
    pattern_id: cover.editorial
    slots:
      title: Autoreport
      subtitle_1: |
        Template-aware PPTX autofill engine
  slides:
    - pattern_id: text_image.edito
""".strip()

MIXED_CHATGPT_STYLE_REPORT_CONTENT = """
report_content:
title_slide:
pattern_id: cover.editorial
slots:
title: 沃섎㈇????? ?겸뫖猷롦?餓λ쵎猷??類ㅺ쉭

```yaml
- pattern_id: text.editorial
  kind: text
  slots:
    title: 筌ㅼ뮄???袁㏃뻣?? ???뼎 ?怨몄젎
    body_1: |
      筌ㅼ뮄???겸뫖猷?? 癰귣벏鍮?怨몄몵嚥??袁㏃뻣??랁???덈뼄.
```
""".strip()

_SHORT_MANUAL_DETAIL_BODY = """        detail_body: |
          Review the starter YAML, note the built-in template mode, and confirm
          the page is ready before moving to the next step."""


def build_long_manual_procedure_example_yaml() -> str:
    replacement = "        detail_body: |\n" + "\n".join(
        [
            "          Review the starter YAML carefully and keep documenting each operator cue before moving forward."
            for _ in range(12)
        ]
    )
    return MANUAL_PROCEDURE_EXAMPLE_YAML.replace(
        _SHORT_MANUAL_DETAIL_BODY,
        replacement,
        1,
    )


def build_single_slide_manual_procedure_example_yaml() -> str:
    payload = parse_yaml_text(MANUAL_PROCEDURE_EXAMPLE_YAML)
    payload["report_content"]["slides"] = payload["report_content"]["slides"][:1]
    return serialize_document(payload, fmt="yaml").strip()


def build_mutated_manual_report_content_yaml(mutate) -> str:
    payload = parse_yaml_text(MANUAL_PROCEDURE_EXAMPLE_YAML)
    report_content = payload["report_content"]
    mutate(report_content)
    return serialize_document(payload, fmt="yaml").strip()


class WebAppTestCase(unittest.TestCase):
    """Verify the public demo page and its public-only API contract."""

    def setUp(self) -> None:
        self.client = TestClient(app)

    def test_demo_page_renders_manual_starter_homepage(self) -> None:
        response = self.client.get("/")

        self.assertEqual(response.status_code, 200)
        self.assertIn("Edit the starter deck", response.text)
        self.assertIn("Starter Deck YAML", response.text)
        self.assertIn("Status", response.text)
        self.assertIn("Reset Starter Example", response.text)
        self.assertIn("Refresh Preview", response.text)
        self.assertIn("Generate PPTX", response.text)
        self.assertIn("Add Slide", response.text)
        self.assertIn("Check Draft", response.text)
        self.assertIn("Delete", response.text)
        self.assertIn("Slide Style Gallery", response.text)
        self.assertIn("Draft Checker", response.text)
        self.assertIn("style-family-chip", response.text)
        self.assertIn("style-preset-card", response.text)
        self.assertIn("slide-delete-button", response.text)
        self.assertIn(
            '["contents_slide", "content_slide"].includes(preview.source_kind)',
            response.text,
        )
        self.assertNotIn(
            '["title_slide", "contents_slide", "content_slide"].includes(preview.source_kind)',
            response.text,
        )
        self.assertIn("Text", response.text)
        self.assertIn("1 Image", response.text)
        self.assertIn("2 Images", response.text)
        self.assertIn("3 Images", response.text)
        self.assertIn("Section Break", response.text)
        self.assertIn("Style 1", response.text)
        self.assertIn("Style 2", response.text)
        self.assertIn("Style 3", response.text)
        self.assertIn("report_content", response.text)
        self.assertIn("Manual Procedure Starter", response.text)
        self.assertIn("screenshot-first manual flow", response.text)
        self.assertIn("Edit the YAML first, then use the style", response.text)
        self.assertIn("Never invent new pattern_id names such as image.manual.step.", response.text)
        self.assertIn("Use step numbers like 2.1, 2.2, 3.1. Do not write 2-1 or 3-1.", response.text)
        self.assertIn("PowerPoint Slide Preview", response.text)
        self.assertIn("matching upload panel on the", response.text)
        self.assertIn("list-style: none;", response.text)
        self.assertIn("font-size: 0.82rem;", response.text)
        self.assertIn("starter-pill", response.text)
        self.assertIn('data-family-filter="all"', response.text)
        self.assertIn('data-family-id="one-image"', response.text)
        self.assertNotIn("Image Order", response.text)
        self.assertNotIn("Screenshot Uploads", response.text)
        self.assertNotIn("Choose or paste screenshots for this slide", response.text)
        self.assertNotIn("Website Intro Starter", response.text)
        self.assertNotIn("Starter Mode", response.text)
        self.assertNotIn("Remove Upload", response.text)
        self.assertNotIn("Thumbnail Preview", response.text)
        self.assertNotIn("How To Use", response.text)
        self.assertNotIn("Image Uploads", response.text)
        self.assertNotIn("starter_app_workspace", response.text)
        self.assertNotIn("starter_app_uploads", response.text)
        self.assertNotIn("app-workspace.png", response.text)
        self.assertNotIn("/starter-assets/app-workspace.png", response.text)
        self.assertNotIn("Advanced Debug: Compiled Report Payload", response.text)
        self.assertNotIn("Optional: View Template Contract", response.text)
        self.assertNotIn("Normalize Draft", response.text)
        self.assertNotIn("Copy AI Package", response.text)
        self.assertNotIn("Optional: AI Prompt Package", response.text)
        self.assertNotIn("Reset To AI Draft Prompt", response.text)
        self.assertNotIn("Refresh Slide Assets", response.text)
        self.assertNotIn("Add Slide Style", response.text)
        self.assertNotIn("Manual Screenshot Workflow", response.text)
        self.assertNotIn("Manual Flow Summary", response.text)
        self.assertNotIn("Selected Slide Style", response.text)
        self.assertNotIn("Built-In Starter", response.text)

    def test_demo_page_defaults_to_prompted_manual_starter(self) -> None:
        response = self.client.get("/")

        self.assertEqual(response.status_code, 200)
        self.assertIn("Autoreport PowerPoint User Guide", response.text)
        self.assertIn("Review The Manual Starter", response.text)
        self.assertIn("Review The Starter Example", response.text)
        self.assertIn("Generate The PowerPoint", response.text)
        self.assertIn("# Paste this brief into another AI", response.text)
        self.assertIn("report_content draft below", response.text)
        self.assertIn("Goal: draft a screenshot-first procedure manual", response.text)

    def test_demo_page_includes_small_screen_layout_guards(self) -> None:
        response = self.client.get("/")

        self.assertEqual(response.status_code, 200)
        self.assertIn("@media (max-width: 640px)", response.text)
        self.assertIn("main { padding: 20px 12px 40px; }", response.text)
        self.assertIn("min-height: 44px;", response.text)
        self.assertIn("overflow-x: auto;", response.text)
        self.assertIn("scrollbar-width: thin;", response.text)
        self.assertIn("justify-content: flex-start;", response.text)

    def test_style_presets_route_returns_manual_catalog(self) -> None:
        response = self.client.get("/api/style-presets?built_in=autoreport_manual")

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["built_in"], "autoreport_manual")
        self.assertEqual(
            payload["families"],
            [
                {"family_id": "text", "label": "Text", "order": 10, "count": 1},
                {"family_id": "one-image", "label": "1 Image", "order": 20, "count": 1},
                {"family_id": "two-images", "label": "2 Images", "order": 30, "count": 1},
                {"family_id": "three-images", "label": "3 Images", "order": 40, "count": 1},
            ],
        )
        self.assertEqual(len(payload["presets"]), 4)
        self.assertEqual(payload["presets"][0]["preset_id"], "manual.section-break")
        self.assertEqual(payload["presets"][0]["pattern_id"], "text.manual.section_break")
        self.assertEqual(payload["presets"][1]["preset_id"], "manual.procedure.one")
        self.assertEqual(payload["presets"][1]["image_count"], 1)
        self.assertEqual(payload["presets"][2]["family_id"], "two-images")
        self.assertEqual(payload["presets"][3]["family_id"], "three-images")
        self.assertIn("thumbnail", payload["presets"][0])
        self.assertIn("default_slot_values", payload["presets"][0])
        self.assertIn("tags", payload["presets"][0])
        self.assertNotIn("4 Images", [family["label"] for family in payload["families"]])

    def test_healthcheck_returns_ok(self) -> None:
        response = self.client.get("/healthz")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json(), {"status": "ok"})

    def test_manual_draft_check_endpoint_passes_prompted_manual_starter(self) -> None:
        response = self.client.post(
            "/api/manual-draft-check",
            data={
                "payload_yaml": PROMPTED_MANUAL_PROCEDURE_EXAMPLE_YAML,
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertTrue(payload["ok"])
        self.assertEqual(payload["errors"], [])
        self.assertEqual(payload["warnings"], [])
        self.assertEqual(payload["summary"]["body_slide_count"], 4)
        self.assertEqual(payload["summary"]["section_break_count"], 1)
        self.assertEqual(payload["summary"]["procedure_slide_count"], 3)

    def test_manual_draft_check_endpoint_repairs_common_manual_indentation_drift(self) -> None:
        response = self.client.post(
            "/api/manual-draft-check",
            data={
                "payload_yaml": BROKEN_MANUAL_INDENTATION_REPORT_CONTENT_YAML,
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertTrue(payload["ok"])
        self.assertEqual(payload["errors"], [])
        self.assertIn(
            "Auto-corrected common manual YAML indentation drift before checking.",
            payload["warnings"],
        )
        self.assertIn(
            "The draft was re-indented automatically. Review the repaired YAML in the editor before generating.",
            payload["hints"],
        )
        self.assertTrue(payload["payload_yaml"].startswith("report_content:\n  title_slide:"))
        repaired = parse_yaml_text(payload["payload_yaml"])
        self.assertEqual(
            repaired["report_content"]["title_slide"]["slots"]["doc_title"],
            "Autoreport PowerPoint User Guide",
        )
        self.assertEqual(len(repaired["report_content"]["slides"]), 2)
        self.assertEqual(payload["summary"]["body_slide_count"], 2)
        self.assertEqual(payload["summary"]["procedure_slide_count"], 1)

    def test_manual_draft_check_endpoint_flags_invalid_manual_pattern_and_step_number(self) -> None:
        response = self.client.post(
            "/api/manual-draft-check",
            data={
                "payload_yaml": INVALID_MANUAL_PATTERN_REPORT_CONTENT_YAML,
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertFalse(payload["ok"])
        self.assertEqual(
            payload["message"],
            "Draft checker found 1 blocking issue(s).",
        )
        self.assertTrue(
            any(
                "Field 'slides[1].pattern_id' uses unsupported manual pattern 'image.manual.step'."
                in error
                and "text_image.manual.procedure.one" in error
                for error in payload["errors"]
            )
        )
        self.assertIn(
            "Field 'slides[1].slots.step_no' uses '2-1'. Prefer dotted numbering such as '2.1'.",
            payload["warnings"],
        )
        self.assertIn(
            "Allowed body pattern_id values: text.manual.section_break, text_image.manual.procedure.one, text_image.manual.procedure.two, text_image.manual.procedure.three.",
            payload["hints"],
        )

    def test_manual_draft_check_endpoint_catches_common_ai_confusions(self) -> None:
        def remove_body_pattern_id(report_content: dict[str, object]) -> None:
            slides = report_content["slides"]
            slides[1].pop("pattern_id", None)

        def mismatch_pattern_to_two_images(report_content: dict[str, object]) -> None:
            slides = report_content["slides"]
            slides[1]["slots"]["image_2"] = "image_2"

        def add_image_to_section_break(report_content: dict[str, object]) -> None:
            slides = report_content["slides"]
            slides[0]["slots"]["image_1"] = "image_1"

        def remove_all_images_from_procedure(report_content: dict[str, object]) -> None:
            slides = report_content["slides"]
            slides[1]["slots"].pop("image_1", None)

        def change_title_pattern_id(report_content: dict[str, object]) -> None:
            report_content["title_slide"]["pattern_id"] = "cover.editorial"

        def change_contents_pattern_id(report_content: dict[str, object]) -> None:
            report_content["contents_slide"]["pattern_id"] = "contents.editorial"

        cases = [
            (
                "missing body pattern id",
                build_mutated_manual_report_content_yaml(remove_body_pattern_id),
                [
                    "Field 'slides[1].pattern_id' is required for manual body slides.",
                    "text_image.manual.procedure.one",
                ],
            ),
            (
                "one-image pattern used with two images",
                build_mutated_manual_report_content_yaml(
                    mismatch_pattern_to_two_images
                ),
                [
                    "Field 'slides[1].pattern_id' is 'text_image.manual.procedure.one', but this slide defines 2 image_* slot(s), so it should use 'text_image.manual.procedure.two'."
                ],
            ),
            (
                "section break carries image slot",
                build_mutated_manual_report_content_yaml(add_image_to_section_break),
                [
                    "Field 'slides[0].pattern_id' is 'text.manual.section_break', so the slide must not define image_* slots."
                ],
            ),
            (
                "procedure slide has no image slots",
                build_mutated_manual_report_content_yaml(
                    remove_all_images_from_procedure
                ),
                [
                    "Field 'slides[1].pattern_id' is 'text_image.manual.procedure.one', but this slide defines 0 image_* slot(s), so it should use 'text.manual.section_break'."
                ],
            ),
            (
                "title slide pattern drifts from cover.manual",
                build_mutated_manual_report_content_yaml(change_title_pattern_id),
                [
                    "Field 'title_slide.pattern_id' must be exactly 'cover.manual'."
                ],
            ),
            (
                "contents slide pattern drifts from contents.manual",
                build_mutated_manual_report_content_yaml(change_contents_pattern_id),
                [
                    "Field 'contents_slide.pattern_id' must be exactly 'contents.manual' when a contents slide is present."
                ],
            ),
        ]

        for case_name, payload_yaml, expected_errors in cases:
            with self.subTest(case_name):
                response = self.client.post(
                    "/api/manual-draft-check",
                    data={
                        "payload_yaml": payload_yaml,
                        "built_in": "autoreport_manual",
                    },
                )

                self.assertEqual(response.status_code, 200)
                payload = response.json()
                self.assertFalse(payload["ok"])
                self.assertEqual(
                    payload["message"],
                    f"Draft checker found {len(payload['errors'])} blocking issue(s).",
                )
                for expected_error in expected_errors:
                    self.assertTrue(
                        any(expected_error in error for error in payload["errors"]),
                        payload,
                    )

    def test_manual_draft_check_endpoint_warns_for_section_number_without_period(self) -> None:
        def remove_section_number_period(report_content: dict[str, object]) -> None:
            slides = report_content["slides"]
            slides[0]["slots"]["section_no"] = "1"

        response = self.client.post(
            "/api/manual-draft-check",
            data={
                "payload_yaml": build_mutated_manual_report_content_yaml(
                    remove_section_number_period
                ),
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertTrue(payload["ok"])
        self.assertEqual(payload["errors"], [])
        self.assertEqual(
            payload["message"],
            "Draft checker passed with warnings. Review the numbering and rule hints before generating.",
        )
        self.assertIn(
            "Field 'slides[0].slots.section_no' should usually end with a trailing period, such as '2.'.",
            payload["warnings"],
        )

    def test_manual_slide_style_endpoint_appends_section_break(self) -> None:
        response = self.client.post(
            "/api/manual-slide-style",
            data={
                "payload_yaml": PROMPTED_MANUAL_PROCEDURE_EXAMPLE_YAML,
                "preset_id": "manual.section-break",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertTrue(
            payload["payload_yaml"].startswith("# Paste this brief into another AI")
        )
        updated = parse_yaml_text(payload["payload_yaml"])
        slides = updated["report_content"]["slides"]
        self.assertEqual(slides[-1]["pattern_id"], "text.manual.section_break")
        self.assertEqual(slides[-1]["slots"]["section_no"], "2.")
        self.assertEqual(payload["added_slide"]["preset_id"], "manual.section-break")
        self.assertEqual(slides[-1]["slots"]["section_title"], "New Section Title")
        self.assertEqual(
            payload["added_slide"]["slide_title"],
            "2. New Section Title",
        )

    def test_manual_slide_style_endpoint_appends_three_image_procedure_with_next_refs(self) -> None:
        response = self.client.post(
            "/api/manual-slide-style",
            data={
                "payload_yaml": PROMPTED_MANUAL_PROCEDURE_EXAMPLE_YAML,
                "preset_id": "manual.procedure.three",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        updated = parse_yaml_text(payload["payload_yaml"])
        last_slide = updated["report_content"]["slides"][-1]
        self.assertEqual(last_slide["pattern_id"], "text_image.manual.procedure.three")
        self.assertEqual(last_slide["slots"]["step_no"], "1.4")
        self.assertEqual(last_slide["slots"]["image_1"], "image_7")
        self.assertEqual(last_slide["slots"]["image_2"], "image_8")
        self.assertEqual(last_slide["slots"]["image_3"], "image_9")
        self.assertEqual(payload["added_slide"]["preset_id"], "manual.procedure.three")

        compile_response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": payload["payload_yaml"],
                "image_manifest": "[]",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(compile_response.status_code, 200)
        compile_payload = compile_response.json()
        self.assertEqual(
            compile_payload["slide_previews"][-1]["pattern_id"],
            "text_image.manual.procedure.three",
        )
        self.assertEqual(
            compile_payload["required_images"][-3:],
            [
                {
                    "slide_no": 5,
                    "slide_title": "1.4 New Procedure Step",
                    "alias": "image_1",
                    "ref": "image_7",
                },
                {
                    "slide_no": 5,
                    "slide_title": "1.4 New Procedure Step",
                    "alias": "image_2",
                    "ref": "image_8",
                },
                {
                    "slide_no": 5,
                    "slide_title": "1.4 New Procedure Step",
                    "alias": "image_3",
                    "ref": "image_9",
                },
            ],
        )

    def test_manual_slide_style_endpoint_rejects_non_manual_template(self) -> None:
        response = self.client.post(
            "/api/manual-slide-style",
            data={
                "payload_yaml": VALID_REPORT_CONTENT_YAML,
                "preset_id": "manual.section-break",
                "built_in": "autoreport_editorial",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertEqual(response.json()["message"], "Slide insertion failed.")
        self.assertIn(
            "Manual slide styles can only be added when the built-in manual starter is active.",
            response.json()["errors"],
        )

    def test_manual_slide_delete_endpoint_removes_contents_slide(self) -> None:
        response = self.client.post(
            "/api/manual-slide-delete",
            data={
                "payload_yaml": PROMPTED_MANUAL_PROCEDURE_EXAMPLE_YAML,
                "source_kind": "contents_slide",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        updated = parse_yaml_text(payload["payload_yaml"])
        self.assertNotIn("contents_slide", updated["report_content"])
        self.assertEqual(payload["deleted_slide"]["source_kind"], "contents_slide")
        self.assertEqual(payload["deleted_slide"]["slide_title"], "Contents")

        compile_response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": payload["payload_yaml"],
                "image_manifest": "[]",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(compile_response.status_code, 200)
        compile_payload = compile_response.json()
        self.assertEqual(
            [preview["slide_title"] for preview in compile_payload["slide_previews"][:2]],
            [
                "Autoreport PowerPoint User Guide",
                "1. Review The Manual Starter",
            ],
        )
        self.assertEqual(compile_payload["slide_previews"][1]["source_kind"], "content_slide")

    def test_manual_slide_delete_endpoint_removes_requested_content_slide(self) -> None:
        response = self.client.post(
            "/api/manual-slide-delete",
            data={
                "payload_yaml": PROMPTED_MANUAL_PROCEDURE_EXAMPLE_YAML,
                "source_kind": "content_slide",
                "source_slide_index": "2",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        updated = parse_yaml_text(payload["payload_yaml"])
        slides = updated["report_content"]["slides"]
        self.assertEqual(len(slides), 3)
        self.assertEqual(payload["deleted_slide"]["source_slide_index"], 2)
        self.assertEqual(
            payload["deleted_slide"]["slide_title"],
            "1.1 Review The Starter Example",
        )
        self.assertNotIn(
            "Review The Starter Example",
            [slide["slots"].get("step_title") for slide in slides],
        )

    def test_manual_slide_delete_endpoint_rejects_last_content_slide(self) -> None:
        response = self.client.post(
            "/api/manual-slide-delete",
            data={
                "payload_yaml": build_single_slide_manual_procedure_example_yaml(),
                "source_kind": "content_slide",
                "source_slide_index": "1",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertEqual(response.json()["message"], "Slide deletion failed.")
        self.assertIn(
            "The manual draft must keep at least 1 content slide.",
            response.json()["errors"],
        )

    def test_compile_endpoint_returns_compiled_runtime_payload(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": VALID_AUTHORING_PAYLOAD_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["payload_kind"], "authoring")
        self.assertIn("report_payload:", response.json()["compiled_yaml"])
        self.assertIn("pattern_id: text.editorial", response.json()["compiled_yaml"])
        self.assertEqual(response.json()["slide_count"], 1)

    def test_compile_endpoint_normalizes_report_content_into_authoring_payload(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": VALID_REPORT_CONTENT_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["payload_kind"], "content")
        self.assertIn("authoring_payload:", response.json()["normalized_authoring_yaml"])
        self.assertIn("goal: What It Does", response.json()["normalized_authoring_yaml"])
        self.assertIn("report_payload:", response.json()["compiled_yaml"])
        self.assertEqual(response.json()["required_images"], [])
        self.assertEqual(
            [preview["slide_title"] for preview in response.json()["slide_previews"]],
            ["Autoreport", "Contents", "What It Does"],
        )

    def test_compile_endpoint_accepts_report_content_without_kind(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": VALID_REPORT_CONTENT_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertIn("kind: text", response.json()["normalized_authoring_yaml"])

    def test_compile_endpoint_accepts_fenced_report_content(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": f"```yaml\n{VALID_REPORT_CONTENT_YAML}\n```",
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["payload_kind"], "content")
        self.assertIn("authoring_payload:", response.json()["normalized_authoring_yaml"])

    def test_compile_endpoint_repairs_common_manual_indentation_drift(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": BROKEN_MANUAL_INDENTATION_REPORT_CONTENT_YAML,
                "image_manifest": "[]",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["payload_kind"], "content")
        self.assertTrue(payload["payload_yaml"].startswith("report_content:\n  title_slide:"))
        self.assertIn(
            "The draft was re-indented automatically. Review the repaired YAML in the editor before generating.",
            payload["hints"],
        )
        self.assertEqual(
            payload["required_images"],
            [
                {
                    "slide_no": 2,
                    "slide_title": "1.1 Review The Starter Example",
                    "alias": "image_1",
                    "ref": "image_1",
                }
            ],
        )
        self.assertEqual(
            [preview["slide_title"] for preview in payload["slide_previews"]],
            [
                "Autoreport PowerPoint User Guide",
                "Contents",
                "1. Review The Manual Starter",
                "1.1 Review The Starter Example",
            ],
        )

    def test_compile_endpoint_returns_manual_required_image_order(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": MANUAL_PROCEDURE_EXAMPLE_YAML,
                "image_manifest": "[]",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["payload_kind"], "content")
        self.assertIn(
            "pattern_id: text_image.manual.procedure.three",
            response.json()["compiled_yaml"],
        )
        self.assertEqual(
            response.json()["required_images"],
            [
                {
                    "slide_no": 2,
                    "slide_title": "1.1 Review The Starter Example",
                    "alias": "image_1",
                    "ref": "image_1",
                },
                {
                    "slide_no": 3,
                    "slide_title": "1.2 Customize The Draft",
                    "alias": "image_1",
                    "ref": "image_2",
                },
                {
                    "slide_no": 3,
                    "slide_title": "1.2 Customize The Draft",
                    "alias": "image_2",
                    "ref": "image_3",
                },
                {
                    "slide_no": 4,
                    "slide_title": "1.3 Generate The PowerPoint",
                    "alias": "image_1",
                    "ref": "image_4",
                },
                {
                    "slide_no": 4,
                    "slide_title": "1.3 Generate The PowerPoint",
                    "alias": "image_2",
                    "ref": "image_5",
                },
                {
                    "slide_no": 4,
                    "slide_title": "1.3 Generate The PowerPoint",
                    "alias": "image_3",
                    "ref": "image_6",
                },
            ],
        )
        self.assertEqual(
            [preview["pattern_id"] for preview in response.json()["slide_previews"]],
            [
                "cover.manual",
                "contents.manual",
                "text.manual.section_break",
                "text_image.manual.procedure.one",
                "text_image.manual.procedure.two",
                "text_image.manual.procedure.three",
            ],
        )
        self.assertEqual(
            response.json()["slide_previews"][3]["image_blocks"],
            [
                {
                    "slot_name": "text_image.image_1",
                    "label": "image_1",
                    "ref": "image_1",
                    "fit": "contain",
                    "uploaded": False,
                    "x_pct": 66.4,
                    "y_pct": 21.4,
                    "w_pct": 26.0,
                    "h_pct": 40.8,
                }
            ],
        )
        self.assertEqual(
            [
                {
                    "source_kind": preview["source_kind"],
                    "source_slide_index": preview["source_slide_index"],
                    "source_can_delete": preview["source_can_delete"],
                    "source_is_primary_preview": preview["source_is_primary_preview"],
                }
                for preview in response.json()["slide_previews"][:4]
            ],
            [
                {
                    "source_kind": "title_slide",
                    "source_slide_index": None,
                    "source_can_delete": False,
                    "source_is_primary_preview": True,
                },
                {
                    "source_kind": "contents_slide",
                    "source_slide_index": None,
                    "source_can_delete": True,
                    "source_is_primary_preview": True,
                },
                {
                    "source_kind": "content_slide",
                    "source_slide_index": 1,
                    "source_can_delete": True,
                    "source_is_primary_preview": True,
                },
                {
                    "source_kind": "content_slide",
                    "source_slide_index": 2,
                    "source_can_delete": True,
                    "source_is_primary_preview": True,
                },
            ],
        )

    def test_compile_endpoint_slide_previews_follow_runtime_continuation(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": build_long_manual_procedure_example_yaml(),
                "image_manifest": "[]",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 200)
        continued_previews = [
            preview
            for preview in response.json()["slide_previews"]
            if preview["slide_title"].startswith("1.1 Review The Starter Example")
        ]
        self.assertGreaterEqual(len(continued_previews), 2)
        self.assertEqual(
            continued_previews[0]["slide_title"],
            "1.1 Review The Starter Example",
        )
        self.assertTrue(
            all(
                preview["slide_title"] == "1.1 Review The Starter Example (cont.)"
                for preview in continued_previews[1:]
            )
        )
        self.assertEqual(len(continued_previews[0]["image_blocks"]), 1)
        self.assertTrue(
            all(preview["image_blocks"] == [] for preview in continued_previews[1:])
        )
        self.assertEqual(continued_previews[0]["source_kind"], "content_slide")
        self.assertEqual(continued_previews[0]["source_slide_index"], 2)
        self.assertTrue(continued_previews[0]["source_can_delete"])
        self.assertTrue(continued_previews[0]["source_is_primary_preview"])
        self.assertTrue(
            all(preview["source_slide_index"] == 2 for preview in continued_previews[1:])
        )
        self.assertTrue(
            all(not preview["source_can_delete"] for preview in continued_previews[1:])
        )
        self.assertTrue(
            all(
                not preview["source_is_primary_preview"]
                for preview in continued_previews[1:]
            )
        )
        self.assertIn("decorations", continued_previews[0])
        self.assertIn("font_size_pt", continued_previews[0]["text_blocks"][0])

    def test_compile_endpoint_rejects_image_backed_payloads_in_public_app(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": VALID_REPORT_CONTENT_WITH_IMAGE_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertEqual(response.json()["message"], "Payload validation failed.")
        self.assertIn(
            "The public web demo currently supports text and metrics slides only.",
            response.json()["errors"],
        )

    def test_compile_endpoint_rejects_non_empty_image_manifest_in_public_app(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": VALID_REPORT_CONTENT_YAML,
                "image_manifest": '[{"ref":"image_1","field_name":"image_1","filename":"workflow.png"}]',
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertIn(
            "The public web demo currently supports text and metrics slides only.",
            response.json()["errors"],
        )

    def test_compile_endpoint_rejects_mixed_partial_fence_ai_output(self) -> None:
        response = self.client.post(
            "/api/compile",
            data={
                "payload_yaml": MIXED_CHATGPT_STYLE_REPORT_CONTENT,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertEqual(response.json()["error_type"], "yaml_parse_error")
        self.assertIn("Mixed fenced and unfenced YAML content", response.json()["message"])

    def test_generate_endpoint_returns_clear_error_for_truncated_pattern_id(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": TRUNCATED_CLAUDE_REPORT_CONTENT_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertIn(
            "Field 'slides[0].slots' must be an object.",
            response.json()["errors"],
        )

    def test_generate_endpoint_returns_pptx_attachment_from_authoring_payload(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": VALID_AUTHORING_PAYLOAD_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.headers["content-type"], MEDIA_TYPE_PPTX)
        self.assertIn(
            'attachment; filename="autoreport_demo.pptx"',
            response.headers["content-disposition"],
        )
        presentation = Presentation(BytesIO(response.content))
        self.assertEqual(len(presentation.slides), 3)

    def test_generate_endpoint_accepts_report_content_draft(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": VALID_REPORT_CONTENT_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 200)
        presentation = Presentation(BytesIO(response.content))
        self.assertEqual(len(presentation.slides), 3)

    def test_generate_endpoint_accepts_manual_built_in_with_uploaded_images(self) -> None:
        manifest = [
            {
                "ref": f"image_{index}",
                "field_name": f"image_{index}",
                "filename": f"step_{index}.png",
            }
            for index in range(1, 7)
        ]
        files = [
            (f"image_{index}", (f"step_{index}.png", PNG_BYTES, "image/png"))
            for index in range(1, 7)
        ]
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": MANUAL_PROCEDURE_EXAMPLE_YAML,
                "image_manifest": json.dumps(manifest),
                "built_in": "autoreport_manual",
            },
            files=files,
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.headers["content-type"], MEDIA_TYPE_PPTX)
        presentation = Presentation(BytesIO(response.content))
        self.assertEqual(len(presentation.slides), 6)

    def test_generate_endpoint_reports_missing_manual_uploads(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": PROMPTED_MANUAL_PROCEDURE_EXAMPLE_YAML,
                "image_manifest": "[]",
                "built_in": "autoreport_manual",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertIn(
            "Slide 2 needs an uploaded image for ref 'image_1'. Upload the matching file below before generating, or replace the draft image note with a real file path/ref.",
            response.json()["errors"],
        )

    def test_generate_endpoint_rejects_duplicate_manual_image_refs(self) -> None:
        with warnings.catch_warnings(record=True) as caught:
            warnings.simplefilter("always", ResourceWarning)
            response = self.client.post(
                "/api/generate",
                data={
                    "payload_yaml": MANUAL_PROCEDURE_EXAMPLE_YAML,
                    "image_manifest": json.dumps(
                        [
                            {
                                "ref": "image_1",
                                "field_name": "upload_1",
                                "filename": "one.png",
                            },
                            {
                                "ref": "image_1",
                                "field_name": "upload_2",
                                "filename": "two.png",
                            },
                        ]
                    ),
                    "built_in": "autoreport_manual",
                },
                files=[
                    ("upload_1", ("one.png", PNG_BYTES, "image/png")),
                    ("upload_2", ("two.png", PNG_BYTES, "image/png")),
                ],
            )
            gc.collect()

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertIn(
            "Field 'image_manifest[1].ref' must be unique.",
            response.json()["errors"],
        )
        resource_warnings = [
            warning
            for warning in caught
            if issubclass(warning.category, ResourceWarning)
        ]
        self.assertEqual(resource_warnings, [])

    def test_generate_endpoint_rejects_manual_upload_refs_with_path_separators(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": MANUAL_PROCEDURE_EXAMPLE_YAML,
                "image_manifest": json.dumps(
                    [
                        {
                            "ref": "../escape",
                            "field_name": "image_1",
                            "filename": "step_1.png",
                        },
                    ]
                ),
                "built_in": "autoreport_manual",
            },
            files=[("image_1", ("step_1.png", PNG_BYTES, "image/png"))],
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertIn(
            "Field 'image_manifest[0].ref' must not contain path separators.",
            response.json()["errors"],
        )

    def test_generate_endpoint_rejects_invalid_manual_upload_file_type(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": MANUAL_PROCEDURE_EXAMPLE_YAML,
                "image_manifest": json.dumps(
                    [
                        {
                            "ref": "image_1",
                            "field_name": "image_1",
                            "filename": "step_1.gif",
                        },
                    ]
                ),
                "built_in": "autoreport_manual",
            },
            files=[("image_1", ("step_1.gif", b"GIF89a", "image/gif"))],
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertIn(
            "Field 'image_manifest[0].ref' has an unsupported file type.",
            response.json()["errors"],
        )

    def test_generate_endpoint_rejects_image_backed_authoring_payload_in_public_app(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": VALID_TEXT_IMAGE_AUTHORING_PAYLOAD_YAML,
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertIn(
            "The public web demo currently supports text and metrics slides only.",
            response.json()["errors"],
        )

    def test_generate_endpoint_rejects_invalid_image_manifest_json(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": VALID_REPORT_CONTENT_YAML,
                "image_manifest": "[",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertEqual(response.json()["message"], "Payload validation failed.")
        self.assertIn(
            "Field 'image_manifest' must be a valid JSON list.",
            response.json()["errors"],
        )

    def test_generate_endpoint_returns_parse_errors(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={"payload_yaml": "authoring_payload: [broken", "image_manifest": "[]"},
        )

        self.assertEqual(response.status_code, 400)
        self.assertEqual(response.json()["error_type"], "yaml_parse_error")
        self.assertIn("Failed to parse YAML:", response.json()["message"])

    def test_generate_endpoint_cleans_temp_dir_after_parse_error_with_upload(self) -> None:
        leaked_temp_dir = Path(tempfile.mkdtemp(prefix="autoreport-web-test-"))
        with patch("autoreport.web.app.tempfile.mkdtemp", return_value=str(leaked_temp_dir)):
            response = self.client.post(
                "/api/generate",
                data={
                    "payload_yaml": "authoring_payload: [broken",
                    "image_manifest": json.dumps(
                        [
                            {
                                "ref": "image_1",
                                "field_name": "image_1",
                                "filename": "step_1.png",
                            },
                        ]
                    ),
                    "built_in": "autoreport_manual",
                },
                files=[("image_1", ("step_1.png", PNG_BYTES, "image/png"))],
            )

        self.assertEqual(response.status_code, 400)
        self.assertFalse(leaked_temp_dir.exists())

    def test_generate_endpoint_returns_validation_errors(self) -> None:
        response = self.client.post(
            "/api/generate",
            data={
                "payload_yaml": """
authoring_payload:
  payload_version: autoreport.authoring.v1
  template_id: autoreport-editorial-v1
  title_slide:
    title: "  "
    subtitle: []
  contents:
    enabled: true
  slides: []
""".strip(),
                "image_manifest": "[]",
            },
        )

        self.assertEqual(response.status_code, 422)
        self.assertEqual(response.json()["error_type"], "validation_error")
        self.assertEqual(response.json()["message"], "Payload validation failed.")
        self.assertIn(
            "Field 'title_slide.title' must be a non-empty string.",
            response.json()["errors"],
        )

    def test_generate_endpoint_returns_generic_internal_errors(self) -> None:
        with patch(
            "autoreport.web.app.generate_report_from_mapping",
            side_effect=RuntimeError("boom"),
        ):
            response = self.client.post(
                "/api/generate",
                data={"payload_yaml": VALID_AUTHORING_PAYLOAD_YAML, "image_manifest": "[]"},
            )

        self.assertEqual(response.status_code, 500)
        self.assertEqual(response.json()["error_type"], "internal_error")
        self.assertEqual(
            response.json()["message"],
            "An unexpected internal error occurred.",
        )


if __name__ == "__main__":
    unittest.main()
