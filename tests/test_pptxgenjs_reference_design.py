"""Branch-local tests for PresentationGo reference analysis and the editorial v0.4 template."""

from __future__ import annotations

import shutil
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from autoreport.engine.generator import generate_report_from_mapping
from autoreport.template_flow import inspect_template_contract
from experiments.v04.prototypes.pptxgenjs_template_spike.presentationgo_reference_analysis import (
    REFERENCE_LIBRARY_DIR,
    collect_analysis,
)
from experiments.v04.prototypes.pptxgenjs_template_spike.demo_payloads import (
    build_demo_payload_document,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
GENERATED_DIR = (
    REPO_ROOT
    / "experiments"
    / "v04"
    / "prototypes"
    / "pptxgenjs_template_spike"
    / "generated"
)
EDITORIAL_TEMPLATE_PATH = GENERATED_DIR / "v04-datawave-editorial-template.pptx"


def _visible_slide_titles(presentation: Presentation) -> list[str]:
    titles: list[str] = []
    for slide in presentation.slides:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            titles.append(slide.shapes.title.text)
            continue
        fallback = "(untitled)"
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                fallback = text.splitlines()[0]
                break
        titles.append(fallback)
    return titles


class PptxGenJSReferenceDesignTestCase(unittest.TestCase):
    """Verify the local reference analysis and the prettier editorial template fixture."""

    @unittest.skipUnless(REFERENCE_LIBRARY_DIR.exists(), "Local PresentationGo library not available")
    def test_reference_analysis_covers_the_local_template_snapshot(self) -> None:
        analysis = collect_analysis()

        self.assertEqual(len(analysis["references"]), 4)
        self.assertEqual(
            analysis["recommended_translation"]["structure_reference"],
            "modern-business-16x9.potx",
        )
        self.assertEqual(
            analysis["recommended_translation"]["palette_reference"],
            "datawave-insights-16x9.potx",
        )

    def test_editorial_template_fixture_exists(self) -> None:
        self.assertTrue(EDITORIAL_TEMPLATE_PATH.exists(), msg=str(EDITORIAL_TEMPLATE_PATH))

    def test_editorial_template_inspection_exposes_title_and_text_patterns(self) -> None:
        contract = inspect_template_contract(template_path=EDITORIAL_TEMPLATE_PATH)

        self.assertEqual(contract.template_source, "user_template")
        self.assertTrue(contract.title_slide.pattern_id)
        self.assertIn("text", [pattern.kind for pattern in contract.slide_patterns])

    def test_editorial_template_prefers_one_primary_body_text_slot(self) -> None:
        contract = inspect_template_contract(template_path=EDITORIAL_TEMPLATE_PATH)
        text_pattern = next(
            pattern
            for pattern in contract.slide_patterns
            if pattern.kind == "text"
        )
        text_slots = [
            slot
            for slot in text_pattern.slots
            if slot.slot_type == "text"
        ]
        image_slots = [
            slot
            for slot in text_pattern.slots
            if slot.slot_type == "image"
        ]

        self.assertEqual(len(text_slots), 1)
        self.assertEqual(len(image_slots), 0)

    def test_editorial_template_generates_an_editorial_output(self) -> None:
        contract = inspect_template_contract(template_path=EDITORIAL_TEMPLATE_PATH)
        payload_document, image_refs = build_demo_payload_document(contract)
        temp_root = Path(tempfile.mkdtemp())
        try:
            output_path = temp_root / "editorial-output.pptx"
            generated_path = generate_report_from_mapping(
                payload_document,
                output_path=output_path,
                template_path=EDITORIAL_TEMPLATE_PATH,
                image_refs=image_refs,
            )
            presentation = Presentation(str(generated_path))
        finally:
            shutil.rmtree(temp_root, ignore_errors=True)

        self.assertEqual(generated_path, output_path)
        self.assertEqual(
            _visible_slide_titles(presentation),
            [
                "Autoreport",
                "Contents",
                "What It Does",
                "Adoption Snapshot",
            ],
        )


if __name__ == "__main__":
    unittest.main()
