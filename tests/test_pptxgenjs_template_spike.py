"""Branch-local tests for the v0.4 PptxGenJS template-authoring spike."""

from __future__ import annotations

import shutil
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from autoreport.engine.generator import generate_report_from_mapping
from autoreport.template_flow import inspect_template_contract
from experiments.v04.prototypes.pptxgenjs_template_spike.demo_payloads import (
    build_demo_payload_document,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
GENERATED_DIR = (
    REPO_ROOT
    / "experiments"
    / "v04"
    / "prototypes"
    / "pptxgenjs_template_spike"
    / "generated"
)
MINIMAL_TEMPLATE_PATH = GENERATED_DIR / "v04-minimal-text-template.pptx"
STACKED_TEMPLATE_PATH = GENERATED_DIR / "v04-stacked-text-template.pptx"
TEXT_IMAGE_TEMPLATE_PATH = GENERATED_DIR / "v04-text-image-template.pptx"


def _visible_slide_titles(presentation: Presentation) -> list[str]:
    titles: list[str] = []
    for slide in presentation.slides:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            titles.append(slide.shapes.title.text)
            continue

        fallback = "(untitled)"
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                fallback = text.splitlines()[0]
                break
        titles.append(fallback)
    return titles


class PptxGenJSTemplateSpikeTestCase(unittest.TestCase):
    """Verify the generated PptxGenJS fixtures against current Python surfaces."""

    def test_fixture_files_exist(self) -> None:
        for template_path in (
            MINIMAL_TEMPLATE_PATH,
            STACKED_TEMPLATE_PATH,
            TEXT_IMAGE_TEMPLATE_PATH,
        ):
            self.assertTrue(template_path.exists(), msg=str(template_path))

    def test_minimal_text_template_inspection_exposes_title_and_text_patterns(self) -> None:
        contract = inspect_template_contract(template_path=MINIMAL_TEMPLATE_PATH)

        self.assertEqual(contract.template_source, "user_template")
        self.assertTrue(contract.title_slide.pattern_id)
        self.assertIn("text", [pattern.kind for pattern in contract.slide_patterns])

    def test_stacked_text_template_contract_has_two_stack_text_slots(self) -> None:
        contract = inspect_template_contract(template_path=STACKED_TEMPLATE_PATH)
        text_pattern = next(
            pattern
            for pattern in contract.slide_patterns
            if pattern.kind == "text"
        )
        text_slots = [
            slot
            for slot in text_pattern.slots
            if slot.slot_type == "text"
        ]

        self.assertEqual(len(text_slots), 2)
        self.assertEqual(
            [slot.orientation for slot in text_slots],
            ["stack", "stack"],
        )

    def test_text_image_template_now_exports_a_text_image_pattern(self) -> None:
        contract = inspect_template_contract(template_path=TEXT_IMAGE_TEMPLATE_PATH)
        pattern_kinds = [pattern.kind for pattern in contract.slide_patterns]
        text_image_pattern = next(
            pattern
            for pattern in contract.slide_patterns
            if pattern.kind == "text_image"
        )
        image_slots = [
            slot
            for slot in text_image_pattern.slots
            if slot.slot_type == "image"
        ]

        self.assertEqual(pattern_kinds, ["text", "metrics", "text_image"])
        self.assertEqual(len(image_slots), 1)

    def test_template_ids_are_stable_across_reinspection(self) -> None:
        for template_path in (
            MINIMAL_TEMPLATE_PATH,
            STACKED_TEMPLATE_PATH,
            TEXT_IMAGE_TEMPLATE_PATH,
        ):
            first = inspect_template_contract(template_path=template_path)
            second = inspect_template_contract(template_path=template_path)
            self.assertEqual(first.template_id, second.template_id)

    def test_generate_report_from_mapping_creates_output_for_minimal_template(self) -> None:
        contract = inspect_template_contract(template_path=MINIMAL_TEMPLATE_PATH)
        payload_document, image_refs = build_demo_payload_document(contract)
        temp_root = Path(tempfile.mkdtemp())
        try:
            output_path = temp_root / "minimal-template-output.pptx"
            generated_path = generate_report_from_mapping(
                payload_document,
                output_path=output_path,
                template_path=MINIMAL_TEMPLATE_PATH,
                image_refs=image_refs,
            )
            presentation = Presentation(str(generated_path))
        finally:
            shutil.rmtree(temp_root, ignore_errors=True)

        self.assertEqual(generated_path, output_path)
        self.assertEqual(
            _visible_slide_titles(presentation),
            [
                "Autoreport",
                "Contents",
                "What It Does",
                "Adoption Snapshot",
            ],
        )


if __name__ == "__main__":
    unittest.main()
