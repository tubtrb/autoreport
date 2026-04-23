from __future__ import annotations

import base64
import csv
from datetime import datetime
from hashlib import sha1
import json
from pathlib import Path
import subprocess
import sys
from time import perf_counter
from typing import Any, Callable

import httpx
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import yaml

from autoreport.web.app import (
    MANUAL_PUBLIC_TEMPLATE_NAME,
    _build_manual_draft_check,
    app as public_app,
)
from autoreport.web.manual_ai_yaml import (
    append_manual_ai_coercion_feedback,
    parse_public_payload_yaml,
)
from tests.verif_test.catalog import (
    PreparedSample,
    REVIEW_CASE_IDS,
    load_prompt_pack_samples,
    prepare_suite_samples,
    validate_yaml_candidate_against_manifest,
)
from tests.verif_test.chatgpt import (
    capture_chatgpt_transport_diagnostics,
    canonical_profile_dir,
    collect_chatgpt_response,
    export_chatgpt_transport_artifacts,
    get_chatgpt_transport_metrics,
    inspect_chatgpt_session,
    normalize_yaml_candidate,
)


REPO_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_BASE_URL = "http://127.0.0.1:8000"
DEFAULT_PREFLIGHT_MODULES = ("tests.test_web_app", "tests.test_web_serve")
PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jXioAAAAASUVORK5CYII="
)
ReviewDecision = dict[str, dict[str, str]]
RETRYABLE_HTTP_STATUS_CODES = {0, 500, 502, 503, 504}
HTTP_STAGE_MAX_RETRIES = 2
RUN_SUMMARY_EXTENSION_KEYS = (
    "runbook",
    "planned_chat_count",
    "completed_chat_count",
    "chunk_results",
    "cooldown_schedule_applied",
    "guard_trip_reason",
    "guard_trip_stage",
    "single_session_browser_relaunches",
)

RUN_GATE_DEFINITIONS = (
    {
        "key": "preflight",
        "label": "Narrow unittest preflight",
    },
    {
        "key": "healthz",
        "label": "Local public app health check",
    },
    {
        "key": "chatgpt_session",
        "label": "ChatGPT session readiness",
    },
)

CASE_STEP_DEFINITIONS = (
    {
        "key": "prompt_artifacts",
        "label": "Starter YAML / prompt artifacts",
    },
    {
        "key": "ai_transport",
        "label": "ChatGPT prompt submit + reply receive",
    },
    {
        "key": "yaml_extract",
        "label": "YAML candidate extract",
    },
    {
        "key": "checker",
        "label": "/api/manual-draft-check",
    },
    {
        "key": "generate",
        "label": "/api/generate",
    },
    {
        "key": "pptx_inspection",
        "label": "PPTX inspection",
    },
)


def execute_suite_run(
    *,
    suite_name: str,
    session: str,
    mode: str,
    output_root: Path,
    base_url: str = DEFAULT_BASE_URL,
    sample_count: int | None = None,
    run_dir: Path | None = None,
    transport: Callable[[PreparedSample], str] | None = None,
    run_preflight: bool = True,
    send_wait_seconds: float = 0.8,
    poll_seconds: float = 3.0,
    max_polls: int = 20,
    python_executable: Path | None = None,
    session_check_only: bool = False,
    progress: Callable[[str], None] | None = None,
) -> Path:
    samples = () if session_check_only else prepare_suite_samples(suite_name, sample_count=sample_count)
    return execute_samples(
        samples=samples,
        suite_name=suite_name,
        session=session,
        mode=mode,
        output_root=output_root,
        base_url=base_url,
        run_dir=run_dir,
        transport=transport,
        run_preflight=run_preflight,
        send_wait_seconds=send_wait_seconds,
        poll_seconds=poll_seconds,
        max_polls=max_polls,
        python_executable=python_executable,
        session_check_only=session_check_only,
        progress=progress,
    )


def execute_prompt_pack_run(
    *,
    prompt_pack_path: Path,
    count: int,
    session: str,
    mode: str,
    output_dir: Path,
    base_url: str = DEFAULT_BASE_URL,
    transport: Callable[[PreparedSample], str] | None = None,
    send_wait_seconds: float = 0.8,
    poll_seconds: float = 3.0,
    max_polls: int = 20,
    python_executable: Path | None = None,
    progress: Callable[[str], None] | None = None,
) -> Path:
    samples = load_prompt_pack_samples(prompt_pack_path, sample_count=count)
    return execute_samples(
        samples=samples,
        suite_name="prompt-pack",
        session=session,
        mode=mode,
        output_root=output_dir.parent,
        base_url=base_url,
        run_dir=output_dir,
        transport=transport,
        run_preflight=True,
        send_wait_seconds=send_wait_seconds,
        poll_seconds=poll_seconds,
        max_polls=max_polls,
        python_executable=python_executable,
        prompt_pack_path=prompt_pack_path,
        progress=progress,
    )


def execute_samples(
    *,
    samples: tuple[PreparedSample, ...],
    suite_name: str,
    session: str,
    mode: str,
    output_root: Path,
    base_url: str = DEFAULT_BASE_URL,
    run_dir: Path | None = None,
    transport: Callable[[PreparedSample], str] | None = None,
    run_preflight: bool = True,
    send_wait_seconds: float = 0.8,
    poll_seconds: float = 3.0,
    max_polls: int = 20,
    python_executable: Path | None = None,
    prompt_pack_path: Path | None = None,
    session_check_only: bool = False,
    progress: Callable[[str], None] | None = None,
) -> Path:
    if mode not in {"http", "local"}:
        raise RuntimeError("mode must be one of: http, local")

    resolved_python = python_executable or Path(sys.executable)
    started_at = _timestamp()
    resolved_output_root = output_root.resolve()
    resolved_run_dir = run_dir.resolve() if run_dir is not None else (
        resolved_output_root / _build_run_id(suite_name)
    )
    resolved_run_dir.mkdir(parents=True, exist_ok=True)
    root_events_path = resolved_run_dir / "events.jsonl"

    run_config = {
        "run_id": resolved_run_dir.name,
        "suite": suite_name,
        "mode": mode,
        "session": session,
        "base_url": base_url,
        "session_check_only": session_check_only,
        "expected_profile_dir": str(canonical_profile_dir(session)),
        "sample_count": len(samples),
        "started_at": started_at,
        "prompt_pack_path": str(prompt_pack_path.resolve()) if prompt_pack_path else None,
        "preflight_modules": list(DEFAULT_PREFLIGHT_MODULES),
        "review_case_ids": list(sorted(REVIEW_CASE_IDS)),
        "cases": [sample.manifest for sample in samples],
    }
    _write_json(resolved_run_dir / "run-config.json", run_config)
    _write_json(resolved_run_dir / "run-metadata.json", run_config)
    _append_event(
        root_events_path,
        run_id=resolved_run_dir.name,
        case_id="_run",
        stage="run",
        status="started",
        duration_ms=0.0,
        message=f"Starting {suite_name} with {len(samples)} sample(s).",
    )
    _report_progress(
        progress,
        f"[manual-ai] starting suite={suite_name} mode={mode} cases={len(samples)} run_dir={resolved_run_dir}",
    )

    preflight_result: dict[str, Any] | None = None
    if run_preflight and not session_check_only:
        stage_started = perf_counter()
        _report_progress(
            progress,
            "[manual-ai] running preflight: python -m unittest tests.test_web_app tests.test_web_serve",
        )
        preflight_result = run_preflight_tests(
            python_executable=resolved_python,
            run_dir=resolved_run_dir,
            modules=DEFAULT_PREFLIGHT_MODULES,
        )
        _append_event(
            root_events_path,
            run_id=resolved_run_dir.name,
            case_id="_run",
            stage="preflight",
            status="success" if preflight_result["ok"] else "failure",
            duration_ms=_duration_ms(stage_started),
            failure_class="preflight_failure" if not preflight_result["ok"] else None,
            message=preflight_result["message"],
            artifact=str((resolved_run_dir / "preflight.stdout.log").relative_to(resolved_run_dir)),
        )
        _report_progress(progress, f"[manual-ai] preflight: {preflight_result['message']}")
        if not preflight_result["ok"]:
            summary = build_summary_payload(
                run_config=run_config,
                cases=[],
                preflight=preflight_result,
                health=None,
                session_check=None,
                started_at=started_at,
                finished_at=_timestamp(),
            )
            summary["overall_status"] = "FAIL"
            summary["message"] = "Preflight tests failed."
            write_summary_artifacts(resolved_run_dir, summary)
            _report_progress(progress, f"[manual-ai] completed with FAIL: {resolved_run_dir}")
            return resolved_run_dir

    health_result: dict[str, Any] | None = None
    if mode == "http" and not session_check_only:
        stage_started = perf_counter()
        _report_progress(progress, f"[manual-ai] checking local app health: {base_url.rstrip('/')}/healthz")
        health_result = check_health(base_url)
        _append_event(
            root_events_path,
            run_id=resolved_run_dir.name,
            case_id="_run",
            stage="healthz",
            status="success" if health_result["ok"] else "failure",
            duration_ms=_duration_ms(stage_started),
            failure_class="ai_transport_failure" if not health_result["ok"] else None,
            message=health_result["message"],
        )
        _report_progress(progress, f"[manual-ai] healthz: {health_result['message']}")
        if not health_result["ok"]:
            summary = build_summary_payload(
                run_config=run_config,
                cases=[],
                preflight=preflight_result,
                health=health_result,
                session_check=None,
                started_at=started_at,
                finished_at=_timestamp(),
            )
            summary["overall_status"] = "FAIL"
            summary["message"] = "Local public app health check failed."
            write_summary_artifacts(resolved_run_dir, summary)
            _report_progress(progress, f"[manual-ai] completed with FAIL: {resolved_run_dir}")
            return resolved_run_dir

    session_check: dict[str, Any] | None = None
    if transport is None or session_check_only:
        stage_started = perf_counter()
        _report_progress(progress, f"[manual-ai] checking ChatGPT session: {session}")
        session_check_kwargs: dict[str, Any] = {
            "session": session,
            "artifact_dir": resolved_run_dir,
        }
        if progress is not None:
            session_check_kwargs["progress"] = progress
        session_check = inspect_chatgpt_session(**session_check_kwargs)
        _append_event(
            root_events_path,
            run_id=resolved_run_dir.name,
            case_id="_run",
            stage="chatgpt_session",
            status="success" if session_check["ok"] else "failure",
            duration_ms=_duration_ms(stage_started),
            failure_class="ai_transport_failure" if not session_check["ok"] else None,
            message=session_check["message"],
        )
        _report_progress(progress, f"[manual-ai] session: {session_check['message']}")
        if session_check_only:
            summary = build_summary_payload(
                run_config=run_config,
                cases=[],
                preflight=None,
                health=None,
                session_check=session_check,
                started_at=started_at,
                finished_at=_timestamp(),
            )
            write_summary_artifacts(resolved_run_dir, summary)
            _report_progress(progress, f"[manual-ai] completed session-check-only: {resolved_run_dir}")
            return resolved_run_dir
        if not session_check["ok"]:
            summary = build_summary_payload(
                run_config=run_config,
                cases=[],
                preflight=preflight_result,
                health=health_result,
                session_check=session_check,
                started_at=started_at,
                finished_at=_timestamp(),
            )
            summary["overall_status"] = "FAIL"
            summary["message"] = "ChatGPT browser session is not ready."
            write_summary_artifacts(resolved_run_dir, summary)
            _report_progress(progress, f"[manual-ai] completed with FAIL: {resolved_run_dir}")
            return resolved_run_dir

    case_rows: list[dict[str, Any]] = []
    transport_fn = transport or (
        lambda sample: collect_chatgpt_response(
            session=session,
            prompt=sample.prompt,
            expected_manifest=sample.manifest,
            send_wait_seconds=send_wait_seconds,
            poll_seconds=poll_seconds,
            max_polls=max_polls,
        )
    )

    for index, sample in enumerate(samples, start=1):
        case_dir = resolved_run_dir / f"{index:03d}-{sample.case_id}"
        case_dir.mkdir(parents=True, exist_ok=True)
        case_events_path = case_dir / "events.jsonl"
        _write_text(case_dir / "starter.yaml", sample.starter_yaml)
        _write_text(case_dir / "prompt.txt", sample.prompt)
        _write_json(case_dir / "case-manifest.json", sample.manifest)
        _report_progress(
            progress,
            f"[manual-ai] case {index}/{len(samples)} {sample.case_id}: sending prompt",
        )

        row: dict[str, Any] = {
            "index": index,
            "sample_id": case_dir.name,
            "case_id": sample.case_id,
            "prompt_id": sample.prompt_id,
            "label": sample.label,
            "case_dir": str(case_dir),
            "review_required": bool(sample.expected_review_required),
            "review_decision": None,
            "status": "success",
            "failure_class": None,
            "checker_ok": False,
            "checker_message": "",
            "blocking_issue_count": 0,
            "warning_count": 0,
            "generate_ok": False,
            "pptx_openable": False,
            "slide_count": 0,
            "shape_count": 0,
            "image_count": 0,
            "artifact_pptx": None,
            "last_completed_step": None,
            "stop_step": None,
            "transport_attempts": 0,
            "session_relaunches": 0,
            "no_sandbox_detected": False,
            "selected_page_url": "",
            "selected_page_title": "",
            "recent_response_failures": [],
            "step_statuses": _build_case_step_statuses(),
        }
        _mark_case_step_status(row, "prompt_artifacts", "success")
        _append_case_event(
            root_events_path,
            case_events_path,
            run_id=resolved_run_dir.name,
            case_id=sample.case_id,
            stage="prompt_artifacts",
            status="success",
            duration_ms=0.0,
            message="Wrote starter.yaml, prompt.txt, and case-manifest.json.",
            artifact=str((case_dir / "case-manifest.json").relative_to(resolved_run_dir)),
        )
        try:
            stage_started = perf_counter()
            raw_turn_text = transport_fn(sample)
            _write_text(case_dir / "ai-raw.txt", raw_turn_text)
            _write_text(case_dir / "raw-turn.txt", raw_turn_text)
            if transport is None:
                transport_state = capture_chatgpt_transport_diagnostics(session)
                _write_json(case_dir / "transport-state.json", transport_state)
                export_chatgpt_transport_artifacts(session, case_dir)
                _apply_transport_metrics(row, transport_state)
            _mark_case_step_status(row, "ai_transport", "success")
            _append_case_event(
                root_events_path,
                case_events_path,
                run_id=resolved_run_dir.name,
                case_id=sample.case_id,
                stage="ai_transport",
                status="success",
                duration_ms=_duration_ms(stage_started),
                message="Received ChatGPT response.",
                artifact=str((case_dir / "ai-raw.txt").relative_to(resolved_run_dir)),
            )
        except Exception as exc:
            row["status"] = "failure"
            row["failure_class"] = "ai_transport_failure"
            row["checker_message"] = str(exc)
            _mark_case_step_status(row, "ai_transport", "failure")
            _write_text(case_dir / "failure.txt", str(exc))
            transport_state_artifact: str | None = None
            if transport is None:
                transport_state = capture_chatgpt_transport_diagnostics(session)
                _write_json(case_dir / "transport-state.json", transport_state)
                export_chatgpt_transport_artifacts(session, case_dir)
                _apply_transport_metrics(row, transport_state)
                transport_state_artifact = str(
                    (case_dir / "transport-state.json").relative_to(resolved_run_dir)
                )
            _append_case_event(
                root_events_path,
                case_events_path,
                run_id=resolved_run_dir.name,
                case_id=sample.case_id,
                stage="ai_transport",
                status="failure",
                duration_ms=0.0,
                failure_class=row["failure_class"],
                message=str(exc),
                artifact=transport_state_artifact or str((case_dir / "failure.txt").relative_to(resolved_run_dir)),
            )
            case_rows.append(row)
            continue

        stage_started = perf_counter()
        yaml_candidate = normalize_yaml_candidate(raw_turn_text)
        _write_text(case_dir / "yaml-candidate.yaml", yaml_candidate)
        if not yaml_candidate or "report_content:" not in yaml_candidate:
            row["status"] = "failure"
            row["failure_class"] = "yaml_extract_failure"
            row["checker_message"] = "The AI reply did not contain a usable report_content YAML block."
            _mark_case_step_status(row, "yaml_extract", "failure")
            _write_text(case_dir / "failure.txt", row["checker_message"])
            _append_case_event(
                root_events_path,
                case_events_path,
                run_id=resolved_run_dir.name,
                case_id=sample.case_id,
                stage="yaml_extract",
                status="failure",
                duration_ms=_duration_ms(stage_started),
                failure_class=row["failure_class"],
                message=row["checker_message"],
                artifact=str((case_dir / "yaml-candidate.yaml").relative_to(resolved_run_dir)),
            )
            case_rows.append(row)
            continue
        _mark_case_step_status(row, "yaml_extract", "success")
        _append_case_event(
            root_events_path,
            case_events_path,
            run_id=resolved_run_dir.name,
            case_id=sample.case_id,
            stage="yaml_extract",
            status="success",
            duration_ms=_duration_ms(stage_started),
            message="Extracted a report_content YAML candidate.",
            artifact=str((case_dir / "yaml-candidate.yaml").relative_to(resolved_run_dir)),
        )

        stage_started = perf_counter()
        checker_payload = _run_checker_with_retry(
            mode=mode,
            base_url=base_url,
            yaml_candidate=yaml_candidate,
        )
        _write_json(case_dir / "checker.json", checker_payload)
        manifest_guard_payload = validate_yaml_candidate_against_manifest(
            yaml_candidate,
            manifest=sample.manifest,
        )
        _write_json(case_dir / "manifest-guard.json", manifest_guard_payload)
        row["checker_ok"] = bool(checker_payload.get("ok"))
        row["checker_message"] = str(checker_payload.get("message", ""))
        summary = checker_payload.get("summary", {}) or {}
        row["blocking_issue_count"] = int(summary.get("blocking_issue_count", 0))
        row["warning_count"] = int(summary.get("warning_count", 0))
        if is_yaml_extract_failure(checker_payload):
            row["status"] = "failure"
            row["failure_class"] = "yaml_extract_failure"
            _mark_case_step_status(row, "checker", "failure")
        elif not row["checker_ok"] or row["blocking_issue_count"] > 0:
            row["status"] = "failure"
            row["failure_class"] = "checker_failure"
            _mark_case_step_status(row, "checker", "failure")
        elif not bool(manifest_guard_payload.get("ok")):
            row["status"] = "failure"
            row["failure_class"] = "checker_failure"
            row["checker_ok"] = False
            row["checker_message"] = str(manifest_guard_payload.get("message", ""))
            row["blocking_issue_count"] = max(
                row["blocking_issue_count"],
                len(list(manifest_guard_payload.get("errors", []) or [])),
            )
            _mark_case_step_status(row, "checker", "failure")
        else:
            _mark_case_step_status(row, "checker", "success")
        _append_case_event(
            root_events_path,
            case_events_path,
            run_id=resolved_run_dir.name,
            case_id=sample.case_id,
            stage="checker",
            status="success" if row["failure_class"] is None else "failure",
            duration_ms=_duration_ms(stage_started),
            failure_class=row["failure_class"],
            message=row["checker_message"],
            artifact=str(
                (
                    case_dir / (
                        "manifest-guard.json"
                        if row["failure_class"] == "checker_failure"
                        and not bool(manifest_guard_payload.get("ok"))
                        else "checker.json"
                    )
                ).relative_to(resolved_run_dir)
            ),
        )
        if row["failure_class"] is not None:
            case_rows.append(row)
            continue

        stage_started = perf_counter()
        generate_payload, pptx_bytes = _run_generate_with_retry(
            mode=mode,
            base_url=base_url,
            payload_yaml=yaml_candidate,
            image_ref_count=sample.image_ref_count,
        )
        _write_json(case_dir / "generate.json", generate_payload)
        row["generate_ok"] = bool(generate_payload.get("ok"))
        if not row["generate_ok"] or not pptx_bytes:
            row["status"] = "failure"
            row["failure_class"] = "generate_failure"
            _mark_case_step_status(row, "generate", "failure")
            _append_case_event(
                root_events_path,
                case_events_path,
                run_id=resolved_run_dir.name,
                case_id=sample.case_id,
                stage="generate",
                status="failure",
                duration_ms=_duration_ms(stage_started),
                failure_class=row["failure_class"],
                message=str(generate_payload.get("message", "Generation failed.")),
                artifact=str((case_dir / "generate.json").relative_to(resolved_run_dir)),
            )
            case_rows.append(row)
            continue
        artifact_path = case_dir / "artifact.pptx"
        artifact_path.write_bytes(pptx_bytes)
        row["artifact_pptx"] = str(artifact_path)
        _mark_case_step_status(row, "generate", "success")
        _append_case_event(
            root_events_path,
            case_events_path,
            run_id=resolved_run_dir.name,
            case_id=sample.case_id,
            stage="generate",
            status="success",
            duration_ms=_duration_ms(stage_started),
            message="Generated PPTX artifact.",
            artifact=str(artifact_path.relative_to(resolved_run_dir)),
        )

        stage_started = perf_counter()
        inspection_payload = inspect_pptx_artifact(artifact_path)
        _write_json(case_dir / "pptx-inspection.json", inspection_payload)
        row["pptx_openable"] = bool(inspection_payload.get("openable"))
        row["slide_count"] = int(inspection_payload.get("slide_count", 0))
        row["shape_count"] = int(inspection_payload.get("shape_count", 0))
        row["image_count"] = int(inspection_payload.get("image_count", 0))
        if not row["pptx_openable"]:
            row["status"] = "failure"
            row["failure_class"] = "pptx_inspection_failure"
            _mark_case_step_status(row, "pptx_inspection", "failure")
        else:
            _mark_case_step_status(row, "pptx_inspection", "success")
        _append_case_event(
            root_events_path,
            case_events_path,
            run_id=resolved_run_dir.name,
            case_id=sample.case_id,
            stage="pptx_inspection",
            status="success" if row["failure_class"] is None else "failure",
            duration_ms=_duration_ms(stage_started),
            failure_class=row["failure_class"],
            message=str(inspection_payload.get("message", "PPTX inspection complete.")),
            artifact=str((case_dir / "pptx-inspection.json").relative_to(resolved_run_dir)),
        )
        case_rows.append(row)

    summary = build_summary_payload(
        run_config=run_config,
        cases=case_rows,
        preflight=preflight_result,
        health=health_result,
        session_check=session_check,
        started_at=started_at,
        finished_at=_timestamp(),
    )
    write_summary_artifacts(resolved_run_dir, summary)
    _append_event(
        root_events_path,
        run_id=resolved_run_dir.name,
        case_id="_run",
        stage="run",
        status="success" if summary["overall_status"] != "FAIL" else "failure",
        duration_ms=0.0,
        message=f"Run completed with status {summary['overall_status']}.",
    )
    _report_progress(
        progress,
        f"[manual-ai] completed with {summary['overall_status']}: {resolved_run_dir}",
    )
    return resolved_run_dir


def _apply_transport_metrics(row: dict[str, Any], payload: dict[str, Any] | None) -> None:
    payload = payload or {}
    row["transport_attempts"] = int(payload.get("transport_attempts", 0) or 0)
    row["session_relaunches"] = int(payload.get("session_relaunches", 0) or 0)
    row["no_sandbox_detected"] = bool(payload.get("no_sandbox_detected"))
    row["selected_page_url"] = str(payload.get("selected_page_url", "") or "")
    row["selected_page_title"] = str(payload.get("selected_page_title", "") or "")
    row["recent_response_failures"] = list(payload.get("recent_response_failures", []) or [])


def _is_retryable_http_failure(payload: dict[str, Any], *, status_key: str) -> bool:
    if payload.get("ok"):
        return False
    try:
        status_code = int(payload.get(status_key, 200) or 0)
    except (TypeError, ValueError):
        status_code = 0
    return status_code in RETRYABLE_HTTP_STATUS_CODES


def _run_checker_with_retry(
    *,
    mode: str,
    base_url: str,
    yaml_candidate: str,
) -> dict[str, Any]:
    attempts: list[dict[str, Any]] = []
    total_attempts = HTTP_STAGE_MAX_RETRIES + 1
    for attempt in range(1, total_attempts + 1):
        try:
            payload = (
                run_manual_draft_check_http(base_url=base_url, yaml_candidate=yaml_candidate)
                if mode == "http"
                else run_manual_draft_check_local(yaml_candidate=yaml_candidate)
            )
        except Exception as exc:
            payload = {
                "ok": False,
                "message": f"Checker request failed: {exc}",
                "errors": [],
                "warnings": [],
                "hints": [],
                "summary": {},
                "_http_status": 0,
            }
        attempts.append(
            {
                "attempt": attempt,
                "status_code": int(payload.get("_http_status", 0) or 0),
                "message": str(payload.get("message", "")),
            }
        )
        payload["attempts"] = attempt
        payload["retry_trace"] = list(attempts)
        if not _is_retryable_http_failure(payload, status_key="_http_status"):
            return payload
        if attempt == total_attempts:
            return payload
    raise RuntimeError("unreachable")


def _run_generate_with_retry(
    *,
    mode: str,
    base_url: str,
    payload_yaml: str,
    image_ref_count: int,
) -> tuple[dict[str, Any], bytes | None]:
    attempts: list[dict[str, Any]] = []
    total_attempts = HTTP_STAGE_MAX_RETRIES + 1
    for attempt in range(1, total_attempts + 1):
        try:
            payload, pptx_bytes = (
                run_generate_http(
                    base_url=base_url,
                    payload_yaml=payload_yaml,
                    image_ref_count=image_ref_count,
                )
                if mode == "http"
                else run_generate_local(
                    payload_yaml=payload_yaml,
                    image_ref_count=image_ref_count,
                )
            )
        except Exception as exc:
            payload, pptx_bytes = (
                {
                    "ok": False,
                    "status_code": 0,
                    "message": f"Generate request failed: {exc}",
                    "errors": [],
                },
                None,
            )
        attempts.append(
            {
                "attempt": attempt,
                "status_code": int(payload.get("status_code", 0) or 0),
                "message": str(payload.get("message", "")),
            }
        )
        payload["attempts"] = attempt
        payload["retry_trace"] = list(attempts)
        if not _is_retryable_http_failure(payload, status_key="status_code"):
            return payload, pptx_bytes
        if attempt == total_attempts:
            return payload, pptx_bytes
    raise RuntimeError("unreachable")


def build_summary_payload(
    *,
    run_config: dict[str, Any],
    cases: list[dict[str, Any]],
    preflight: dict[str, Any] | None,
    health: dict[str, Any] | None,
    session_check: dict[str, Any] | None,
    started_at: str,
    finished_at: str,
    review_decisions: ReviewDecision | None = None,
    summary_overrides: dict[str, Any] | None = None,
) -> dict[str, Any]:
    decisions = review_decisions or {}
    updated_cases = [dict(case) for case in cases]
    for case in updated_cases:
        decision = decisions.get(case["case_id"])
        if decision is not None:
            case["review_decision"] = decision["decision"]
            case["review_note"] = decision.get("note", "")
            case["review_recorded_at"] = decision["recorded_at"]
    failure_counts: dict[str, int] = {}
    for case in updated_cases:
        failure_class = case.get("failure_class")
        if not failure_class:
            continue
        failure_counts[str(failure_class)] = failure_counts.get(str(failure_class), 0) + 1
    run_gate_stats = _build_run_gate_stats(
        preflight=preflight,
        health=health,
        session_check=session_check,
    )
    case_step_stats = _build_case_step_stats(
        cases=updated_cases,
        total_cases=int(run_config.get("sample_count", len(updated_cases))),
    )
    primary_blocker = _select_primary_blocker(
        run_gate_stats=run_gate_stats,
        case_step_stats=case_step_stats,
    )
    review_queue = select_review_queue(updated_cases, suite_name=str(run_config["suite"]))
    overall_status = compute_overall_status(
        updated_cases,
        preflight=preflight,
        health=health,
        session_check=session_check,
    )
    payload = {
        "run_id": run_config["run_id"],
        "suite": run_config["suite"],
        "mode": run_config["mode"],
        "session": run_config["session"],
        "base_url": run_config["base_url"],
        "session_check_only": bool(run_config.get("session_check_only")),
        "expected_profile_dir": run_config.get("expected_profile_dir"),
        "sample_count": int(run_config.get("sample_count", len(updated_cases))),
        "started_at": started_at,
        "finished_at": finished_at,
        "overall_status": overall_status,
        "launcher": (session_check or {}).get("launcher"),
        "no_sandbox_detected": bool((session_check or {}).get("no_sandbox_detected")),
        "selected_page_url": str((session_check or {}).get("selected_page_url", "") or ""),
        "selected_page_title": str((session_check or {}).get("selected_page_title", "") or ""),
        "recent_response_failures": list((session_check or {}).get("recent_response_failures", []) or []),
        "preflight": preflight,
        "health": health,
        "session_check": session_check,
        "run_gate_stats": run_gate_stats,
        "case_step_stats": case_step_stats,
        "primary_blocker": primary_blocker,
        "failure_counts": dict(sorted(failure_counts.items())),
        "review_decisions": decisions,
        "review_queue": review_queue,
        "cases": updated_cases,
        "message": _overall_message(
            overall_status,
            session_check_only=bool(run_config.get("session_check_only")),
        ),
    }
    if summary_overrides:
        payload.update(summary_overrides)
    return payload


def _build_case_step_statuses() -> dict[str, str]:
    return {spec["key"]: "pending" for spec in CASE_STEP_DEFINITIONS}


def _mark_case_step_status(row: dict[str, Any], step_key: str, status: str) -> None:
    step_statuses = row.setdefault("step_statuses", _build_case_step_statuses())
    step_statuses[step_key] = status
    if status == "success":
        row["last_completed_step"] = step_key
    elif status == "failure":
        row["stop_step"] = step_key


def _build_run_gate_stats(
    *,
    preflight: dict[str, Any] | None,
    health: dict[str, Any] | None,
    session_check: dict[str, Any] | None,
) -> list[dict[str, Any]]:
    payloads = {
        "preflight": preflight,
        "healthz": health,
        "chatgpt_session": session_check,
    }
    stats: list[dict[str, Any]] = []
    for spec in RUN_GATE_DEFINITIONS:
        payload = payloads[spec["key"]]
        attempted = payload is not None
        ok = bool(payload.get("ok")) if attempted else None
        stats.append(
            {
                "step": spec["key"],
                "label": spec["label"],
                "attempted": attempted,
                "status": "PASS" if ok else ("FAIL" if attempted else "SKIP"),
                "reason": payload.get("reason") if attempted else None,
                "message": payload.get("message") if attempted else "",
            }
        )
    return stats


def _build_case_step_stats(
    *,
    cases: list[dict[str, Any]],
    total_cases: int,
) -> list[dict[str, Any]]:
    stats: list[dict[str, Any]] = []
    for spec in CASE_STEP_DEFINITIONS:
        success_cases = 0
        failure_cases = 0
        for case in cases:
            step_status = (case.get("step_statuses") or {}).get(spec["key"], "pending")
            if step_status == "success":
                success_cases += 1
            elif step_status == "failure":
                failure_cases += 1
        attempted_cases = success_cases + failure_cases
        blocked_cases = max(total_cases - attempted_cases, 0)
        failure_rate_pct = round((failure_cases / attempted_cases) * 100.0, 1) if attempted_cases else None
        reach_rate_pct = round((attempted_cases / total_cases) * 100.0, 1) if total_cases else None
        stats.append(
            {
                "step": spec["key"],
                "label": spec["label"],
                "total_cases": total_cases,
                "attempted_cases": attempted_cases,
                "success_cases": success_cases,
                "failure_cases": failure_cases,
                "blocked_cases": blocked_cases,
                "failure_rate_pct": failure_rate_pct,
                "reach_rate_pct": reach_rate_pct,
            }
        )
    return stats


def _select_primary_blocker(
    *,
    run_gate_stats: list[dict[str, Any]],
    case_step_stats: list[dict[str, Any]],
) -> dict[str, Any] | None:
    for gate in run_gate_stats:
        if gate["status"] == "FAIL":
            return {
                "scope": "run_gate",
                "step": gate["step"],
                "label": gate["label"],
                "reason": gate.get("reason"),
                "message": gate.get("message", ""),
            }

    failing_steps = [item for item in case_step_stats if item["failure_cases"] > 0]
    if not failing_steps:
        return None
    primary = max(
        failing_steps,
        key=lambda item: (item["failure_cases"], item["failure_rate_pct"] or 0.0, -item["blocked_cases"]),
    )
    return {
        "scope": "case_step",
        "step": primary["step"],
        "label": primary["label"],
        "failure_cases": primary["failure_cases"],
        "attempted_cases": primary["attempted_cases"],
        "failure_rate_pct": primary["failure_rate_pct"],
        "blocked_cases": primary["blocked_cases"],
    }


def write_summary_artifacts(run_dir: Path, summary: dict[str, Any]) -> None:
    _write_json(run_dir / "summary.json", summary)
    _write_summary_csv(run_dir / "summary.csv", summary["cases"])
    summary_md = render_summary_markdown(summary)
    _write_text(run_dir / "summary.md", summary_md)
    _write_text(run_dir / "summary.txt", summary_md)
    if summary.get("session_check") is not None:
        _write_json(run_dir / "session-check.json", summary["session_check"])
        _write_text(
            run_dir / "session-check.md",
            render_session_check_markdown(summary["session_check"]),
        )
    review_queue_payload = {
        "run_id": summary["run_id"],
        "overall_status": summary["overall_status"],
        "items": summary["review_queue"],
    }
    _write_json(run_dir / "review-queue.json", review_queue_payload)
    _write_text(run_dir / "review-queue.md", render_review_queue_markdown(review_queue_payload))


def record_visual_review(
    *,
    run_dir: Path,
    case_id: str,
    decision: str,
    note: str,
) -> dict[str, Any]:
    if decision not in {"pass", "fail"}:
        raise RuntimeError("decision must be one of: pass, fail")

    summary_path = run_dir / "summary.json"
    if not summary_path.exists():
        raise RuntimeError(f"Missing summary.json under {run_dir}")
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    cases = summary.get("cases")
    if not isinstance(cases, list):
        raise RuntimeError("summary.json does not contain cases.")

    matching_case = next(
        (
            case
            for case in cases
            if case.get("case_id") == case_id and case.get("review_required")
        ),
        None,
    )
    if matching_case is None:
        raise RuntimeError(f"Case '{case_id}' is not waiting for visual review in {run_dir}.")

    decisions = summary.get("review_decisions", {}) or {}
    decisions[case_id] = {
        "decision": decision,
        "note": note,
        "recorded_at": _timestamp(),
    }
    refreshed = build_summary_payload(
        run_config={
            "run_id": summary["run_id"],
            "suite": summary["suite"],
            "mode": summary["mode"],
            "session": summary["session"],
            "base_url": summary["base_url"],
            "session_check_only": summary.get("session_check_only", False),
            "expected_profile_dir": summary.get("expected_profile_dir"),
        },
        cases=cases,
        preflight=summary.get("preflight"),
        health=summary.get("health"),
        session_check=summary.get("session_check"),
        started_at=summary["started_at"],
        finished_at=_timestamp(),
        review_decisions=decisions,
        summary_overrides={
            key: summary.get(key)
            for key in RUN_SUMMARY_EXTENSION_KEYS
            if key in summary
        },
    )
    _write_json(run_dir / "review-decisions.json", decisions)
    write_summary_artifacts(run_dir, refreshed)
    return refreshed


def recheck_saved_artifacts(*, artifact_dir: Path) -> dict[str, Any]:
    resolved_artifact_dir = artifact_dir.resolve()
    sample_dirs = sorted(
        path
        for path in resolved_artifact_dir.iterdir()
        if path.is_dir() and (path / "yaml-candidate.yaml").exists()
    )
    if not sample_dirs:
        raise RuntimeError(
            f"No sample folders with yaml-candidate.yaml were found under: {resolved_artifact_dir}"
        )

    rows: list[dict[str, Any]] = []
    counts: dict[str, int] = {}
    for sample_dir in sample_dirs:
        candidate_text = (sample_dir / "yaml-candidate.yaml").read_text(encoding="utf-8")
        checker_payload = run_manual_draft_check_local(yaml_candidate=candidate_text)
        category = classify_recheck_category(checker_payload)
        counts[category] = counts.get(category, 0) + 1
        rows.append(
            {
                "sample_dir": str(sample_dir),
                "category": category,
                "repaired": bool(checker_payload.get("payload_yaml")),
                "message": checker_payload.get("message", ""),
                "errors": checker_payload.get("errors", []),
                "warnings": checker_payload.get("warnings", []),
                "summary": checker_payload.get("summary", {}),
            }
        )

    payload = {
        "artifact_dir": str(resolved_artifact_dir),
        "sample_count": len(rows),
        "category_counts": dict(sorted(counts.items())),
        "results": rows,
    }
    _write_json(resolved_artifact_dir / "recheck-summary.json", payload)
    lines = [
        "Manual YAML Recheck Summary",
        f"Artifact directory: {resolved_artifact_dir}",
        f"Collected samples: {len(rows)}",
        "",
        "Category counts:",
    ]
    for category, count in sorted(counts.items()):
        lines.append(f"- {category}: {count}")
    lines.append("")
    lines.append("Per-sample results:")
    for index, row in enumerate(rows, start=1):
        repaired_flag = " repaired" if row["repaired"] else ""
        lines.append(
            f"- #{index:03d} {Path(row['sample_dir']).name}: {row['category']}{repaired_flag} | {row['message']}"
        )
    summary_text = "\n".join(lines) + "\n"
    _write_text(resolved_artifact_dir / "recheck-summary.txt", summary_text)
    return payload


def run_preflight_tests(
    *,
    python_executable: Path,
    run_dir: Path,
    modules: tuple[str, ...] = DEFAULT_PREFLIGHT_MODULES,
    timeout_seconds: float = 120.0,
) -> dict[str, Any]:
    command = [str(python_executable), "-m", "unittest", *modules]
    stdout_path = run_dir / "preflight.stdout.log"
    stderr_path = run_dir / "preflight.stderr.log"
    try:
        completed = subprocess.run(
            command,
            cwd=REPO_ROOT,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            check=False,
            timeout=timeout_seconds,
        )
    except subprocess.TimeoutExpired as exc:
        _write_text(stdout_path, exc.stdout or "")
        _write_text(stderr_path, exc.stderr or "")
        return {
            "ok": False,
            "command": " ".join(command),
            "exit_code": None,
            "stdout_path": str(stdout_path),
            "stderr_path": str(stderr_path),
            "message": f"Preflight tests timed out after {timeout_seconds:.0f}s.",
        }
    _write_text(stdout_path, completed.stdout or "")
    _write_text(stderr_path, completed.stderr or "")
    return {
        "ok": completed.returncode == 0,
        "command": " ".join(command),
        "exit_code": completed.returncode,
        "stdout_path": str(stdout_path),
        "stderr_path": str(stderr_path),
        "message": (
            "Preflight tests passed."
            if completed.returncode == 0
            else f"Preflight tests failed with exit code {completed.returncode}."
        ),
    }


def check_health(base_url: str) -> dict[str, Any]:
    health_url = f"{base_url.rstrip('/')}/healthz"
    try:
        response = httpx.get(health_url, timeout=15.0)
        payload = response.json()
    except Exception as exc:
        return {
            "ok": False,
            "message": f"Health check request failed: {exc}",
            "status_code": 0,
        }
    return {
        "ok": response.status_code == 200 and payload.get("status") == "ok",
        "message": f"Health check returned HTTP {response.status_code}.",
        "status_code": response.status_code,
        "payload": payload,
    }


def run_manual_draft_check_local(*, yaml_candidate: str) -> dict[str, Any]:
    try:
        raw_data, coercion_result = parse_public_payload_yaml(
            yaml_candidate,
            built_in=MANUAL_PUBLIC_TEMPLATE_NAME,
        )
    except yaml.YAMLError as exc:
        return {
            "ok": False,
            "error_type": "yaml_parse_error",
            "message": f"Failed to parse YAML: {exc}",
            "errors": [],
            "warnings": [],
            "hints": [],
            "summary": {
                "payload_kind": "unknown",
                "body_slide_count": 0,
                "section_break_count": 0,
                "procedure_slide_count": 0,
                "blocking_issue_count": 1,
                "warning_count": 0,
            },
            "_http_status": 400,
        }
    payload = {
        **_build_manual_draft_check(raw_data, built_in=MANUAL_PUBLIC_TEMPLATE_NAME),
        "_http_status": 200,
    }
    if coercion_result is not None:
        payload = append_manual_ai_coercion_feedback(
            payload,
            coercion=coercion_result,
        )
    return payload


def run_manual_draft_check_http(*, base_url: str, yaml_candidate: str) -> dict[str, Any]:
    checker_url = f"{base_url.rstrip('/')}/api/manual-draft-check"
    try:
        response = httpx.post(
            checker_url,
            data={
                "payload_yaml": yaml_candidate,
                "built_in": MANUAL_PUBLIC_TEMPLATE_NAME,
            },
            timeout=30.0,
        )
    except Exception as exc:
        return {
            "ok": False,
            "message": f"Checker request failed: {exc}",
            "errors": [],
            "warnings": [],
            "hints": [],
            "summary": {},
            "_http_status": 0,
        }
    try:
        payload = response.json()
    except ValueError:
        payload = {
            "ok": False,
            "message": f"Checker HTTP {response.status_code} returned invalid JSON.",
            "errors": [],
            "warnings": [],
            "hints": [],
            "summary": {},
        }
    payload["_http_status"] = response.status_code
    return payload


def run_generate_local(
    *,
    payload_yaml: str,
    image_ref_count: int,
) -> tuple[dict[str, Any], bytes | None]:
    manifest, files = build_placeholder_uploads(image_ref_count)
    with TestClient(public_app) as client:
        response = client.post(
            "/api/generate",
            data={
                "payload_yaml": payload_yaml,
                "image_manifest": json.dumps(manifest),
                "built_in": MANUAL_PUBLIC_TEMPLATE_NAME,
            },
            files=files,
        )
    if response.status_code != 200:
        payload = response.json()
        return (
            {
                "ok": False,
                "status_code": response.status_code,
                "error_type": payload.get("error_type"),
                "message": payload.get("message"),
                "errors": payload.get("errors", []),
            },
            None,
        )
    return (
        {
            "ok": True,
            "status_code": response.status_code,
            "content_type": response.headers.get("content-type", ""),
            "content_disposition": response.headers.get("content-disposition", ""),
            "size_bytes": len(response.content),
            "message": "Generated PPTX successfully.",
        },
        response.content,
    )


def run_generate_http(
    *,
    base_url: str,
    payload_yaml: str,
    image_ref_count: int,
) -> tuple[dict[str, Any], bytes | None]:
    manifest, files = build_placeholder_uploads(image_ref_count)
    generate_url = f"{base_url.rstrip('/')}/api/generate"
    try:
        with httpx.Client(timeout=60.0) as client:
            response = client.post(
                generate_url,
                data={
                    "payload_yaml": payload_yaml,
                    "image_manifest": json.dumps(manifest),
                    "built_in": MANUAL_PUBLIC_TEMPLATE_NAME,
                },
                files=files,
            )
    except Exception as exc:
        return (
            {
                "ok": False,
                "status_code": 0,
                "message": f"Generate request failed: {exc}",
                "errors": [],
            },
            None,
        )
    if response.status_code != 200:
        try:
            payload = response.json()
        except ValueError:
            payload = {
                "error_type": "generate_http_error",
                "message": f"Generate HTTP {response.status_code} returned invalid JSON.",
                "errors": [],
            }
        return (
            {
                "ok": False,
                "status_code": response.status_code,
                "error_type": payload.get("error_type"),
                "message": payload.get("message"),
                "errors": payload.get("errors", []),
            },
            None,
        )
    return (
        {
            "ok": True,
            "status_code": response.status_code,
            "content_type": response.headers.get("content-type", ""),
            "content_disposition": response.headers.get("content-disposition", ""),
            "size_bytes": len(response.content),
            "message": "Generated PPTX successfully.",
        },
        response.content,
    )


def build_placeholder_uploads(
    image_ref_count: int,
) -> tuple[list[dict[str, str]], list[tuple[str, tuple[str, bytes, str]]]]:
    manifest: list[dict[str, str]] = []
    files: list[tuple[str, tuple[str, bytes, str]]] = []
    for index in range(1, image_ref_count + 1):
        ref = f"image_{index}"
        filename = f"{ref}.png"
        manifest.append(
            {
                "ref": ref,
                "field_name": ref,
                "filename": filename,
            }
        )
        files.append((ref, (filename, PNG_BYTES, "image/png")))
    return manifest, files


def inspect_pptx_artifact(artifact_path: Path) -> dict[str, Any]:
    try:
        presentation = Presentation(str(artifact_path))
    except Exception as exc:
        return {
            "openable": False,
            "message": f"Failed to open PPTX: {exc}",
            "slide_count": 0,
            "shape_count": 0,
            "image_count": 0,
            "file_size_bytes": artifact_path.stat().st_size if artifact_path.exists() else 0,
            "slides": [],
        }

    slides: list[dict[str, Any]] = []
    total_shape_count = 0
    total_image_count = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_text = _extract_slide_title(slide)
        shape_count = len(slide.shapes)
        image_count = sum(
            1
            for shape in slide.shapes
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
        )
        total_shape_count += shape_count
        total_image_count += image_count
        slides.append(
            {
                "slide_no": slide_index,
                "title": title_text,
                "title_sha1": sha1(title_text.encode("utf-8")).hexdigest()[:12],
                "shape_count": shape_count,
                "image_count": image_count,
            }
        )

    return {
        "openable": True,
        "message": "PPTX inspection complete.",
        "slide_count": len(slides),
        "shape_count": total_shape_count,
        "image_count": total_image_count,
        "file_size_bytes": artifact_path.stat().st_size,
        "slides": slides,
    }


def select_review_queue(cases: list[dict[str, Any]], *, suite_name: str) -> list[dict[str, Any]]:
    required_ids = {"01_one_image_canary"}
    if suite_name in {"regression", "full", "release-gate"}:
        required_ids.add("05_balanced_canary")
    if suite_name in {"full", "release-gate"}:
        required_ids.add("10_full_family_canary")

    queue: list[dict[str, Any]] = []
    for case_id in sorted(required_ids):
        matching_case = next(
            (
                case
                for case in cases
                if case.get("case_id") == case_id
                and case.get("review_required")
                and case.get("failure_class") is None
            ),
            None,
        )
        if matching_case is None:
            continue
        if matching_case.get("review_decision"):
            continue
        queue.append(
            {
                "case_id": matching_case["case_id"],
                "label": matching_case["label"],
                "artifact_pptx": matching_case.get("artifact_pptx"),
                "slide_count": matching_case.get("slide_count", 0),
            }
        )
    return queue


def compute_overall_status(
    cases: list[dict[str, Any]],
    *,
    preflight: dict[str, Any] | None,
    health: dict[str, Any] | None,
    session_check: dict[str, Any] | None,
) -> str:
    if preflight is not None and not preflight.get("ok"):
        return "FAIL"
    if health is not None and not health.get("ok"):
        return "FAIL"
    if session_check is not None and not session_check.get("ok"):
        return "FAIL"
    if any(case.get("failure_class") for case in cases):
        return "FAIL"
    if any(case.get("review_decision") == "fail" for case in cases):
        return "FAIL"
    if any(case.get("review_required") and not case.get("review_decision") for case in cases):
        return "REVIEW"
    return "PASS"


def is_yaml_extract_failure(checker_payload: dict[str, Any]) -> bool:
    message = str(checker_payload.get("message", ""))
    error_type = str(checker_payload.get("error_type", ""))
    http_status = int(checker_payload.get("_http_status", 200))
    return (
        error_type == "yaml_parse_error"
        or "Failed to parse YAML:" in message
        or http_status == 400
    )


def classify_recheck_category(checker_payload: dict[str, Any]) -> str:
    if is_yaml_extract_failure(checker_payload):
        return "yaml-parse-failure"
    if checker_payload.get("ok") and not checker_payload.get("errors"):
        return "manual-pass"
    if checker_payload.get("errors"):
        return "manual-fail"
    return "unknown"


def render_summary_markdown(summary: dict[str, Any]) -> str:
    lines = [
        "# Manual AI Regression Summary",
        "",
        f"- Run ID: `{summary['run_id']}`",
        f"- Suite: `{summary['suite']}`",
        f"- Mode: `{summary['mode']}`",
        f"- Planned Cases: `{summary.get('sample_count', len(summary.get('cases', [])))}`",
        f"- Overall Status: `{summary['overall_status']}`",
        f"- Started At: `{summary['started_at']}`",
        f"- Finished At: `{summary['finished_at']}`",
        "",
    ]
    if summary.get("primary_blocker"):
        blocker = summary["primary_blocker"]
        lines.append("## Primary Blocker")
        lines.append("")
        lines.append(f"- Scope: `{blocker['scope']}`")
        lines.append(f"- Step: `{blocker['step']}`")
        lines.append(f"- Label: `{blocker['label']}`")
        if blocker.get("reason"):
            lines.append(f"- Reason: `{blocker['reason']}`")
        if blocker.get("failure_cases") is not None:
            lines.append(
                f"- Failed Cases: `{blocker['failure_cases']}/{blocker['attempted_cases']}`"
            )
        if blocker.get("failure_rate_pct") is not None:
            lines.append(f"- Failure Rate: `{blocker['failure_rate_pct']}%`")
        if blocker.get("blocked_cases") is not None:
            lines.append(f"- Blocked Cases: `{blocker['blocked_cases']}`")
        if blocker.get("message"):
            lines.append(f"- Message: {blocker['message']}")
        lines.append("")
    if summary.get("runbook"):
        lines.append("## Release Gate")
        lines.append("")
        lines.append(f"- Runbook: `{summary['runbook']}`")
        lines.append(f"- Planned Chats: `{summary.get('planned_chat_count', 0)}`")
        lines.append(f"- Completed Chats: `{summary.get('completed_chat_count', 0)}`")
        lines.append(
            f"- Single Session Browser Relaunches: `{summary.get('single_session_browser_relaunches', 0)}`"
        )
        guard_reason = str(summary.get("guard_trip_reason", "") or "")
        guard_stage = str(summary.get("guard_trip_stage", "") or "")
        lines.append(f"- Guard Trip Stage: `{guard_stage or 'none'}`")
        lines.append(f"- Guard Trip Reason: `{guard_reason or 'none'}`")
        cooldowns = summary.get("cooldown_schedule_applied") or []
        lines.append(f"- Cooldowns Applied: `{len(cooldowns)}`")
        lines.append("")
    if summary.get("chunk_results"):
        lines.append("## Chunk Results")
        lines.append("")
        for chunk in summary["chunk_results"]:
            lines.append(
                f"- `{chunk.get('name', '')}`: `{chunk.get('status', 'unknown')}` | completed={chunk.get('completed_chat_count', 0)}/{chunk.get('planned_chat_count', 0)} | cooldown_after={chunk.get('cooldown_after_seconds', 0)}s"
            )
            if chunk.get("guard_trip_reason"):
                lines.append(f"  - guard: {chunk['guard_trip_reason']}")
        lines.append("")
    if summary.get("failure_counts"):
        lines.append("## Failure Counts")
        lines.append("")
        for failure_class, count in summary["failure_counts"].items():
            lines.append(f"- `{failure_class}`: {count}")
        lines.append("")
    if summary.get("run_gate_stats"):
        lines.append("## Run Gates")
        lines.append("")
        for item in summary["run_gate_stats"]:
            line = f"- `{item['step']}`: `{item['status']}`"
            if item.get("reason"):
                line += f" | reason=`{item['reason']}`"
            lines.append(line)
        lines.append("")
    if summary.get("case_step_stats"):
        lines.append("## Case Step Stats")
        lines.append("")
        for item in summary["case_step_stats"]:
            failure_rate = "n/a" if item["failure_rate_pct"] is None else f"{item['failure_rate_pct']}%"
            reach_rate = "n/a" if item["reach_rate_pct"] is None else f"{item['reach_rate_pct']}%"
            lines.append(
                f"- `{item['step']}`: attempted={item['attempted_cases']}/{item['total_cases']} | failed={item['failure_cases']} | blocked={item['blocked_cases']} | failure_rate={failure_rate} | reach_rate={reach_rate}"
            )
        lines.append("")
    if summary.get("preflight"):
        preflight = summary["preflight"]
        lines.append("## Preflight")
        lines.append("")
        lines.append(f"- Status: `{'PASS' if preflight['ok'] else 'FAIL'}`")
        lines.append(f"- Command: `{preflight['command']}`")
        lines.append("")
    if summary.get("session_check"):
        session_check = summary["session_check"]
        lines.append("## ChatGPT Session")
        lines.append("")
        lines.append(f"- Status: `{'PASS' if session_check['ok'] else 'FAIL'}`")
        lines.append(f"- Reason: `{session_check['reason']}`")
        if session_check.get("launcher"):
            lines.append(f"- Launcher: `{session_check['launcher']}`")
        if session_check.get("browser_pid") is not None:
            lines.append(f"- Browser PID: `{session_check['browser_pid']}`")
        lines.append(f"- Expected Profile: `{session_check['expected_profile_dir']}`")
        if session_check.get("actual_profile_dir"):
            lines.append(f"- Actual Profile: `{session_check['actual_profile_dir']}`")
        if session_check.get("page_url"):
            lines.append(f"- Page URL: `{session_check['page_url']}`")
        if session_check.get("page_title"):
            lines.append(f"- Page Title: `{session_check['page_title']}`")
        lines.append(
            f"- No Sandbox Detected: `{bool(session_check.get('no_sandbox_detected'))}`"
        )
        if session_check.get("manual_intervention_used"):
            lines.append(
                f"- Manual Intervention: `{session_check.get('manual_intervention_elapsed_seconds', 0.0)}`s / `{session_check.get('manual_intervention_timeout_seconds', 0.0)}`s"
            )
        if session_check.get("recent_response_failures"):
            lines.append(
                f"- Recent Response Failures: `{len(session_check['recent_response_failures'])}`"
            )
        lines.append(f"- Message: {session_check['message']}")
        lines.append("")
    if summary.get("session_check_only"):
        lines.append("## Cases")
        lines.append("")
        lines.append("No case execution ran. This run only validated the ChatGPT session readiness contract.")
        lines.append("")
        return "\n".join(lines)
    if summary.get("review_queue"):
        lines.append("## Review Queue")
        lines.append("")
        for item in summary["review_queue"]:
            lines.append(
                f"- `{item['case_id']}`: review [artifact]({item['artifact_pptx']}) after automated checks."
            )
        lines.append("")
    lines.append("## Cases")
    lines.append("")
    if not summary["cases"]:
        lines.append("No case execution ran.")
    else:
        for case in summary["cases"]:
            status = case["failure_class"] or case.get("review_decision") or "ok"
            lines.append(
                f"- `{case['sample_id']}` `{case['case_id']}`: `{status}` | stop={case.get('stop_step') or 'completed'} | last={case.get('last_completed_step') or 'none'} | transport_attempts={case.get('transport_attempts', 0)} | session_relaunches={case.get('session_relaunches', 0)} | slides={case['slide_count']} | images={case['image_count']} | warnings={case['warning_count']}"
            )
    lines.append("")
    return "\n".join(lines)


def render_review_queue_markdown(payload: dict[str, Any]) -> str:
    lines = [
        "# Manual AI Review Queue",
        "",
        f"- Run ID: `{payload['run_id']}`",
        f"- Overall Status: `{payload['overall_status']}`",
        "",
    ]
    items = payload.get("items") or []
    if not items:
        lines.append("No pending visual review items.")
        lines.append("")
        return "\n".join(lines)
    for item in items:
        lines.append(
            f"- `{item['case_id']}`: open [artifact]({item['artifact_pptx']}) and record the review decision."
        )
    lines.append("")
    return "\n".join(lines)


def render_session_check_markdown(payload: dict[str, Any]) -> str:
    lines = [
        "# ChatGPT Session Check",
        "",
        f"- Session: `{payload['session']}`",
        f"- Status: `{'PASS' if payload['ok'] else 'FAIL'}`",
        f"- Reason: `{payload['reason']}`",
        f"- Launcher: `{payload.get('launcher', '')}`",
        f"- Expected Profile Dir: `{payload['expected_profile_dir']}`",
        f"- Legacy Recovery Source: `{payload['legacy_recovery_source']}`",
    ]
    if payload.get("cdp_endpoint"):
        lines.append(f"- CDP Endpoint: `{payload['cdp_endpoint']}`")
    if payload.get("actual_profile_dir"):
        lines.append(f"- Actual Profile Dir: `{payload['actual_profile_dir']}`")
    if payload.get("browser_pid") is not None:
        lines.append(f"- Browser PID: `{payload['browser_pid']}`")
    if payload.get("page_url"):
        lines.append(f"- Page URL: `{payload['page_url']}`")
    if payload.get("page_title"):
        lines.append(f"- Page Title: `{payload['page_title']}`")
    if payload.get("manual_intervention_used"):
        lines.append(
            f"- Manual Intervention Elapsed: `{payload.get('manual_intervention_elapsed_seconds', 0.0)}`s / `{payload.get('manual_intervention_timeout_seconds', 0.0)}`s"
        )
    lines.extend(
        [
            f"- Has Composer: `{bool(payload.get('has_composer'))}`",
            f"- Has New Chat Button: `{bool(payload.get('has_new_chat_button'))}`",
            f"- No Sandbox Detected: `{bool(payload.get('no_sandbox_detected'))}`",
            f"- Manual Intervention Used: `{bool(payload.get('manual_intervention_used'))}`",
            "",
            "## Tabs",
            "",
        ]
    )
    tabs = payload.get("tabs") or []
    if not tabs:
        lines.append("No tab data was captured.")
    else:
        for tab in tabs:
            current = " current" if tab.get("current") else ""
            if tab.get("url"):
                lines.append(f"- `{tab['index']}`{current}: [{tab['title']}]({tab['url']})")
            else:
                lines.append(f"- `{tab['index']}`{current}: `{tab['title']}`")
    if payload.get("recent_response_failures"):
        lines.extend(["", "## Recent Failures", ""])
        for item in payload["recent_response_failures"]:
            lines.append(
                f"- `{item.get('status', 0)}` `{item.get('url', '')}`"
            )
    lines.extend(
        [
            "",
            "## Message",
            "",
            payload["message"],
            "",
            "## Manual Browser Preparation",
            "",
            "```powershell",
            *payload.get("manual_prepare_commands", []),
            "```",
            "",
        ]
    )
    return "\n".join(lines)


def _write_summary_csv(path: Path, cases: list[dict[str, Any]]) -> None:
    fieldnames = [
        "index",
        "chunk_index",
        "chunk_name",
        "chunk_label",
        "chat_index",
        "sample_id",
        "case_id",
        "prompt_id",
        "label",
        "status",
        "failure_class",
        "checker_ok",
        "checker_message",
        "blocking_issue_count",
        "warning_count",
        "generate_ok",
        "pptx_openable",
        "slide_count",
        "shape_count",
        "image_count",
        "last_completed_step",
        "stop_step",
        "transport_attempts",
        "session_relaunches",
        "no_sandbox_detected",
        "selected_page_url",
        "selected_page_title",
        "recent_response_failures",
        "review_required",
        "review_decision",
        "artifact_pptx",
        "case_dir",
    ]
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        for case in cases:
            writer.writerow({name: case.get(name) for name in fieldnames})


def _append_case_event(
    root_events_path: Path,
    case_events_path: Path,
    *,
    run_id: str,
    case_id: str,
    stage: str,
    status: str,
    duration_ms: float,
    message: str,
    artifact: str | None = None,
    failure_class: str | None = None,
) -> None:
    payload = {
        "run_id": run_id,
        "case_id": case_id,
        "stage": stage,
        "status": status,
        "duration_ms": duration_ms,
        "failure_class": failure_class,
        "message": message,
        "artifact": artifact,
    }
    _append_event(root_events_path, ts=_timestamp(), **payload)
    _append_event(case_events_path, ts=_timestamp(), **payload)


def _append_event(
    path: Path,
    *,
    run_id: str,
    case_id: str,
    stage: str,
    status: str,
    duration_ms: float,
    message: str,
    ts: str | None = None,
    failure_class: str | None = None,
    artifact: str | None = None,
) -> None:
    payload = {
        "ts": ts or _timestamp(),
        "run_id": run_id,
        "case_id": case_id,
        "stage": stage,
        "status": status,
        "duration_ms": round(duration_ms, 2),
        "failure_class": failure_class,
        "message": message,
        "artifact": artifact,
    }
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("a", encoding="utf-8") as handle:
        handle.write(json.dumps(payload, ensure_ascii=False) + "\n")


def _extract_slide_title(slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is not None and getattr(title_shape, "text", "").strip():
        return title_shape.text.strip()
    for shape in slide.shapes:
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            continue
        text = text_frame.text.strip()
        if text:
            return text
    return ""


def _write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def _write_json(path: Path, payload: Any) -> None:
    _write_text(path, json.dumps(payload, ensure_ascii=False, indent=2) + "\n")


def _build_run_id(suite_name: str) -> str:
    return f"manual-ai-{suite_name}-{datetime.now().astimezone().strftime('%Y%m%d-%H%M%S')}"


def _timestamp() -> str:
    return datetime.now().astimezone().isoformat(timespec="seconds")


def _duration_ms(started_at: float) -> float:
    return round((perf_counter() - started_at) * 1000, 2)


def _overall_message(status: str, *, session_check_only: bool = False) -> str:
    if session_check_only:
        if status == "FAIL":
            return "ChatGPT session readiness validation found a blocking issue."
        return "ChatGPT session readiness validation passed."
    if status == "FAIL":
        return "Automated verification found a blocking failure."
    if status == "REVIEW":
        return "Automated verification passed, but fixed representative visual review is still pending."
    return "Automated verification and representative visual review are complete."


def _report_progress(progress: Callable[[str], None] | None, message: str) -> None:
    if progress is not None:
        progress(message)
